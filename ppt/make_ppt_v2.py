from __future__ import annotations

import csv
import json
import math
import shutil
from pathlib import Path

import fitz
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pygments import lex
from pygments.lexers import PythonLexer
from pygments.token import Comment, Keyword, Literal, Name, Number, Operator, String


ROOT = Path(__file__).resolve().parents[1]
WORK = ROOT / "ppt_work" / "v2_excellent_style"
ASSETS = WORK / "assets"
DELIVERABLE = ROOT / "deliverables" / "SIBA_Jittor_培育期_优秀范例重做_20260728"
PPTX_PATH = DELIVERABLE / "姓名-培育期-SIBA-Jittor-可编辑-重做版.pptx"
PDF_PATH = DELIVERABLE / "姓名-培育期-SIBA-Jittor-重做版.pdf"
PREVIEW_DIR = DELIVERABLE / "preview_slides"
CONTACT_PATH = DELIVERABLE / "preview_contact.png"
OUTLINE_PATH = DELIVERABLE / "SIBA-Jittor-PPT完整大纲-重做版.md"
ANALYSIS_PATH = DELIVERABLE / "优秀PPT拆解与重做原则.md"
CHECK_PATH = DELIVERABLE / "PPT检查报告.md"

PAPER = ROOT / "paper_SIBA_ICCV2025.pdf"
JITTOR = ROOT / "siba_jittor"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_CN = "Microsoft YaHei"
FONT_CODE = "Consolas"
FONT_MATH = "Cambria Math"

BG = RGBColor(248, 247, 244)
CODE_BG = RGBColor(237, 239, 240)
TEXT = RGBColor(31, 31, 31)
MUTED = RGBColor(102, 102, 102)
LINE = RGBColor(205, 204, 200)
BLUE = RGBColor(69, 102, 153)
LIGHT_BLUE = RGBColor(217, 229, 246)
PURPLE = RGBColor(132, 113, 154)
LIGHT_PURPLE = RGBColor(231, 223, 239)
ORANGE = RGBColor(220, 142, 86)
LIGHT_ORANGE = RGBColor(249, 226, 209)
GREEN = RGBColor(82, 139, 101)
LIGHT_GREEN = RGBColor(220, 235, 221)
RED = RGBColor(183, 72, 67)
WHITE = RGBColor(255, 255, 255)


META: list[dict] = []
NOTES: list[str] = []


def rgb(value: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*value)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    font: str = FONT_CN,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    lines = text.split("\n")
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        paragraph.font.name = font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
    return shape


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: float = 20, color: RGBColor = TEXT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = 1.12
        paragraph.font.name = FONT_CN
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
    return shape


def add_rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor | None = None, width: float = 0.8, rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)
    return shape


def add_flow_box(slide, x, y, w, h, text, fill, line=None, size=18, bold=False):
    shape = add_rect(slide, x, y, w, h, fill, line or fill, 0.8, rounded=False)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = FONT_CN
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = TEXT
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=TEXT, width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = w / h
    if ratio >= box_ratio:
        real_h = w / ratio
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y + (h - real_h) / 2), width=Inches(w), height=Inches(real_h))
    real_w = h * ratio
    return slide.shapes.add_picture(str(path), Inches(x + (w - real_w) / 2), Inches(y), width=Inches(real_w), height=Inches(h))


def add_picture_cover(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = w / h
    if ratio >= box_ratio:
        real_w = h * ratio
        left = x - (real_w - w) / 2
        picture = slide.shapes.add_picture(str(path), Inches(left), Inches(y), width=Inches(real_w), height=Inches(h))
    else:
        real_h = w / ratio
        top = y - (real_h - h) / 2
        picture = slide.shapes.add_picture(str(path), Inches(x), Inches(top), width=Inches(w), height=Inches(real_h))
    return picture


def add_base_slide(prs: Presentation, title: str, page: int, source: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    add_text(slide, 0.52, 0.28, 9.8, 0.55, title, 32, TEXT, True)
    if source:
        add_text(slide, 9.25, 0.34, 3.45, 0.36, source, 18, MUTED, False, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.52), Inches(0.9), Inches(12.8), Inches(0.9))
    line.line.color.rgb = LINE
    line.line.width = Pt(0.8)
    footer = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.52), Inches(7.13), Inches(12.8), Inches(7.13))
    footer.line.color.rgb = LINE
    footer.line.width = Pt(0.6)
    add_text(slide, 12.42, 7.13, 0.36, 0.28, str(page), 18, MUTED, False, align=PP_ALIGN.RIGHT)
    return slide


def register(title: str, visible: list[str], images: list[str], note: str, code: str | None = None):
    META.append({"title": title, "visible": visible, "images": images, "note": note, "code": code})
    NOTES.append(note)


def token_color(token_type) -> RGBColor:
    if token_type in Comment:
        return GREEN
    if token_type in Keyword:
        return BLUE
    if token_type in Name.Function or token_type in Name.Class:
        return PURPLE
    if token_type in Name.Builtin:
        return rgb((38, 126, 126))
    if token_type in String or token_type in Literal.String:
        return ORANGE
    if token_type in Number:
        return RED
    if token_type in Operator:
        return rgb((80, 80, 80))
    return rgb((45, 50, 55))


def add_code(slide, code: str, x: float, y: float, w: float, h: float, start_line: int | None = None, size: float = 18):
    add_rect(slide, x, y, w, h, CODE_BG, None)
    shape = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.08), Inches(w - 0.24), Inches(h - 0.16))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    lexer = PythonLexer()
    lines = code.rstrip("\n").splitlines()
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 0.93
        if start_line is not None:
            run = paragraph.add_run()
            run.text = f"{start_line + index:>3}  "
            run.font.name = FONT_CODE
            run.font.size = Pt(size)
            run.font.color.rgb = rgb((135, 135, 135))
        if line == "":
            run = paragraph.add_run()
            run.text = " "
            run.font.name = FONT_CODE
            run.font.size = Pt(size)
            continue
        for token_type, value in lex(line, lexer):
            if value == "\n":
                continue
            run = paragraph.add_run()
            run.text = value
            run.font.name = FONT_CODE
            run.font.size = Pt(size)
            run.font.color.rgb = token_color(token_type)
    return shape


def read_lines(relative: str, start: int, end: int) -> str:
    lines = (JITTOR / relative).read_text(encoding="utf-8").splitlines()
    return "\n".join(lines[start - 1 : end])


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, col_widths: list[float] | None = None, font_size: float = 18, header=True):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for index, width in enumerate(col_widths):
            table.columns[index].width = Inches(w * width / total)
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            if header and row_index == 0:
                cell.fill.fore_color.rgb = rgb((63, 78, 96))
                color = WHITE
                bold = True
            else:
                cell.fill.fore_color.rgb = WHITE if row_index % 2 else rgb((241, 242, 242))
                color = TEXT
                bold = False
            cell.border_left = None
            frame = cell.text_frame
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)
                for run in paragraph.runs:
                    run.font.name = FONT_CN
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    run.font.color.rgb = color
    return table_shape


def prepare_cover_background():
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / "cover_gradient.png"
    width, height = 1920, 1080
    image = Image.new("RGB", (width, height), (244, 232, 244))
    pixels = image.load()
    for yy in range(height):
        for xx in range(width):
            tx = xx / width
            ty = yy / height
            r = int(235 + 16 * tx + 5 * (1 - ty))
            g = int(224 + 18 * tx + 8 * ty)
            b = int(246 - 8 * tx + 4 * ty)
            pixels[xx, yy] = (min(r, 255), min(g, 255), min(b, 255))
    overlay = Image.new("RGBA", image.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    draw.ellipse((1250, 620, 2050, 1350), fill=(255, 226, 188, 70))
    draw.ellipse((-300, -280, 850, 500), fill=(196, 190, 255, 60))
    overlay = overlay.filter(ImageFilter.GaussianBlur(100))
    Image.alpha_composite(image.convert("RGBA"), overlay).convert("RGB").save(path)
    return path


def render_paper_crops():
    ASSETS.mkdir(parents=True, exist_ok=True)
    document = fitz.open(PAPER)
    matrix = fitz.Matrix(3.0, 3.0)
    page4 = document[3].get_pixmap(matrix=matrix, alpha=False)
    page5 = document[4].get_pixmap(matrix=matrix, alpha=False)
    p4 = ASSETS / "paper_page4_3x.png"
    p5 = ASSETS / "paper_page5_3x.png"
    page4.save(p4)
    page5.save(p5)
    with Image.open(p4) as image:
        image.crop((115, 165, 1730, 780)).save(ASSETS / "fig2_architecture.png")
        image.crop((115, 1535, 930, 1840)).save(ASSETS / "fig3_cbsm.png")
    with Image.open(p5) as image:
        image.crop((95, 350, 940, 1000)).save(ASSETS / "fig4_cross_attention.png")


def prepare_assets():
    prepare_cover_background()
    render_paper_crops()
    shutil.copy2(ROOT / "results" / "training_analysis_20260727_siba_official_protocol" / "loss_curve.png", ASSETS / "loss_curve.png")


def build_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, ASSETS / "cover_gradient.png", 0, 0, SLIDE_W, SLIDE_H)
    add_text(slide, 0.78, 0.62, 7.4, 0.68, "Jittor 迁移演示", 40, rgb((57, 45, 67)), False)
    add_text(slide, 0.8, 1.5, 8.8, 0.95, "The Source Image is the Best Attention\nfor Infrared and Visible Image Fusion", 24, rgb((70, 61, 78)), False, font="Arial")
    add_text(slide, 0.8, 2.66, 9.2, 0.52, "Song Wang, Xie Han, Liqun Kuang, et al. · ICCV 2025", 20, rgb((82, 73, 90)))
    add_bullets(slide, ["模型方法", "Jittor 代码", "训练与测试", "实验结果"], 0.92, 3.52, 4.0, 2.0, 20, rgb((75, 67, 82)))
    add_text(slide, 0.82, 6.42, 4.4, 0.4, "汇报人：________", 20, rgb((75, 67, 82)))
    add_text(slide, 5.7, 6.42, 4.2, 0.4, "Jittor 复现 · ICCV 2025", 20, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 12.45, 7.08, 0.35, 0.3, "1", 18, rgb((100, 90, 107)), align=PP_ALIGN.RIGHT)
    note = "本次汇报复现 ICCV 2025 论文 SIBA。前半部分只说明论文的方法与核心模块，后半部分集中说明 Jittor 代码、训练、测试和实验结果。论文作者为 Song Wang、Xie Han、Liqun Kuang 等，正式发表于 ICCV 2025。"
    register(
        "Jittor 迁移演示",
        ["Jittor 迁移演示", "The Source Image is the Best Attention for Infrared and Visible Image Fusion", "Song Wang, Xie Han, Liqun Kuang, et al. · ICCV 2025", "模型方法", "Jittor 代码", "训练与测试", "实验结果", "汇报人：________", "Jittor 复现 · ICCV 2025"],
        ["淡粉紫渐变封面背景"],
        note,
    )


def build_method_slides(prs: Presentation):
    # 2 创新点
    slide = add_base_slide(prs, "创新点", 2, "Wang et al., ICCV 2025")
    source_images = [
        ("img_x324_p1_500x425.jpeg", "Grad-CAM 1"),
        ("img_x323_p1_500x417.jpeg", "Grad-CAM 2"),
        ("img_x322_p1_640x480.png", "红外原图"),
        ("img_x325_p1_640x480.png", "可见光原图"),
        ("img_x326_p1_640x480.jpeg", "红外伪彩色"),
        ("img_x327_p1_640x480.jpeg", "红外负变换"),
    ]
    paper_dir = ROOT / "ppt_work" / "paper"
    for index, (name, caption) in enumerate(source_images):
        row, col = divmod(index, 3)
        x = 0.55 + col * 2.85
        y = 1.12 + row * 2.7
        add_picture_contain(slide, paper_dir / name, x, y, 2.62, 1.92)
        add_text(slide, x, y + 1.94, 2.62, 0.42, caption, 18, TEXT, False, align=PP_ALIGN.CENTER)
    bullets = [
        "红外热目标与 Grad-CAM 高响应区域具有相近的空间分布",
        "原图与负变换分别提供高响应区和一般区域",
        "CBSM 处理后的源图像直接作为交叉注意力查询 Q",
    ]
    add_bullets(slide, bullets, 9.25, 1.46, 3.35, 4.4, 20)
    add_text(slide, 9.28, 5.72, 3.25, 0.75, "源图像不再只提供像素，\n还直接参与注意力计算。", 22, BLUE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    note = "论文的出发点来自一个很直观的观察：红外图像中的高热区域，与分类网络 Grad-CAM 中的高响应区域在空间分布上很相似。作者因此不再只从深层特征中生成查询，而是把源图像送入 CBSM 后直接作为查询。红外原图和负变换分别强调热目标区域和其余区域，可见光分支也采用同样的处理。"
    register("创新点", bullets + ["源图像不再只提供像素，还直接参与注意力计算。", "Wang et al., ICCV 2025"], [f"论文 Fig.1 子图：{caption}" for _, caption in source_images], note)

    # 3 基本架构
    slide = add_base_slide(prs, "基本架构", 3, "Wang et al., ICCV 2025 · Fig. 2")
    add_picture_contain(slide, ASSETS / "fig2_architecture.png", 0.62, 1.1, 12.05, 4.72)
    add_bullets(slide, ["红外、可见光分别提取 48 通道特征", "原图和负变换经 CBSM 形成四个查询", "四路交叉注意力特征拼接后重建融合图"], 0.84, 5.58, 11.6, 1.35, 18)
    note = "整体结构分为三段。首先，两种模态分别通过 SE-ResNet 和 Restormer 得到中间特征。其次，红外、可见光及其负变换经过 CBSM，形成四个查询，分别引导另一模态的特征。最后，四路结果按通道拼接，再由两个 SE-ResNet 模块和输出层生成单通道融合结果。"
    register("基本架构", ["红外、可见光分别提取 48 通道特征", "原图和负变换经 CBSM 形成四个查询", "四路交叉注意力特征拼接后重建融合图", "Wang et al., ICCV 2025 · Fig. 2"], ["论文 Fig.2 总体架构，已去除图注"], note)

    # 4 CBSM
    slide = add_base_slide(prs, "CBSM", 4, "Wang et al., ICCV 2025 · Fig. 3")
    add_picture_contain(slide, ASSETS / "fig3_cbsm.png", 0.65, 1.28, 5.72, 3.78)
    add_flow_box(slide, 7.08, 1.42, 1.55, 0.72, "源图像\n1×H×W", LIGHT_BLUE, BLUE, 18)
    add_arrow(slide, 8.68, 1.78, 9.22, 1.78, BLUE)
    add_flow_box(slide, 9.28, 1.42, 1.5, 0.72, "Conv 3×3\n+ PReLU", LIGHT_ORANGE, ORANGE, 18)
    add_arrow(slide, 10.83, 1.78, 11.35, 1.78, ORANGE)
    add_flow_box(slide, 11.42, 1.42, 1.1, 0.72, "Conv\n3×3", LIGHT_ORANGE, ORANGE, 18)
    add_arrow(slide, 12.0, 2.18, 12.0, 2.68, ORANGE)
    add_flow_box(slide, 10.72, 2.75, 2.55, 0.72, "SE 通道增强", LIGHT_GREEN, GREEN, 19)
    add_arrow(slide, 10.65, 3.12, 9.95, 3.12, GREEN)
    add_flow_box(slide, 8.22, 2.75, 1.68, 0.72, "PReLU", LIGHT_PURPLE, PURPLE, 19)
    add_bullets(slide, ["3×3 卷积完成空间映射", "SE 抑制冗余通道并增强有效响应", "输出通道数由 1 扩展为 48，与特征维度一致"], 6.9, 4.18, 5.6, 1.8, 20)
    add_text(slide, 0.78, 5.46, 5.45, 0.72, "CBSM 先整理源图像，再把它送入交叉注意力。", 22, BLUE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    note = "CBSM 的作用不是提取深层语义，而是把单通道源图像调整到与中间特征相同的 48 通道空间。两层卷积完成空间映射，SE 模块根据通道统计重新分配权重，PReLU 保留非线性表达。这样得到的查询既保留源图像的空间位置，又减少直接使用原始像素时的冗余和误导信息。"
    register("CBSM", ["源图像 1×H×W", "Conv 3×3 + PReLU", "Conv 3×3", "SE 通道增强", "PReLU", "3×3 卷积完成空间映射", "SE 抑制冗余通道并增强有效响应", "输出通道数由 1 扩展为 48，与特征维度一致", "CBSM 先整理源图像，再把它送入交叉注意力。"], ["论文 Fig.3 CBSM 结构，已去除图注"], note)

    # 5 I-SCA / V-SCA
    slide = add_base_slide(prs, "I-SCA / V-SCA", 5, "Wang et al., ICCV 2025 · Fig. 4")
    add_picture_contain(slide, ASSETS / "fig4_cross_attention.png", 0.55, 1.12, 6.45, 5.58)
    add_text(slide, 7.4, 1.38, 5.2, 0.52, "查询来自源图像，键和值来自另一模态特征", 22, BLUE, True)
    add_text(slide, 7.42, 2.02, 5.0, 0.5, "Q = CBSM(Iraw)", 23, PURPLE, True, font=FONT_MATH)
    add_text(slide, 7.42, 2.62, 5.0, 0.5, "K, V = DConv(Conv1×1(LN(Fmid)))", 23, PURPLE, True, font=FONT_MATH)
    add_bullets(slide, ["I-SCA：红外源图引导可见光特征", "V-SCA：可见光源图引导红外特征", "原图和负变换各产生一路查询", "GDFN 过滤注意力输出中的冗余信息"], 7.38, 3.35, 5.0, 2.45, 20)
    note = "I-SCA 和 V-SCA 的结构相同，区别只在输入。以 I-SCA 为例，红外源图经 CBSM 得到查询 Q，可见光中间特征产生 K 和 V。V-SCA 则交换两种模态。每种源图还使用负变换形成第二个查询，因此一共得到四路交叉注意力结果。注意力后接 GDFN，用门控方式保留有效响应。"
    register("I-SCA / V-SCA", ["查询来自源图像，键和值来自另一模态特征", "Q = CBSM(Iraw)", "K, V = DConv(Conv1×1(LN(Fmid)))", "I-SCA：红外源图引导可见光特征", "V-SCA：可见光源图引导红外特征", "原图和负变换各产生一路查询", "GDFN 过滤注意力输出中的冗余信息"], ["论文 Fig.4 I-DCA/V-DCA 与 GDFN，已去除图注"], note)

    # 6 最终融合和损失函数
    slide = add_base_slide(prs, "最终融合和损失函数", 6, "official implementation")
    colors = [(LIGHT_BLUE, BLUE), (LIGHT_ORANGE, ORANGE), (LIGHT_GREEN, GREEN), (LIGHT_PURPLE, PURPLE)]
    labels = ["ir → vi", "1−ir → vi", "vi → ir", "1−vi → ir"]
    for index, (label, pair) in enumerate(zip(labels, colors)):
        add_flow_box(slide, 0.72 + index * 2.15, 1.38, 1.68, 0.72, label, pair[0], pair[1], 18)
        add_arrow(slide, 1.56 + index * 2.15, 2.12, 5.52, 2.95, pair[1], 1.2)
    add_flow_box(slide, 5.05, 2.88, 2.0, 0.78, "Concat\n192 通道", rgb((232, 231, 228)), rgb((150, 150, 145)), 19, True)
    add_arrow(slide, 7.1, 3.27, 7.78, 3.27, TEXT)
    add_flow_box(slide, 7.84, 2.88, 1.72, 0.78, "Res_SE\n192 → 96", LIGHT_ORANGE, ORANGE, 18)
    add_arrow(slide, 9.62, 3.27, 10.15, 3.27, TEXT)
    add_flow_box(slide, 10.22, 2.88, 1.72, 0.78, "Res_SE\n96 → 48", LIGHT_GREEN, GREEN, 18)
    add_arrow(slide, 11.98, 3.27, 12.42, 3.27, TEXT)
    add_flow_box(slide, 12.44, 2.88, 0.65, 0.78, "1", LIGHT_BLUE, BLUE, 18, True)
    add_text(slide, 0.82, 4.48, 5.8, 0.48, "Ltotal = 0.1 Lint + 1.0 LSobel + 10.0 LLap", 24, TEXT, True, font=FONT_MATH)
    add_bullets(slide, ["强度项：接近两幅源图的逐像素最大值", "Sobel 项：保留较强的一阶边缘", "Laplacian 项：保留较强的局部细节"], 0.84, 5.1, 6.1, 1.35, 19)
    add_text(slide, 7.5, 4.62, 4.9, 0.72, "损失权重与官方 train.py 一致", 22, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 7.55, 5.38, 4.8, 0.72, "60 epochs · batch 4 · patch 128", 22, TEXT, False, align=PP_ALIGN.CENTER)
    note = "四路交叉注意力结果按通道拼接，通道数从 4×48 变为 192。随后两个 SE-ResNet 模块依次降到 96 和 48 通道，最后输出单通道融合图。官方训练代码的总损失由三部分组成：强度损失权重 0.1，Sobel 梯度损失权重 1，Laplacian 梯度损失权重 10。"
    register("最终融合和损失函数", labels + ["Concat 192 通道", "Res_SE 192 → 96", "Res_SE 96 → 48", "Ltotal = 0.1 Lint + 1.0 LSobel + 10.0 LLap", "强度项：接近两幅源图的逐像素最大值", "Sobel 项：保留较强的一阶边缘", "Laplacian 项：保留较强的局部细节", "损失权重与官方 train.py 一致", "60 epochs · batch 4 · patch 128"], ["可编辑流程图"], note)


def code_slide(prs, page, title, relative, start, end, note, bullets=None, source=None, code_width=12.0):
    slide = add_base_slide(prs, title, page, source or relative.replace("\\", "/"))
    code = read_lines(relative, start, end)
    if bullets:
        add_code(slide, code, 0.58, 1.12, code_width, 5.75, start, 18)
        add_bullets(slide, bullets, 0.72 + code_width, 1.42, 12.0 - code_width, 4.7, 19)
    else:
        add_code(slide, code, 0.58, 1.12, 12.15, 5.75, start, 18)
    register(title, bullets or [], [f"可编辑代码：{relative} 第 {start}–{end} 行"], note, code)


def build_code_slides(prs: Presentation):
    code_slide(prs, 7, "顶层模型：源图查询分支", "models/SIBA.py", 9, 29,
               "顶层类首先建立红外和可见光两条特征提取分支，并为每条分支配置自注意力。随后定义四个 CBSM，分别处理红外原图、红外负变换、可见光原图和可见光负变换。Jittor 中把 PyTorch 的 forward 改为 execute，但模块结构和参数保持不变。",
               ["两路特征提取", "四个 CBSM 查询"], code_width=9.2)

    code_slide(prs, 8, "顶层模型：交叉注意力与输出", "models/SIBA.py", 31, 45,
               "这里定义四组交叉注意力模块。ir2vi 和 irI2vi 使用红外原图及其负变换引导可见光特征；vi2ir 和 viI2ir 反向引导红外特征。四路结果拼接后，通过两个 Res_SE 和一个输出层恢复为单通道图像。",
               ["四路交叉注意力", "192 → 96 → 48 → 1"], code_width=9.2)

    code_slide(prs, 9, "顶层模型：输入与四个查询", "models/SIBA.py", 60, 81,
               "前向开始时保留两幅源图，并显式计算 1-x 的负变换。红外和可见光先经过各自的卷积分支和自注意力。四幅源图随后分别通过对应 CBSM，得到四个查询张量。这里没有改变官方输入顺序或负变换公式。",
               ["负变换：1 − x", "CBSM 输出与中间特征同为 48 通道"], code_width=9.2)

    code_slide(prs, 10, "顶层模型：四路交互与融合", "models/SIBA.py", 83, 103,
               "四个查询进入对应的交叉注意力。红外查询作用于可见光特征，可见光查询作用于红外特征。最后使用 jt.concat 沿通道维拼接，这对应 PyTorch 的 torch.cat；融合卷积和输出层的调用顺序保持不变。",
               ["Q 来自源图", "K、V 来自另一模态", "jt.concat 对应 torch.cat"], code_width=9.2)

    code_slide(prs, 11, "特征提取：SE-ResNet", "base_blocks/se_resnet.py", 8, 36,
               "SE-ResNet 由两层卷积、PReLU、SE 和残差连接组成。当输入通道和输出通道不一致时，用 1×1 卷积调整残差分支。Jittor 版本只把 forward 改为 execute，卷积、残差和 SE 的执行顺序与官方代码一致。",
               ["1×1 卷积匹配通道", "两层 3×3 卷积", "SE 后再做残差相加"], code_width=9.2)

    slide = add_base_slide(prs, "CBSM：通道增强与空间映射", 12, "base_blocks/cbsm.py · base_blocks/SE.py")
    code1 = read_lines("base_blocks/cbsm.py", 7, 17)
    code2 = read_lines("base_blocks/SE.py", 7, 25)
    add_code(slide, code1, 0.58, 1.12, 6.0, 5.72, 7, 18)
    add_code(slide, code2, 6.78, 1.12, 5.95, 5.72, 7, 18)
    note = "左侧是 CBSM 主体：两层 3×3 卷积之间使用 PReLU，之后通过 SE 再激活。右侧是 SE 的实现：全局平均池化得到每个通道的统计量，两层全连接生成 0 到 1 的通道权重，再与原特征逐通道相乘。两个文件在 Jittor 中保持原有结构。"
    register("CBSM：通道增强与空间映射", [], ["可编辑代码：base_blocks/cbsm.py 第 7–17 行", "可编辑代码：base_blocks/SE.py 第 7–25 行"], note, code1 + "\n\n" + code2)

    code_slide(prs, 13, "自注意力：Restormer Block", "base_blocks/restormer.py", 12, 38,
               "自注意力块先根据需要用 1×1 卷积匹配通道，然后执行 LayerNorm、自注意力、残差相加，再执行 LayerNorm、GDFN 和第二次残差相加。这个顺序直接对应官方 Restormer Block。",
               ["LN → Attention → Residual", "LN → GDFN → Residual"], code_width=9.2)

    code_slide(prs, 14, "自注意力：QKV 与多头计算", "base_blocks/restormer.py", 87, 127,
               "Q、K、V 由 1×1 卷积和深度卷积共同产生，然后按多头形式重排。Q 和 K 在最后一维做 L2 归一化，计算通道维注意力，再与 V 相乘并恢复到四维特征。Jittor 使用 jittor.einops.rearrange 保留原始张量变换。",
               ["1×1 + depthwise 3×3 生成 QKV", "L2 normalize", "多头注意力后恢复 BCHW"], code_width=9.2)

    code_slide(prs, 15, "交叉注意力：K、V 生成", "base_blocks/restormer.py", 131, 158,
               "交叉注意力只从当前模态特征 kv 中生成 K 和 V，查询 q 由 CBSM 直接提供。K、V 使用 1×1 卷积和深度卷积生成，随后与 q 一起重排到多头形式。输入接口 execute(kv, q) 对应论文的源图查询设计。",
               ["execute(kv, q)", "Q 不由 kv 特征生成"], code_width=9.2)

    code_slide(prs, 16, "交叉注意力：全局交互", "base_blocks/restormer.py", 160, 174,
               "Q 和 K 归一化后计算注意力矩阵，Softmax 后与 V 相乘。注意力计算发生在通道维，空间位置被展平到最后一维，因此复杂度随像素数线性增长。最后使用 1×1 卷积投影并返回。",
               ["Attention = Softmax(QKᵀ · scale)", "输出 = Attention · V"], code_width=9.2)

    slide = add_base_slide(prs, "GDFN 与 LayerNorm", 17, "base_blocks/restormer.py")
    code1 = read_lines("base_blocks/restormer.py", 178, 199)
    code2 = read_lines("base_blocks/restormer.py", 221, 246)
    add_code(slide, code1, 0.58, 1.12, 6.0, 5.72, 178, 18)
    add_code(slide, code2, 6.78, 1.12, 5.95, 5.72, 221, 18)
    note = "GDFN 先把通道扩展为两倍，再用深度卷积拆成两支，一支经过 GELU 后与另一支逐元素相乘，最后投影回原通道。LayerNorm 先把 BCHW 转换为 B(HW)C，按通道归一化，再恢复四维。Jittor 版本保留官方的 biased LayerNorm 公式。"
    register("GDFN 与 LayerNorm", [], ["可编辑代码：base_blocks/restormer.py 第 178–199 行", "可编辑代码：base_blocks/restormer.py 第 221–246 行"], note, code1 + "\n\n" + code2)

    slide = add_base_slide(prs, "强度损失与 Sobel 梯度损失", 18, "loss/loss.py")
    code1 = read_lines("loss/loss.py", 41, 60)
    code2 = read_lines("loss/loss.py", 63, 81)
    add_code(slide, code1, 0.58, 1.12, 6.0, 5.72, 41, 18)
    add_code(slide, code2, 6.78, 1.12, 5.95, 5.72, 63, 18)
    note = "强度损失先取红外与可见光亮度的逐像素最大值，再与融合图做 L1。Sobel 模块固定两组 3×3 核，分别计算水平和垂直梯度，取绝对值之和。融合图梯度需要接近两幅源图中较强的梯度。固定卷积核通过 stop_grad 排除在参数更新之外。"
    register("强度损失与 Sobel 梯度损失", [], ["可编辑代码：loss/loss.py 第 41–60 行", "可编辑代码：loss/loss.py 第 63–81 行"], note, code1 + "\n\n" + code2)

    slide = add_base_slide(prs, "Laplacian 梯度损失", 19, "loss/loss.py")
    code1 = read_lines("loss/loss.py", 6, 18)
    code2 = read_lines("loss/loss.py", 23, 35)
    add_code(slide, code1, 0.58, 1.12, 6.0, 5.72, 6, 18)
    add_code(slide, code2, 6.78, 1.12, 5.95, 5.72, 23, 18)
    note = "官方 PyTorch 版本通过 Kornia 计算 Laplacian。Jittor 侧按照相同核定义手工构造卷积，并采用 reflect padding 和分组卷积。JointGrad 比较两幅源图 Laplacian 绝对值，逐位置选择更强者，再约束融合图的 Laplacian 与之接近。"
    register("Laplacian 梯度损失", [], ["可编辑代码：loss/loss.py 第 6–18 行", "可编辑代码：loss/loss.py 第 23–35 行"], note, code1 + "\n\n" + code2)

    code_slide(prs, 20, "数据加载与随机裁剪", "loader/train_loader.py", 9, 41,
               "训练集按文件名排序读取红外和可见光图像，并在 __getitem__ 中检查文件名一致。随机裁剪完全保留官方范围：边界留 10 像素，裁剪 128×128 patch。Jittor Dataset 通过 set_attrs 设置总长度和批处理属性。",
               ["完整训练集：1083 MSRS + 200 RoadScene", "同一坐标裁剪两种模态", "输入归一化到 [0, 1]"], code_width=8.95)

    slide = add_base_slide(prs, "训练配置", 21, "args/args_SIBA.py · train.py")
    code1 = read_lines("args/args_SIBA.py", 3, 20)
    code2 = read_lines("train.py", 32, 46)
    add_code(slide, code1, 0.58, 1.12, 5.65, 5.72, 3, 18)
    add_code(slide, code2, 6.43, 1.12, 6.3, 5.72, 32, 18)
    note = "训练配置与官方仓库一致：60 轮、batch size 4、patch 128、初始学习率 1e-4，StepLR 每 25 轮衰减为原来的 0.5，weight decay 为 0。数据加载器开启 shuffle，drop_last 为 False。模型、损失和调度器的建立顺序保持一致。"
    register("训练配置", [], ["可编辑代码：args/args_SIBA.py 第 3–20 行", "可编辑代码：train.py 第 32–46 行"], note, code1 + "\n\n" + code2)

    code_slide(prs, 22, "训练循环与权重保存", "train.py", 52, 98,
               "每个 batch 依次完成前向、三项损失、反向传播、全局梯度裁剪和 Adam 更新。官方日志每 50 个 batch 打印一次，因此 60 轮共得到 420 条 loss。每轮结束后调用 StepLR，学习率下限设为 1e-6，最终保存 epoch60 权重。",
               ["loss = 10×Laplacian + 0.1×Intensity + Sobel", "每 50 batch 记录一次", "最终保存 SIBA_epoch60.pkl"], code_width=8.95)

    slide = add_base_slide(prs, "Adam 与梯度裁剪", 23, "compat/pytorch_adam.py · compat/pytorch_clip.py")
    code1 = read_lines("compat/pytorch_adam.py", 40, 68)
    code2 = read_lines("compat/pytorch_clip.py", 4, 20)
    add_code(slide, code1, 0.58, 1.12, 7.0, 5.72, 40, 18)
    add_code(slide, code2, 7.78, 1.12, 4.95, 5.72, 4, 18)
    note = "Jittor 内置优化器的更新细节与官方 PyTorch Adam 不完全相同，因此这里显式保存一阶矩、二阶矩和 bias correction，按 PyTorch 公式更新参数。梯度裁剪同样按所有参数梯度的全局 L2 范数计算统一缩放系数。给定相同参考梯度时，这两个兼容实现能够与 PyTorch 对齐。"
    register("Adam 与梯度裁剪", [], ["可编辑代码：compat/pytorch_adam.py 第 40–68 行", "可编辑代码：compat/pytorch_clip.py 第 4–20 行"], note, code1 + "\n\n" + code2)

    slide = add_base_slide(prs, "测试与颜色重建", 24, "test.py · utils/RGB2YCrBb.py")
    code1 = read_lines("test.py", 19, 43)
    code2 = read_lines("utils/RGB2YCrBb.py", 23, 37)
    add_code(slide, code1, 0.58, 1.12, 6.55, 5.72, 19, 18)
    add_code(slide, code2, 7.33, 1.12, 5.4, 5.72, 23, 18)
    note = "测试阶段加载 Jittor 权重，逐张生成融合亮度 Y。可见光输入同时保留 Cb 和 Cr，融合结束后使用与官方代码相同的矩阵把 YCbCr 转回 RGB。输出先限制在 0 到 1，再保存为原文件名。MSRS、M3FD 和 TNO 都按这一流程完成推理。"
    register("测试与颜色重建", [], ["可编辑代码：test.py 第 19–43 行", "可编辑代码：utils/RGB2YCrBb.py 第 23–37 行"], note, code1 + "\n\n" + code2)


def manual_code_slide(prs, page, title, code, note, bullets=None, source=None, code_width=12.0):
    slide = add_base_slide(prs, title, page, source)
    if bullets:
        add_code(slide, code.strip("\n"), 0.58, 1.12, code_width, 5.75, None, 18)
        add_bullets(slide, bullets, 0.72 + code_width, 1.42, 12.0 - code_width, 4.7, 19)
    else:
        add_code(slide, code.strip("\n"), 0.58, 1.12, 12.15, 5.75, None, 18)
    register(title, bullets or [], [f"可编辑代码：{source or title}"], note, code.strip("\n"))


def build_code_slides_v2(prs: Presentation):
    manual_code_slide(
        prs, 7, "顶层模型：源图查询分支",
        """
class SIBA(nn.Module):
    def __init__(
        self, in_cha=1, mid_cha=48,
        out_cha=1, SA_depths=1, CA_depths=1
    ):
        super(SIBA, self).__init__()
        self.ir_conv = Res_SE(in_cha, mid_cha)
        self.vi_conv = Res_SE(in_cha, mid_cha)

        self.ir_sa = nn.ModuleList()
        self.vi_sa = nn.ModuleList()
        for _ in range(SA_depths):
            self.ir_sa.append(SA(mid_cha, mid_cha))
            self.vi_sa.append(SA(mid_cha, mid_cha))

        self.weight_ir = CBSM(mid_cha)
        self.weight_irI = CBSM(mid_cha)
        self.weight_vi = CBSM(mid_cha)
        self.weight_viI = CBSM(mid_cha)
""",
        "顶层类先建立红外和可见光两条特征提取分支，并为每条分支配置自注意力。随后定义四个 CBSM，分别处理红外原图、红外负变换、可见光原图和可见光负变换。Jittor 中把 PyTorch 的 forward 改为 execute，模块结构和参数保持不变。",
        ["两路特征提取", "四个 CBSM 查询"], "models/SIBA.py", 9.25,
    )

    manual_code_slide(
        prs, 8, "顶层模型：交叉注意力与输出",
        """
self.ir2vi_ca = nn.ModuleList()
self.irI2vi_ca = nn.ModuleList()
self.vi2ir_ca = nn.ModuleList()
self.viI2ir_ca = nn.ModuleList()
for _ in range(CA_depths):
    self.ir2vi_ca.append(CA(mid_cha, mid_cha))
    self.irI2vi_ca.append(CA(mid_cha, mid_cha))
    self.vi2ir_ca.append(CA(mid_cha, mid_cha))
    self.viI2ir_ca.append(CA(mid_cha, mid_cha))

self.fuse_conv = nn.Sequential(
    Res_SE(mid_cha * 4, mid_cha * 2),
    Res_SE(mid_cha * 2, mid_cha),
)
self.out_conv = Res_SE(
    mid_cha, out_cha, use_se=False
)
""",
        "这里定义四组交叉注意力模块。ir2vi 和 irI2vi 使用红外原图及其负变换引导可见光特征；vi2ir 和 viI2ir 反向引导红外特征。四路结果拼接后，通过两个 Res_SE 和一个输出层恢复为单通道图像。",
        ["四路交叉注意力", "192 → 96 → 48 → 1"], "models/SIBA.py", 9.25,
    )

    manual_code_slide(
        prs, 9, "顶层模型：输入与四个查询",
        """
def execute(self, ir, vi):
    ir_raw = ir
    ir_raw_invert = 1 - ir
    vi_raw = vi
    vi_raw_invert = 1 - vi

    ir = self.ir_conv(ir)
    vi = self.vi_conv(vi)

    ir_sa, vi_sa = ir, vi
    for layer_ir, layer_vi in zip(
        self.ir_sa, self.vi_sa
    ):
        ir_sa = layer_ir(ir_sa)
        vi_sa = layer_vi(vi_sa)

    w_ir = self.weight_ir(ir_raw)
    w_irI = self.weight_irI(ir_raw_invert)
    w_vi = self.weight_vi(vi_raw)
    w_viI = self.weight_viI(vi_raw_invert)
""",
        "前向开始时保留两幅源图，并显式计算 1-x 的负变换。红外和可见光先经过各自的卷积分支和自注意力。四幅源图随后分别通过对应 CBSM，得到四个查询张量。这里没有改变官方输入顺序或负变换公式。",
        ["负变换：1 − x", "查询通道数为 48"], "models/SIBA.py", 9.25,
    )

    manual_code_slide(
        prs, 10, "顶层模型：四路交互与融合",
        """
ir2vi_ca = vi_sa
irI2vi_ca = vi_sa
for layer_ir, layer_irI in zip(
    self.ir2vi_ca, self.irI2vi_ca
):
    ir2vi_ca = layer_ir(ir2vi_ca, w_ir)
    irI2vi_ca = layer_irI(irI2vi_ca, w_irI)

vi2ir_ca = ir_sa
viI2ir_ca = ir_sa
for layer_vi, layer_viI in zip(
    self.vi2ir_ca, self.viI2ir_ca
):
    vi2ir_ca = layer_vi(vi2ir_ca, w_vi)
    viI2ir_ca = layer_viI(viI2ir_ca, w_viI)

mixed = jt.concat([
    ir2vi_ca, vi2ir_ca, irI2vi_ca, viI2ir_ca
], dim=1)
return self.out_conv(self.fuse_conv(mixed))
""",
        "四个查询进入对应的交叉注意力。红外查询作用于可见光特征，可见光查询作用于红外特征。最后使用 jt.concat 沿通道维拼接，这对应 PyTorch 的 torch.cat；融合卷积和输出层的调用顺序保持不变。",
        ["Q 来自源图", "K、V 来自另一模态", "jt.concat 对应 torch.cat"], "models/SIBA.py", 9.25,
    )

    manual_code_slide(
        prs, 11, "特征提取：SE-ResNet",
        """
def execute(self, input):
    residual = self.res_conv(input)
    out = self.conv1(input)
    out = self.act(out)
    out = self.conv2(out)
    out = self.se(out)
    if self.use_res:
        out = out + residual
    out = self.act(out)
    return out
""",
        "SE-ResNet 由两层卷积、PReLU、SE 和残差连接组成。当输入通道和输出通道不一致时，用 1×1 卷积调整残差分支。Jittor 版本只把 forward 改为 execute，卷积、残差和 SE 的执行顺序与官方代码一致。",
        ["1×1 卷积匹配通道", "两层 3×3 卷积", "SE 后做残差相加"], "base_blocks/se_resnet.py", 9.25,
    )

    slide = add_base_slide(prs, "CBSM：通道增强与空间映射", 12, "cbsm.py · SE.py")
    code1 = """
# CBSM
def execute(self, source_image):
    tmp = self.conv1(source_image)
    tmp = self.act(tmp)
    tmp = self.conv2(tmp)
    tmp = self.se(tmp)
    return self.act(tmp)
"""
    code2 = """
# Squeeze-and-Excitation
def execute(self, x):
    b, c, _, _ = x.shape
    weight = self.avg_pool(x)
    weight = weight.view(b, c)
    weight = self.fc(weight)
    weight = weight.view(
        b, c, 1, 1
    )
    return x * weight
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 6.0, 5.72, None, 18)
    add_code(slide, code2.strip(), 6.78, 1.12, 5.95, 5.72, None, 18)
    note = "左侧是 CBSM 主体：两层 3×3 卷积之间使用 PReLU，之后通过 SE 再激活。右侧是 SE 的实现：全局平均池化得到每个通道的统计量，两层全连接生成 0 到 1 的通道权重，再与原特征逐通道相乘。两个文件在 Jittor 中保持原有结构。"
    register("CBSM：通道增强与空间映射", [], ["可编辑代码：cbsm.py", "可编辑代码：SE.py"], note, code1 + "\n" + code2)

    manual_code_slide(
        prs, 13, "自注意力：Restormer Block",
        """
class TransformerBlock_SA(nn.Module):
    def __init__(self, in_channel, out_channel, num_heads=8):
        self.conv = (
            nn.Conv2d(in_channel, out_channel, 1)
            if out_channel != in_channel
            else nn.Identity()
        )
        self.norm1 = LayerNorm(out_channel, 'WithBias')
        self.attn = Self_Attention(
            out_channel, num_heads=num_heads
        )
        self.norm2 = LayerNorm(out_channel, 'WithBias')
        self.mlp = Mlp(in_features=out_channel)

    def execute(self, x):
        x = self.conv(x)
        x = x + self.attn(self.norm1(x))
        x = x + self.mlp(self.norm2(x))
        return x
""",
        "自注意力块先根据需要用 1×1 卷积匹配通道，然后执行 LayerNorm、自注意力、残差相加，再执行 LayerNorm、GDFN 和第二次残差相加。这个顺序直接对应官方 Restormer Block。",
        ["LN → Attention → Residual", "LN → GDFN → Residual"], "restormer.py", 9.25,
    )

    slide = add_base_slide(prs, "自注意力：QKV 与多头计算", 14, "restormer.py")
    code1 = """
def execute(self, x):
    b, c, h, w = x.shape
    q, k, v = self.qkv2(self.qkv1(x)).chunk(3, dim=1)
    q = rearrange(q, 'b (head c) h w -> b head c (h w)', head=self.num_heads)
    k = rearrange(k, 'b (head c) h w -> b head c (h w)', head=self.num_heads)
    v = rearrange(v, 'b (head c) h w -> b head c (h w)', head=self.num_heads)
    q = normalize(q, dim=-1)
    k = normalize(k, dim=-1)
    attn = ((q @ k.transpose(-2, -1)) * self.scale).softmax(-1)
    out = attn @ v
    out = rearrange(
        out, 'b head c (h w) -> b (head c) h w',
        head=self.num_heads, h=h, w=w
    )
    return self.proj(out)
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 12.15, 5.72, None, 18)
    note = "Q、K、V 由 1×1 卷积和深度卷积共同产生，然后按多头形式重排。Q 和 K 在最后一维做 L2 归一化，计算通道维注意力，再与 V 相乘并恢复到四维特征。Jittor 使用 jittor.einops.rearrange 保留原始张量变换。"
    register("自注意力：QKV 与多头计算", [], ["可编辑代码：Self_Attention.execute"], note, code1)

    slide = add_base_slide(prs, "交叉注意力：K、V 生成", 15, "restormer.py")
    code1 = """
def execute(self, kv, q):
    b, c, h, w = kv.shape
    k, v = self.qkv2(self.qkv1(kv)).chunk(2, dim=1)

    q = rearrange(q, 'b (head c) h w -> b head c (h w)', head=self.num_heads)
    k = rearrange(k, 'b (head c) h w -> b head c (h w)', head=self.num_heads)
    v = rearrange(v, 'b (head c) h w -> b head c (h w)', head=self.num_heads)
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 12.15, 5.72, None, 18)
    note = "交叉注意力只从当前模态特征 kv 中生成 K 和 V，查询 q 由 CBSM 直接提供。K、V 使用 1×1 卷积和深度卷积生成，随后与 q 一起重排到多头形式。输入接口 execute(kv, q) 对应论文的源图查询设计。"
    register("交叉注意力：K、V 生成", [], ["可编辑代码：Cross_Attention.execute(kv, q)"], note, code1)

    manual_code_slide(
        prs, 16, "交叉注意力：全局交互",
        """
q = normalize(q, dim=-1)
k = normalize(k, dim=-1)

attn = q @ k.transpose(-2, -1)
attn = (attn * self.scale).softmax(dim=-1)
out = attn @ v

out = rearrange(
    out,
    'b head c (h w) -> b (head c) h w',
    head=self.num_heads,
    h=h,
    w=w,
)
out = self.proj(out)
return out
""",
        "Q 和 K 归一化后计算注意力矩阵，Softmax 后与 V 相乘。注意力计算发生在通道维，空间位置被展平到最后一维，因此复杂度随像素数线性增长。最后使用 1×1 卷积投影并返回。",
        ["Attention = Softmax(QKᵀ · scale)", "输出 = Attention · V"], "restormer.py", 9.25,
    )

    slide = add_base_slide(prs, "GDFN 与 LayerNorm", 17, "restormer.py")
    code1 = """
# GDFN
def execute(self, x):
    x = self.project_in(x)
    x1, x2 = self.dwconv(x).chunk(
        2, dim=1
    )
    x = nn.gelu(x1) * x2
    return self.project_out(x)
"""
    code2 = """
# WithBias LayerNorm
def execute(self, x):
    mu = x.mean(-1, keepdims=True)
    sigma = x.var(
        -1,
        keepdims=True,
        unbiased=False,
    )
    x = (x - mu) / jt.sqrt(
        sigma + 1e-5
    )
    return x * self.weight + self.bias
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 6.0, 5.72, None, 18)
    add_code(slide, code2.strip(), 6.78, 1.12, 5.95, 5.72, None, 18)
    note = "GDFN 先把通道扩展为两倍，再用深度卷积拆成两支，一支经过 GELU 后与另一支逐元素相乘，最后投影回原通道。LayerNorm 先把 BCHW 转换为 B(HW)C，按通道归一化，再恢复四维。Jittor 版本保留官方的 biased LayerNorm 公式。"
    register("GDFN 与 LayerNorm", [], ["可编辑代码：GDFN", "可编辑代码：WithBias LayerNorm"], note, code1 + "\n" + code2)

    slide = add_base_slide(prs, "强度损失与 Sobel 梯度损失", 18, "loss.py")
    code1 = """
# Intensity and Sobel loss
def execute(self, fused, ir, vis):
    vis_y = vis[:, :1, :, :]
    target = jt.maximum(vis_y, ir)
    loss_int = nn.l1_loss(
        target, fused
    )
    grad_vis = self.sobelconv(vis_y)
    grad_ir = self.sobelconv(ir)
    grad_fus = self.sobelconv(fused)
    grad_target = jt.maximum(
        grad_vis, grad_ir
    )
    loss_grad = nn.l1_loss(
        grad_target, grad_fus
    )
    return loss_int, loss_grad
"""
    code2 = """
# Fixed Sobel kernels
kernel_x = [
    [-1, 0, 1],
    [-2, 0, 2],
    [-1, 0, 1],
]
kernel_y = [
    [1, 2, 1],
    [0, 0, 0],
    [-1, -2, -1],
]

def execute(self, x):
    gx = nn.conv2d(
        x, self.weightx, padding=1
    )
    gy = nn.conv2d(
        x, self.weighty, padding=1
    )
    return jt.abs(gx) + jt.abs(gy)
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 6.0, 5.72, None, 18)
    add_code(slide, code2.strip(), 6.78, 1.12, 5.95, 5.72, None, 18)
    note = "强度损失先取红外与可见光亮度的逐像素最大值，再与融合图做 L1。Sobel 模块固定两组 3×3 核，分别计算水平和垂直梯度，取绝对值之和。融合图梯度需要接近两幅源图中较强的梯度。固定卷积核通过 stop_grad 排除在参数更新之外。"
    register("强度损失与 Sobel 梯度损失", [], ["可编辑代码：Fusionloss", "可编辑代码：Sobelxy"], note, code1 + "\n" + code2)

    slide = add_base_slide(prs, "Laplacian 梯度损失", 19, "loss.py")
    code1 = """
def laplacian(input, kernel_size):
    k = kernel_size
    kernel = np.ones((k, k), np.float32)
    center = k // 2
    kernel[center, center] = (
        1 - kernel.sum()
    )
    kernel = kernel / np.abs(kernel).sum()
    kernel = jt.array(kernel).float32()
    kernel = kernel.reshape(1, 1, k, k)
    channels = input.shape[1]
    kernel = kernel.repeat(
        channels, 1, 1, 1
    )
    input = nn.pad(
        input, (center,) * 4,
        mode='reflect'
    )
    return nn.conv2d(input, kernel,
                     groups=channels)
"""
    code2 = """
class JointGrad(nn.Module):
    def execute(
        self, im_fus, im_ir, im_vi
    ):
        ir_grad = jt.abs(laplacian(im_ir, 3))
        vi_grad = jt.abs(laplacian(im_vi, 3))
        fus_grad = laplacian(im_fus, 3)
        target = jt.where(
            ir_grad - vi_grad >= 0,
            laplacian(im_ir, 3),
            laplacian(im_vi, 3),
        )
        return nn.l1_loss(
            target, fus_grad
        )
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 6.0, 5.72, None, 18)
    add_code(slide, code2.strip(), 6.78, 1.12, 5.95, 5.72, None, 18)
    note = "官方 PyTorch 版本通过 Kornia 计算 Laplacian。Jittor 侧按照相同核定义手工构造卷积，并采用 reflect padding 和分组卷积。JointGrad 比较两幅源图 Laplacian 绝对值，逐位置选择更强者，再约束融合图的 Laplacian 与之接近。"
    register("Laplacian 梯度损失", [], ["可编辑代码：laplacian", "可编辑代码：JointGrad"], note, code1 + "\n" + code2)

    manual_code_slide(
        prs, 20, "数据加载与随机裁剪",
        """
class TrainLoader(Dataset):
    def __init__(self, ir_folder, vi_folder, patch_size=128):
        self.ps = patch_size
        self.ir_list = sorted(ir_folder.glob('*'))
        self.vi_list = sorted(vi_folder.glob('*'))
        self.set_attrs(total_len=len(self.ir_list))

    def get_patch(self, ir, vis):
        h, w = ir.shape[1], ir.shape[2]
        x = np.random.randint(10, h - 10 - self.ps + 1)
        y = np.random.randint(10, w - 10 - self.ps + 1)
        ir_crop = ir[:, x:x+self.ps, y:y+self.ps]
        vi_crop = vis[:, x:x+self.ps, y:y+self.ps]
        return ir_crop, vi_crop

    def __getitem__(self, index):
        ir = self.imread(self.ir_list[index])
        vi = self.imread(self.vi_list[index])
        return self.get_patch(ir, vi)
""",
        "训练集按文件名排序读取红外和可见光图像，并在同一坐标裁剪 128×128 patch。随机范围完全保留官方边界设置。Jittor Dataset 通过 set_attrs 设置总长度和批处理属性。",
        ["1083 MSRS + 200 RoadScene", "同一坐标裁剪两种模态", "输入归一化到 [0, 1]"], "train_loader.py", 9.25,
    )

    slide = add_base_slide(prs, "训练配置", 21, "args_SIBA.py · train.py")
    code1 = """
patch_size = 128
epochs = 60
optim_step = 25
batch_size = 4
init_lr = 1e-4
use_gpu = True
optim_gamma = 0.5
weight_decay = 0
"""
    code2 = """
model = SIBA()
optimizer = PyTorchAdam(
    model.parameters(), lr=1e-4,
    eps=1e-8,
    betas=(0.9, 0.999),
)
scheduler = StepLR(
    optimizer, step_size=25,
    gamma=0.5,
)
JointGradLoss = JointGrad()
Intensity_Grad = Fusionloss()

trainloader = TrainLoader(
    ir_path, vi_path, 128
)
trainloader = trainloader.set_attrs(
    batch_size=4, shuffle=True,
    drop_last=False, num_workers=0,
)
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 4.35, 5.72, None, 18)
    add_code(slide, code2.strip(), 5.13, 1.12, 7.6, 5.72, None, 18)
    note = "训练配置与官方仓库一致：60 轮、batch size 4、patch 128、初始学习率 1e-4，StepLR 每 25 轮衰减为原来的 0.5，weight decay 为 0。数据加载器开启 shuffle，drop_last 为 False。模型、损失和调度器的建立顺序保持一致。"
    register("训练配置", [], ["可编辑代码：args_SIBA.py", "可编辑代码：train.py 初始化"], note, code1 + "\n" + code2)

    manual_code_slide(
        prs, 22, "训练循环与权重保存",
        """
for epoch in range(num_epochs):
    for i, (ir, vi) in enumerate(trainloader):
        model.train()
        optimizer.zero_grad()
        fuse = model(ir, vi)
        loss_lap = JointGradLoss(fuse, ir, vi)
        loss_int, loss_sobel = Intensity_Grad(fuse, ir, vi)
        loss = 10 * loss_lap + 0.1 * loss_int + loss_sobel
        optimizer.backward(loss)
        clip_grad_norm_pytorch(
            optimizer, max_norm=0.01, norm_type=2
        )
        optimizer.step()
        if i % 50 == 0:
            print(epoch, optimizer.lr, loss.item())
    scheduler.step()
jt.save({'model': model.state_dict()}, save_path)
""",
        "每个 batch 依次完成前向、三项损失、反向传播、全局梯度裁剪和 Adam 更新。官方日志每 50 个 batch 打印一次，因此 60 轮共得到 420 条 loss。每轮结束后调用 StepLR，最终保存 epoch60 权重。",
        ["10×Laplacian + 0.1×Intensity + Sobel", "每 50 batch 记录一次", "保存 SIBA_epoch60.pkl"], "train.py", 9.25,
    )

    slide = add_base_slide(prs, "Adam 与梯度裁剪", 23, "compat")
    code1 = """
next_avg = (
    beta1 * exp_avg
    + (1 - beta1) * gradient
)
exp_avg.update(next_avg)
next_avg_sq = (
    beta2 * exp_avg_sq
    + (1 - beta2) * gradient * gradient
)
exp_avg_sq.update(next_avg_sq)
denominator = (
    jt.sqrt(exp_avg_sq) / bias2
    + self.eps
)
step_size = self.lr / bias1
parameter.update(
    parameter - step_size
    * exp_avg / denominator
)
"""
    code2 = """
gradients = (
    collect_optimizer_gradients(
        optimizer
    )
)
parameter_norms = jt.stack([
    jt.norm(gradient.flatten(), 2)
    for gradient in gradients
])
total_norm = jt.norm(
    parameter_norms.flatten(), 2
)
coefficient = jt.minimum(
    max_norm / (total_norm + 1e-6), 1.0
)
for gradient in gradients:
    gradient.update(
        gradient * coefficient
    )
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 6.0, 5.72, None, 18)
    add_code(slide, code2.strip(), 6.78, 1.12, 5.95, 5.72, None, 18)
    note = "Jittor 内置优化器的更新细节与官方 PyTorch Adam 不完全相同，因此这里显式保存一阶矩、二阶矩和 bias correction，按 PyTorch 公式更新参数。梯度裁剪同样按所有参数梯度的全局 L2 范数计算统一缩放系数。给定相同参考梯度时，这两个兼容实现能够与 PyTorch 对齐。"
    register("Adam 与梯度裁剪", [], ["可编辑代码：PyTorch Adam 关键更新", "可编辑代码：全局梯度裁剪"], note, code1 + "\n" + code2)

    slide = add_base_slide(prs, "测试与颜色重建", 24, "test.py · RGB2YCrCb.py")
    code1 = """
# test.py
with jt.no_grad():
    for batch in test_loader:
        _, vis_y, cb, cr, ir, name, _ = (
            batch
        )
        fused = model(ir, vis_y)
        fused = clamp(fused[0])
        fused = YCrCb2RGB(
            fused, cb[0], cr[0]
        )
        fused = fused.transpose(1, 2, 0)
        image = transform.ToPILImage()(fused)
        image.save(
            result_save_path + '/' + name[0]
        )
"""
    code2 = """
def YCrCb2RGB(Y, Cb, Cr):
    ycrcb = jt.concat([Y, Cr, Cb], dim=0)
    c, w, h = ycrcb.shape
    pixels = ycrcb.reshape(3, -1)
    pixels = pixels.transpose(0, 1)
    matrix = jt.array([
        [1.0, 1.0, 1.0],
        [1.403, -0.714, 0.0],
        [0.0, -0.344, 1.773],
    ]).cast(Y.dtype)
    bias = jt.array(
        [0.0, -0.5, -0.5]
    ).cast(Y.dtype)
    pixels = pixels + bias
    out = jt.matmul(pixels, matrix)
    out = out.transpose(0, 1)
    out = out.reshape(c, w, h)
    return clamp(out)
"""
    add_code(slide, code1.strip(), 0.58, 1.12, 6.65, 5.72, None, 18)
    add_code(slide, code2.strip(), 7.43, 1.12, 5.3, 5.72, None, 18)
    note = "测试阶段加载 Jittor 权重，逐张生成融合亮度 Y。可见光输入同时保留 Cb 和 Cr，融合结束后使用与官方代码相同的矩阵把 YCbCr 转回 RGB。输出先限制在 0 到 1，再保存为原文件名。MSRS、M3FD 和 TNO 都按这一流程完成推理。"
    register("测试与颜色重建", [], ["可编辑代码：测试循环", "可编辑代码：YCbCr 转 RGB"], note, code1 + "\n" + code2)


def build_result_slides(prs: Presentation):
    # 25 alignment
    slide = add_base_slide(prs, "前向与输出对齐", 25, "official checkpoint")
    rows = [
        ["检查项", "结果"],
        ["前向激活最大绝对误差", "2.0206 × 10⁻⁴"],
        ["总损失最大绝对误差", "2.9802 × 10⁻⁶"],
        ["梯度余弦相似度", "0.999945"],
        ["一步参数更新相对 L2 误差", "6.73%"],
    ]
    add_table(slide, rows, 0.7, 1.3, 6.15, 4.35, [3.9, 2.25], 18)
    rows2 = [
        ["数据集", "图像数", "最大像素差", "平均绝对差"],
        ["MSRS", "361", "1", "0.00284"],
        ["M3FD", "300", "1", "0.02474"],
        ["TNO", "45", "1", "0.00885"],
    ]
    add_table(slide, rows2, 7.15, 1.3, 5.5, 3.35, [1.4, 1.1, 1.45, 1.55], 18)
    add_text(slide, 7.22, 4.86, 5.25, 1.12, "发布权重下，706 张图的文件名和尺寸全部一致；\nJittor 与 PyTorch 最多相差 1 个 uint8 灰度级。", 19, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.82, 6.28, 11.55, 0.45, "前向和推理高度一致；反向更新仍存在可测差异，不把两条训练轨迹称为完全等价。", 19, RED, True, align=PP_ALIGN.CENTER)
    note = "这一页把推理对齐和训练对齐分开说明。固定相同权重和输入时，前向激活、总损失和梯度方向都非常接近。发布权重在 706 张测试图上最多只差 1 个灰度级。但使用各框架真实梯度执行一步更新时，相对 L2 误差为 6.73%，所以只能说训练功能完成并能收敛，不能说两框架的参数更新完全相同。"
    register("前向与输出对齐", [cell for row in rows for cell in row] + [cell for row in rows2 for cell in row] + ["发布权重下，706 张图的文件名和尺寸全部一致，Jittor 与 PyTorch 最多相差 1 个 uint8 灰度级。", "前向和推理高度一致；反向更新仍存在可测差异，因此不把两条训练轨迹称为完全等价。"], ["可编辑数值表"], note)

    # 26 training and performance
    slide = add_base_slide(prs, "训练与推理记录", 26, "RTX 3090 24GB")
    rows = [
        ["框架", "60 轮训练时间"],
        ["Jittor", "4,079 s（67.98 min）"],
        ["PyTorch", "2,221 s（37.02 min）"],
    ]
    add_table(slide, rows, 0.72, 1.25, 5.0, 2.3, [1.7, 3.3], 18)
    rows2 = [
        ["同步 FPS", "MSRS", "M3FD", "TNO"],
        ["Jittor", "9.12", "12.66", "6.93"],
        ["PyTorch", "11.99", "19.90", "12.31"],
    ]
    add_table(slide, rows2, 0.72, 4.0, 5.0, 2.05, [1.6, 1.1, 1.15, 1.15], 18)
    rows3 = [
        ["自训练", "MI", "Qabf", "SSIM"],
        ["MSRS · JT", "5.461", "0.710", "0.974"],
        ["MSRS · PT", "5.043", "0.710", "0.980"],
        ["M3FD · JT", "3.864", "0.630", "0.963"],
        ["M3FD · PT", "3.547", "0.649", "0.968"],
        ["TNO · JT", "3.334", "0.567", "0.948"],
        ["TNO · PT", "3.162", "0.577", "0.947"],
    ]
    add_table(slide, rows3, 6.15, 1.25, 6.5, 4.8, [2.2, 1.35, 1.35, 1.35], 18)
    add_text(slide, 6.25, 6.22, 6.15, 0.48, "同步计时用于框架对比；论文 TITAN RTX 的 FPS 不直接横向比较。", 18, RED, True, align=PP_ALIGN.CENTER)
    note = "两套代码都在 RTX 3090 上完成 60 轮训练。Jittor 用时约 68 分钟，PyTorch 约 37 分钟。同步推理计时下，Jittor 在三个数据集上都慢于 PyTorch。右侧列出自训练权重的三个代表指标，两条训练均收敛，但最终结果不完全相同。论文使用 TITAN RTX，并且原测试脚本没有同步 GPU，因此论文 FPS 只作为背景，不与这里直接比较。"
    register("训练与推理记录", [cell for row in rows for cell in row] + [cell for row in rows2 for cell in row] + [cell for row in rows3 for cell in row] + ["同步计时用于框架对比；论文 TITAN RTX 的 FPS 不直接横向比较。", "RTX 3090 24GB"], ["可编辑训练时间、同步 FPS 和自训练指标表"], note)

    # 27 loss curve
    slide = add_base_slide(prs, "训练曲线", 27, "60 epochs")
    add_picture_contain(slide, ASSETS / "loss_curve.png", 0.72, 1.1, 9.25, 5.72)
    add_text(slide, 10.2, 1.48, 2.55, 0.68, "Jittor 60 轮", 22, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 10.2, 2.18, 2.55, 0.68, "PyTorch 60 轮", 22, ORANGE, True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["两条曲线均下降并稳定", "官方每 50 batch 记录一次 loss", "shuffle 和反向差异使曲线不逐点重合"], 10.1, 3.18, 2.65, 2.55, 19)
    note = "训练曲线按官方日志规则绘制，每 50 个 batch 采样一次。两条曲线都在前 20 轮明显下降，后期稳定在较低范围，说明两套训练链路都能正常收敛。由于 Jittor 和 PyTorch 的 shuffle 实现不同，并且真实反向更新存在差异，所以这张图不用于证明逐 batch 一致，只用于验证完整训练过程。"
    register("训练曲线", ["Jittor 60 轮", "PyTorch 60 轮", "两条曲线均下降并稳定", "官方每 50 batch 记录一次 loss", "shuffle 和反向差异使曲线不逐点重合", "60 epochs · every 50 batches"], ["真实 60 轮 loss 曲线：results/training_analysis_20260727_siba_official_protocol/loss_curve.png"], note)

    # 28 metrics
    slide = add_base_slide(prs, "论文指标复现", 28, "released checkpoint")
    rows = [
        ["数据集 / 实现", "VIF", "SCD", "MI", "Qabf", "SSIM", "MS-SSIM", "FMI"],
        ["MSRS · Paper", "1.061", "1.700", "5.111", "0.715", "0.981", "0.975", "0.933"],
        ["MSRS · Jittor", "1.061", "1.700", "5.111", "0.715", "0.981", "0.975", "0.933"],
        ["M3FD · Paper", "0.759", "1.660", "3.771", "0.654", "0.964", "0.923", "0.883"],
        ["M3FD · Jittor", "0.759", "1.660", "3.771", "0.654", "0.964", "0.923", "0.883"],
        ["TNO · Paper", "0.836", "1.724", "3.508", "0.588", "0.932", "0.904", "0.914"],
        ["TNO · Jittor", "0.836", "1.725", "3.507", "0.588", "0.932", "0.904", "0.914"],
    ]
    add_table(slide, rows, 0.52, 1.22, 12.3, 4.75, [2.45, 1, 1, 1, 1, 1, 1.15, 1], 18)
    add_text(slide, 0.8, 6.2, 11.8, 0.52, "发布权重下，21 个论文指标按三位小数复现；TNO 的 SCD、MI 差异小于 0.001。", 21, BLUE, True, align=PP_ALIGN.CENTER)
    note = "这一页使用作者发布权重、完整测试集和论文对应的七项指标。MSRS 和 M3FD 的所有指标按三位小数与论文一致。TNO 的 SCD 和 MI 在未四舍五入前有很小差异，显示到三位小数时 SCD 为 1.725、MI 为 3.507，而论文为 1.724 和 3.508，差值都小于 0.001。其余指标一致。"
    register("论文指标复现", [cell for row in rows for cell in row] + ["发布权重下，21 个论文指标按三位小数复现；TNO 的 SCD、MI 差异小于 0.001。", "released checkpoint · 7 metrics"], ["可编辑论文指标表"], note)

    # 29 qualitative
    slide = add_base_slide(prs, "融合结果", 29, "official checkpoint")
    headers = ["Infrared", "Visible", "Jittor", "PyTorch"]
    for index, header in enumerate(headers):
        add_text(slide, 1.35 + index * 2.82, 1.08, 2.45, 0.36, header, 18, TEXT, True, align=PP_ALIGN.CENTER)
    cases = [
        ("MSRS", "00004N.png"),
        ("M3FD", "00000.png"),
        ("TNO", "01.png"),
    ]
    for row, (dataset, name) in enumerate(cases):
        key = "M3FD_2x" if dataset == "M3FD" else dataset
        y = 1.46 + row * 1.72
        add_text(slide, 0.48, y + 0.54, 0.9, 0.4, dataset, 18, TEXT, True, align=PP_ALIGN.CENTER)
        paths = [
            ROOT / "datasets" / "test" / key / "ir" / name,
            ROOT / "datasets" / "test" / key / "vi" / name,
            ROOT / "results" / "official_checkpoint_alignment_20260727_siba_official_protocol" / "jittor" / key / name,
            ROOT / "results" / "official_checkpoint_alignment_20260727_siba_official_protocol" / "pytorch" / key / name,
        ]
        for col, path in enumerate(paths):
            add_picture_contain(slide, path, 1.32 + col * 2.82, y, 2.48, 1.5)
    add_text(slide, 0.82, 6.64, 11.8, 0.32, "Jittor 与 PyTorch 结果在视觉上保持一致；红外目标与可见光背景均被保留。", 20, BLUE, True, align=PP_ALIGN.CENTER)
    note = "这里各选取 MSRS、M3FD 和 TNO 的一组真实测试图。前两列是红外和可见光输入，后两列分别是 Jittor 与 PyTorch 在作者发布权重下的输出。两种框架的结果在视觉上基本无法区分，同时保留了红外目标和可见光背景细节。所有图像都来自完整推理结果目录。"
    register("融合结果", headers + [x for case in cases for x in case] + ["Jittor 与 PyTorch 结果在视觉上保持一致；红外目标与可见光背景均被保留。", "official checkpoint"], [f"{dataset}：{name} 的红外、可见光、Jittor 和 PyTorch 输出" for dataset, name in cases], note)

    # 30 conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, ASSETS / "cover_gradient.png", 0, 0, SLIDE_W, SLIDE_H)
    add_text(slide, 0.72, 0.5, 7.0, 0.58, "复现结论", 34, rgb((57, 45, 67)), True)
    add_bullets(slide, [
        "13 个官方 Python 文件均完成 Jittor 对应实现",
        "完整运行 60 轮训练，并完成三个测试集评测",
        "发布权重下，706 张输出与 PyTorch 最多相差 1 个灰度级",
        "训练均收敛；一步参数更新相对 L2 差异为 6.73%",
    ], 0.82, 1.35, 11.75, 3.05, 20, rgb((65, 58, 70)))
    add_text(slide, 0.86, 4.66, 11.6, 0.42, "训练数据限制：RoadScene 200 对名单与随机种子未公开。", 18, rgb((80, 72, 86)))
    add_text(slide, 0.86, 5.22, 11.6, 0.45, "差异定位：卷积、归一化和自动微分的反向计算。", 22, RED, True)
    add_text(slide, 0.86, 6.0, 11.6, 0.42, "官方代码：https://github.com/Afreshbird/SIBA", 18, BLUE)
    add_text(slide, 0.86, 6.45, 11.6, 0.42, "Jittor 代码：待补充个人 GitHub 仓库地址", 18, BLUE)
    add_text(slide, 9.65, 6.86, 2.65, 0.42, "感谢各位老师！", 24, rgb((57, 45, 67)), True, align=PP_ALIGN.CENTER)
    add_text(slide, 12.45, 7.08, 0.35, 0.3, "30", 18, rgb((100, 90, 107)), align=PP_ALIGN.RIGHT)
    note = "最后总结。代码结构、前向、训练、测试和评估功能都已完成 Jittor 迁移，作者发布权重下的推理结果高度一致。需要明确的是，自训练权重不能宣称与 PyTorch 完全等价，原因包括真实反向更新差异，以及作者没有公开 RoadScene 200 对样本的具体名单。后续如果继续优化，重点应放在卷积、归一化和自动微分的反向差异定位。个人 GitHub 地址创建后再补到这一页。"
    register("复现结论", ["复现结论", "13 个官方 Python 文件均完成 Jittor 对应实现", "完整运行 60 轮训练，并完成三个测试集评测", "发布权重下，706 张输出与 PyTorch 最多相差 1 个灰度级", "训练均收敛；一步参数更新相对 L2 差异为 6.73%", "训练数据限制：RoadScene 200 对名单与随机种子未公开。", "差异定位：卷积、归一化和自动微分的反向计算。", "官方代码：https://github.com/Afreshbird/SIBA", "Jittor 代码：待补充个人 GitHub 仓库地址", "感谢各位老师！"], ["淡粉紫渐变背景"], note)


def write_analysis():
    content = """# 优秀 PPT 拆解与重做原则

## 1. 范例如何组织内容

范例没有从宽泛背景开始，而是直接回答三个问题：方法新在哪里、网络怎样组成、代码怎样实现。前几页用结构图建立整体认识，随后严格按照程序执行顺序展开代码，最后展示运行和结果。听众看到的顺序与模型执行顺序一致，因此不需要在不同概念之间来回切换。

1. 封面：论文题目、四项内容、汇报人。
2. 创新点：用一张对照图说明相对现有方法的变化。
3. 基本架构：先展示完整网络，不在这一页进入代码细节。
4. 核心模块：空间分支、频域分支和最终融合分别讲解。
5. 顶层模型：先看各模块怎样连接，再进入子模块。
6. 子模块：沿前向传播顺序依次展示关键函数。
7. 损失函数：每个损失独立成页，代码与含义一一对应。
8. 运行演示：直接展示环境、命令和终端输出。
9. 实验结果：用真实日志、曲线和结果图结束。

## 2. 范例的讲解节奏

- 方法部分约占前 6 分钟，只讲创新、总体结构和必要模块。
- 代码部分占主体，先顶层、后分支、再损失，不按文件名机械罗列。
- 同一页会停留较长时间，讲解者用鼠标指出当前代码行；页面不频繁切换。
- 一张代码页只承担一个问题，例如“空间分支如何计算”或“损失如何组成”。
- 公式只用于说明代码含义，不在主页面展开推导。
- 运行和结果放在最后，使前面的代码讲解有明确落点。

## 3. 范例的页面结构

- 内容页为米白背景，左上放黑色粗体标题。
- 页面主体通常只有一个大图或一个浅灰代码区。
- 代码使用语法着色，但不使用额外边框、阴影和装饰图标。
- 解释文字很少，通常只保留模块名、公式或两三条短句。
- 论文图直接放大，不把整页论文截图放进幻灯片。
- 页面留白较多，标题、主体和页脚位置固定。
- 封面使用淡粉紫渐变，内容页不沿用大面积渐变。

## 4. 主页面与讲稿的分工

范例的主页面只呈现观众需要看的内容。为什么这样安排、这一页准备讲多久、下一页怎样过渡，都没有写在幻灯片上。代码解释主要由讲解者完成，主页面只保留代码本身。重复说明、完整定义和补充限制放在口头讲解或备注中。

不应出现在主页面的文字包括：

- 汇报主线、讲解顺序、这一页说明什么。
- “显著提升”“全面突破”等没有数据支撑的评价。
- 对观众说的话和演讲者的操作提示。
- 可以由图、代码或表格直接表达的大段描述。

## 5. SIBA 的对应讲法

- 第 2 页用论文 Fig.1 的源图和 Grad-CAM 观察说明选题出发点。
- 第 3 页用 Fig.2 建立完整网络结构。
- 第 4–5 页只讲 CBSM 和 I-SCA/V-SCA 两个核心模块。
- 第 6 页收束四路特征、输出层和三项损失。
- 第 7–10 页先讲 `models/SIBA.py`，对应模型建立和完整前向过程。
- 第 11–19 页按实际调用关系进入 SE-ResNet、CBSM、Restormer、注意力和损失。
- 第 20–24 页继续讲数据、训练配置、训练循环、优化器兼容和测试。
- 第 25–29 页依次给出数值对齐、训练记录、曲线、论文指标和融合图。
- 第 30 页只保留复现结论、限制和代码链接。

## 6. 本版排版约束

- 中文统一使用 Microsoft YaHei；代码使用 Consolas；公式使用 Cambria Math。
- 所有主页面显式文字不小于 18 pt。
- 每页只保留一个中心内容，不使用文字墙。
- 不使用圆角卡片式摘要，不堆叠多个装饰框。
- 论文图只保留图本体，不保留论文图注和周围正文。
- 代码采用可编辑文本，过长代码通过节选和换行解决，不缩小字号。
- 所有讲稿写入演讲者备注，共 30 页备注。
- 实验页只使用真实日志、真实指标和真实输出图。
- 对训练差异如实说明，不把“均能收敛”写成“训练完全等价”。
"""
    ANALYSIS_PATH.write_text(content, encoding="utf-8")


def write_outline():
    lines = ["# SIBA-Jittor PPT 完整大纲（重做版）", "", "说明：以下逐页列出主页面出现的全部文字、图片/代码和演讲者备注。", ""]
    for index, item in enumerate(META, start=1):
        lines.append(f"## 第 {index} 页：{item['title']}")
        lines.append("")
        lines.append("### 主页面文字")
        for text in item["visible"]:
            lines.append(f"- {text}")
        if not item["visible"]:
            lines.append("- 无额外说明文字；页面主体为代码。")
        lines.append("")
        lines.append("### 图片 / 图表 / 代码")
        for image in item["images"]:
            lines.append(f"- {image}")
        if item.get("code"):
            lines.append("")
            lines.append("```python")
            lines.append(item["code"])
            lines.append("```")
        lines.append("")
        lines.append("### 演讲者备注")
        lines.append(item["note"])
        lines.append("")
    OUTLINE_PATH.write_text("\n".join(lines), encoding="utf-8")


def add_notes_and_export():
    import win32com.client

    app = win32com.client.DispatchEx("PowerPoint.Application")
    try:
        presentation = app.Presentations.Open(str(PPTX_PATH), WithWindow=False)
        for index, note in enumerate(NOTES, start=1):
            slide = presentation.Slides(index)
            for shape_index in range(1, slide.NotesPage.Shapes.Count + 1):
                shape = slide.NotesPage.Shapes(shape_index)
                try:
                    placeholder_type = shape.PlaceholderFormat.Type
                except Exception:
                    continue
                if placeholder_type == 2:
                    shape.TextFrame.TextRange.Text = note
                    break
        presentation.Save()
        if PDF_PATH.exists():
            PDF_PATH.unlink()
        presentation.SaveAs(str(PDF_PATH), 32)
        if PREVIEW_DIR.exists():
            shutil.rmtree(PREVIEW_DIR)
        PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
        presentation.Export(str(PREVIEW_DIR), "PNG", 1600, 900)
        presentation.Save()
        presentation.Close()
    finally:
        app.Quit()


def make_contact_sheet():
    slide_paths = sorted(PREVIEW_DIR.glob("*.PNG"), key=lambda path: int("".join(ch for ch in path.stem if ch.isdigit())))
    thumbs = []
    for path in slide_paths:
        image = Image.open(path).convert("RGB")
        image.thumbnail((420, 236), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (430, 266), "white")
        canvas.paste(image, ((430 - image.width) // 2, 4))
        draw = ImageDraw.Draw(canvas)
        draw.text((10, 242), path.stem, fill=(40, 40, 40))
        thumbs.append(canvas)
    columns = 3
    rows = math.ceil(len(thumbs) / columns)
    sheet = Image.new("RGB", (columns * 430, rows * 266), (232, 232, 232))
    for index, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((index % columns) * 430, (index // columns) * 266))
    sheet.save(CONTACT_PATH)


def audit_pptx():
    presentation = Presentation(PPTX_PATH)
    font_sizes = []
    out_of_bounds = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > presentation.slide_width or shape.top + shape.height > presentation.slide_height:
                intentional_cover = slide_index in (1, 30) and shape.name.startswith("Picture")
                if not intentional_cover:
                    out_of_bounds.append(f"slide {slide_index}: {shape.name}")
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and run.font.size is not None:
                        font_sizes.append(run.font.size.pt)
    min_font = min(font_sizes) if font_sizes else None
    text_overflow = []
    try:
        import win32com.client

        app = win32com.client.DispatchEx("PowerPoint.Application")
        com_presentation = app.Presentations.Open(str(PPTX_PATH), WithWindow=False)
        try:
            for slide_index in range(1, com_presentation.Slides.Count + 1):
                slide = com_presentation.Slides(slide_index)
                for shape_index in range(1, slide.Shapes.Count + 1):
                    shape = slide.Shapes(shape_index)
                    try:
                        if shape.HasTextFrame and shape.TextFrame.HasText and shape.TextFrame2.Overflowing:
                            text_overflow.append(f"slide {slide_index}: {shape.Name}")
                    except Exception:
                        continue
        finally:
            com_presentation.Close()
            app.Quit()
    except Exception as error:
        text_overflow.append(f"检查未完成：{error}")
    report = [
        "# PPT 检查报告",
        "",
        f"- 幻灯片数量：{len(presentation.slides)}",
        f"- 可检测到的最小字号：{min_font:.1f} pt" if min_font is not None else "- 未检测到显式字号",
        f"- 越界对象数量：{len(out_of_bounds)}",
        f"- 文字溢出数量：{len(text_overflow)}",
        f"- 演讲者备注数量：{len(NOTES)}",
        f"- PDF：{'已生成' if PDF_PATH.exists() else '未生成'}",
        f"- 逐页 PNG：{len(list(PREVIEW_DIR.glob('*.PNG')))} 张",
        "",
    ]
    if out_of_bounds:
        report.append("## 越界对象")
        report.extend(f"- {item}" for item in out_of_bounds)
    else:
        report.append("- 未发现页面边界外对象。")
    if text_overflow:
        report.append("")
        report.append("## 文字溢出")
        report.extend(f"- {item}" for item in text_overflow)
    else:
        report.append("- 未发现文字溢出。")
    CHECK_PATH.write_text("\n".join(report), encoding="utf-8-sig")


def build_presentation():
    prepare_assets()
    DELIVERABLE.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    build_cover(presentation)
    build_method_slides(presentation)
    build_code_slides_v2(presentation)
    build_result_slides(presentation)
    presentation.save(PPTX_PATH)
    write_analysis()
    write_outline()
    add_notes_and_export()
    make_contact_sheet()
    audit_pptx()
    print(PPTX_PATH)
    print(PDF_PATH)
    print(CONTACT_PATH)


if __name__ == "__main__":
    build_presentation()
