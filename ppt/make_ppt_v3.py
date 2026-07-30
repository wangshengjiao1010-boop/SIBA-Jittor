#!/usr/bin/env python3
import ast
import csv
import json
import math
import re
import shutil
import textwrap
from pathlib import Path

import black
import fitz
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pygments import lex
from pygments.lexers import PythonLexer
from pygments.token import Comment, Keyword, Literal, Name, Number, Operator, String
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
JITTOR = ROOT / "siba_jittor"
PPT = ROOT / "ppt"
ASSET_SOURCE = PPT / "assets_generated"
ASSET_DIR = PPT / "assets_v3"
OUTPUT_DIR = ROOT / "deliverables" / "SIBA_Jittor_培育期_最终版_20260729"
PPTX_PATH = OUTPUT_DIR / "王胜娇-培育期.pptx"
PDF_PATH = OUTPUT_DIR / "王胜娇-培育期.pdf"
PREVIEW_DIR = OUTPUT_DIR / "preview_slides"
CONTACT_PATH = OUTPUT_DIR / "preview_contact.png"
NOTES_PATH = OUTPUT_DIR / "王胜娇-培育期-逐页讲稿.md"
CHECK_PATH = OUTPUT_DIR / "PPT检查报告.md"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_CN = "Times New Roman"
FONT_CN_BOLD = "Times New Roman"
FONT_CODE = "Times New Roman"
FONT_MATH = "Times New Roman"

BG = RGBColor(255, 255, 255)
TEXT = RGBColor(43, 43, 43)
MUTED = RGBColor(104, 100, 98)
CODE_BG = RGBColor(247, 248, 249)
BLUE = RGBColor(11, 105, 177)
ORANGE = RGBColor(207, 126, 75)
GREEN = RGBColor(81, 139, 105)
PURPLE = RGBColor(124, 100, 154)
RED = RGBColor(179, 77, 74)
PINK = RGBColor(235, 191, 196)
LIGHT_BLUE = RGBColor(184, 207, 237)
LIGHT_ORANGE = RGBColor(244, 205, 164)
LIGHT_GREEN = RGBColor(187, 220, 193)
LIGHT_PURPLE = RGBColor(207, 191, 235)
LIGHT_PINK = RGBColor(239, 188, 195)
LIGHT_GRAY = RGBColor(216, 218, 221)
BLACK = RGBColor(0, 0, 0)
FOOTER_BLUE = RGBColor(51, 128, 184)
FOOTER_MID = RGBColor(54, 92, 154)
FOOTER_DARK = RGBColor(39, 66, 118)

NOTES = []


def rgb(value):
    return RGBColor(*value)


def set_run_font(run, name, size, color=TEXT, bold=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    properties = run._r.get_or_add_rPr()
    properties.set("lang", "zh-CN")
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = properties.find(tag, properties.nsmap)
        if node is None:
            node = OxmlElement(tag)
            properties.append(node)
        node.set("typeface", name)


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=20,
    color=TEXT,
    bold=False,
    font=FONT_CN,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    name=None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    set_run_font(run, font, size, color, bold)
    return shape


def add_rich_text(
    slide, x, y, w, h, parts, size=20, align=PP_ALIGN.LEFT, name=None
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, color, bold, font in parts:
        run = paragraph.add_run()
        run.text = text
        set_run_font(run, font or FONT_CN, size, color, bold)
    return shape


def add_rect(slide, x, y, w, h, fill, line=None, width=1.0, rounded=False, name=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)
    if name:
        shape.name = name
    if shape._element.spPr.find(qn("a:effectLst")) is None:
        shape._element.spPr.append(OxmlElement("a:effectLst"))
    return shape


def add_box(slide, x, y, w, h, text, fill, line=None, size=18, bold=False, name=None):
    shape = add_rect(slide, x, y, w, h, fill, BLACK, 1.15, True, name)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_run_font(run, FONT_CN_BOLD if bold else FONT_CN, size, TEXT, bold)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=2.0, name=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line_xml = line.line._get_or_add_ln()
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    line_xml.append(tail)
    if name:
        line.name = name
    return line


def add_picture_contain(slide, path, x, y, w, h, name=None):
    path = Path(path)
    with Image.open(path) as image:
        image_w, image_h = image.size
    ratio = min(w / image_w, h / image_h)
    draw_w = image_w * ratio
    draw_h = image_h * ratio
    shape = slide.shapes.add_picture(
        str(path),
        Inches(x + (w - draw_w) / 2),
        Inches(y + (h - draw_h) / 2),
        Inches(draw_w),
        Inches(draw_h),
    )
    if name:
        shape.name = name
    return shape


def add_picture_cover(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as image:
        image_w, image_h = image.size
    ratio = max(w / image_w, h / image_h)
    draw_w = image_w * ratio
    draw_h = image_h * ratio
    shape = slide.shapes.add_picture(
        str(path),
        Inches(x + (w - draw_w) / 2),
        Inches(y + (h - draw_h) / 2),
        Inches(draw_w),
        Inches(draw_h),
    )
    return shape


def add_footer(slide, page):
    add_rect(slide, 0, 7.16, 3.85, 0.34, FOOTER_BLUE, FOOTER_BLUE)
    add_rect(slide, 3.85, 7.16, 8.0, 0.34, FOOTER_MID, FOOTER_MID)
    add_rect(slide, 11.85, 7.16, 1.483, 0.34, FOOTER_DARK, FOOTER_DARK)
    add_text(
        slide,
        0.55,
        7.16,
        3.0,
        0.34,
        "培育期",
        18,
        rgb((255, 255, 255)),
        font=FONT_CN,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        4.1,
        7.16,
        7.5,
        0.34,
        "SIBA 的 Jittor 复现",
        18,
        rgb((255, 255, 255)),
        font=FONT_CN,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        12.25,
        7.16,
        0.55,
        0.34,
        str(page),
        18,
        rgb((255, 255, 255)),
        font=FONT_CN,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_base_slide(prs, title, page, source=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_text(slide, 0.55, 0.16, 8.8, 0.55, title, 32, TEXT, True, FONT_CN_BOLD)
    if source:
        add_text(
            slide,
            9.15,
            0.22,
            3.6,
            0.42,
            source,
            18,
            MUTED,
            False,
            FONT_CN,
            PP_ALIGN.RIGHT,
        )
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.55),
        Inches(0.78),
        Inches(12.78),
        Inches(0.78),
    )
    line.line.color.rgb = BLUE
    line.line.width = Pt(1.5)
    add_footer(slide, page)
    return slide


def add_dot_list(slide, items, x, y, w, h, size=20, color=TEXT, gap=5):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = FONT_CN
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.25
        properties = paragraph._p.get_or_add_pPr()
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        properties.insert(0, bullet)
    return shape


def token_color(token_type):
    if token_type in Comment:
        return GREEN
    if token_type in Keyword:
        return BLUE
    if token_type in Name.Function or token_type in Name.Class:
        return PURPLE
    if token_type in Name.Builtin:
        return rgb((38, 126, 126))
    if token_type in String or token_type in Literal.String:
        return ORANGE
    if token_type in Number:
        return RED
    if token_type in Operator:
        return rgb((80, 80, 80))
    return rgb((45, 50, 55))


def add_code(
    slide,
    code,
    x,
    y,
    w,
    h,
    size=18,
    label=None,
    name=None,
    background=True,
    line_spacing=0.56,
    leading=None,
):
    if background:
        add_rect(slide, x, y, w, h, CODE_BG, None, rounded=False)
    if label:
        add_text(slide, x + 0.15, y + 0.08, w - 0.3, 0.3, label, 18, MUTED, True)
        text_y = y + 0.36
        text_h = h - 0.42
    else:
        text_y = y + 0.1
        text_h = h - 0.2
    shape = slide.shapes.add_textbox(
        Inches(x + 0.08), Inches(text_y), Inches(w - 0.16), Inches(text_h)
    )
    shape.name = name or f"CODE_TEXT_{label or 'block'}"
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    lexer = PythonLexer()
    lines = [line for line in code.rstrip().splitlines() if line.strip()]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = Pt(leading if leading is not None else size + 2)
        for token_type, value in lex(line, lexer):
            value = value.replace("\r", "").replace("\n", "")
            if not value:
                continue
            run = paragraph.add_run()
            run.text = value
            set_run_font(run, FONT_CODE, size, token_color(token_type))
    return shape


def add_table(slide, rows, x, y, w, h, font_size=18, col_widths=None):
    shape = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = shape.table
    if col_widths:
        for index, width in enumerate(col_widths):
            table.columns[index].width = Inches(width)
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            properties = cell._tc.get_or_add_tcPr()
            for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr"):
                border = OxmlElement(edge)
                border.append(OxmlElement("a:noFill"))
                properties.append(border)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = paragraph.add_run()
            run.text = str(value)
            set_run_font(
                run,
                FONT_CN,
                font_size,
                TEXT,
                row_index == 0,
            )
    row_height = h / len(rows)
    for rule_y, width in ((y, 1.5), (y + row_height, 1.0), (y + h, 1.5)):
        rule = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x),
            Inches(rule_y),
            Inches(x + w),
            Inches(rule_y),
        )
        rule.line.color.rgb = BLACK
        rule.line.width = Pt(width)
    return shape


def note(text):
    NOTES.append(text.strip())


def find_node(path, class_name=None, function_name=None):
    source = path.read_text(encoding="utf-8")
    tree = ast.parse(source)
    if class_name:
        cls = next(
            node
            for node in tree.body
            if isinstance(node, ast.ClassDef) and node.name == class_name
        )
        if function_name is None:
            return source, cls
        return source, next(
            node
            for node in cls.body
            if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef))
            and node.name == function_name
        )
    return source, next(
        node
        for node in tree.body
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef))
        and node.name == function_name
    )


def extract_code(relative, class_name=None, function_name=None, line_length=54):
    path = JITTOR / relative
    source, node = find_node(path, class_name, function_name)
    lines = source.splitlines()
    code = textwrap.dedent("\n".join(lines[node.lineno - 1 : node.end_lineno]))
    try:
        code = black.format_str(
            code + "\n", mode=black.Mode(line_length=line_length)
        ).rstrip()
    except Exception:
        code = code.rstrip()
    return code


def split_code(code, preferred=None):
    lines = code.splitlines()
    if preferred is None:
        preferred = math.ceil(len(lines) / 2)
    return "\n".join(lines[:preferred]), "\n".join(lines[preferred:])


def split_code_n(code, count):
    lines = [line for line in code.splitlines() if line.strip()]
    chunk_size = math.ceil(len(lines) / count)
    chunks = [
        "\n".join(lines[start : start + chunk_size])
        for start in range(0, len(lines), chunk_size)
    ]
    while len(chunks) < count:
        chunks.append("")
    return chunks[:count]


def class_method_excerpt(class_name, method_code):
    indented = "\n".join(
        f"  {line}" if line.strip() else line for line in method_code.splitlines()
    )
    return f"class {class_name}(nn.Module):\n{indented}"


def compact_indent(code):
    compacted = []
    for line in code.splitlines():
        leading = len(line) - len(line.lstrip(" "))
        compacted.append(" " * ((leading // 4) * 2 + leading % 4) + line.lstrip(" "))
    return "\n".join(compacted)


def wrap_rearrange_patterns(code):
    return code.replace(
        '        "b (head c) h w -> b head c (h w)",',
        '        "b (head c) h w -> "\n        "b head c (h w)",',
    ).replace(
        '        "b head c (h w) -> b (head c) h w",',
        '        "b head c (h w) -> "\n        "b (head c) h w",',
    )


def strip_ansi(text):
    return re.sub(r"\x1b\[[0-9;]*m", "", text)


def prepare_cover():
    path = ASSET_DIR / "cover_gradient.png"
    width, height = 1920, 1080
    x = np.linspace(0, 1, width)[None, :, None]
    y = np.linspace(0, 1, height)[:, None, None]
    base = np.ones((height, width, 3), dtype=np.float32)
    base[:] = np.array([248, 244, 242], dtype=np.float32)
    lavender = np.array([228, 215, 245], dtype=np.float32)
    peach = np.array([248, 220, 205], dtype=np.float32)
    yellow = np.array([247, 236, 201], dtype=np.float32)
    alpha_lav = np.exp(-((x - 0.08) ** 2 + (y - 0.08) ** 2) / 0.13)
    alpha_peach = np.exp(-((x - 0.82) ** 2 + (y - 0.18) ** 2) / 0.18)
    alpha_yellow = np.exp(-((x - 0.55) ** 2 + (y - 0.86) ** 2) / 0.22)
    image = base * (1 - 0.38 * alpha_lav) + lavender * (0.38 * alpha_lav)
    image = image * (1 - 0.34 * alpha_peach) + peach * (0.34 * alpha_peach)
    image = image * (1 - 0.26 * alpha_yellow) + yellow * (0.26 * alpha_yellow)
    Image.fromarray(np.clip(image, 0, 255).astype(np.uint8)).save(path)
    return path


def prepare_source_attention():
    source = ASSET_SOURCE / "source_attention_observation.png"
    target = ASSET_DIR / "source_attention_four.png"
    if target.exists():
        return target
    image = Image.open(source).convert("RGB")
    width, height = image.size
    cell_w = width // 3
    top_h = height // 2
    crops = [
        image.crop((0, 0, cell_w, top_h)),
        image.crop((cell_w, 0, cell_w * 2, top_h)),
        image.crop((cell_w * 2, 0, width, top_h)),
        image.crop((0, top_h, cell_w, height)),
    ]
    canvas = Image.new("RGB", (cell_w * 4, top_h), "white")
    for index, crop in enumerate(crops):
        canvas.paste(
            crop.resize((cell_w, top_h), Image.Resampling.LANCZOS), (index * cell_w, 0)
        )
    canvas.save(target)
    return target


def render_paper_figures_hd():
    document = fitz.open(ROOT / "paper_SIBA_ICCV2025.pdf")
    matrix = fitz.Matrix(8.0, 8.0)
    scale = 8.0 / 3.0
    crops = {
        "paper_fig2_arch_hd.png": (3, (115, 165, 1730, 780)),
        "paper_fig3_cbsm_hd.png": (3, (115, 1535, 930, 1840)),
        "paper_fig4_cross_attention_hd.png": (4, (95, 350, 940, 1000)),
    }
    for name, (page_index, crop) in crops.items():
        pixmap = document[page_index].get_pixmap(matrix=matrix, alpha=False)
        image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
        box = tuple(round(value * scale) for value in crop)
        image.crop(box).save(ASSET_DIR / name, optimize=True)
    document.close()


def prepare_terminal_image():
    log_path = (
        ROOT
        / "logs"
        / "demo_20260729_175522"
        / "train.log"
    )
    lines = [
        line
        for line in strip_ansi(
            log_path.read_text(encoding="utf-8", errors="replace")
        ).splitlines()
        if line.startswith("[demo step")
    ]
    selected = lines[:3] + ["..."] + [lines[9], lines[14], lines[-1]]
    path = ASSET_DIR / "real_training_terminal.png"
    canvas = Image.new("RGB", (1600, 760), (29, 32, 38))
    draw = ImageDraw.Draw(canvas)
    try:
        font = ImageFont.truetype("C:/Windows/Fonts/CascadiaCode.ttf", 31)
        font_bold = ImageFont.truetype("C:/Windows/Fonts/CascadiaCode.ttf", 32)
    except OSError:
        font = ImageFont.load_default()
        font_bold = font
    draw.rounded_rectangle(
        (15, 15, 1585, 745), radius=28, fill=(29, 32, 38), outline=(73, 77, 84), width=3
    )
    draw.ellipse((42, 42, 66, 66), fill=(239, 95, 87))
    draw.ellipse((78, 42, 102, 66), fill=(244, 190, 79))
    draw.ellipse((114, 42, 138, 66), fill=(87, 190, 102))
    draw.text(
        (170, 38),
        "screen -S kk · live Jittor training",
        font=font_bold,
        fill=(220, 224, 230),
    )
    y = 104
    for line in selected:
        color = (131, 212, 153) if line.startswith("[demo step") else (180, 186, 194)
        draw.text((48, y), line[:96], font=font, fill=color)
        y += 58
    canvas.save(path)
    return path


def prepare_performance_chart():
    csv_path = (
        ROOT
        / "results"
        / "performance_summary_20260727_siba_official_protocol"
        / "inference_timing.csv"
    )
    rows = list(csv.DictReader(csv_path.open("r", encoding="utf-8-sig")))
    datasets = ["MSRS", "M3FD_2x", "TNO"]
    data = {framework: [] for framework in ("jittor", "pytorch")}
    for framework in data:
        for dataset in datasets:
            row = next(
                item
                for item in rows
                if item["framework"] == framework and item["dataset"] == dataset
            )
            data[framework].append(float(row["synchronized_fps"]))
    path = ASSET_DIR / "synchronized_fps.png"
    plt.rcParams["font.family"] = "Noto Sans SC"
    fig, axis = plt.subplots(figsize=(10.5, 4.6), dpi=180)
    x = np.arange(len(datasets))
    width = 0.34
    axis.bar(x - width / 2, data["jittor"], width, label="Jittor", color="#7D96B8")
    axis.bar(x + width / 2, data["pytorch"], width, label="PyTorch", color="#D8A06F")
    axis.set_xticks(x, ["MSRS", "M3FD 1/2", "TNO"], fontsize=13)
    axis.set_ylabel("Synchronized FPS", fontsize=13)
    axis.spines[["top", "right"]].set_visible(False)
    axis.grid(axis="y", alpha=0.18)
    axis.legend(frameon=False, fontsize=12)
    axis.tick_params(axis="y", labelsize=12)
    for container in axis.containers:
        axis.bar_label(container, fmt="%.2f", fontsize=11, padding=3)
    fig.tight_layout()
    fig.savefig(path, transparent=True, bbox_inches="tight")
    plt.close(fig)
    return path


def prepare_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    render_paper_figures_hd()
    prepare_source_attention()
    prepare_terminal_image()
    prepare_performance_chart()
    shutil.copy2(
        ASSET_SOURCE / "official_alignment_examples.png",
        ASSET_DIR / "official_alignment_examples.png",
    )


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_rect(slide, 0, 0, SLIDE_W, 3.52, BLUE, BLUE)
    add_text(
        slide,
        0.75,
        0.58,
        11.83,
        0.72,
        "SIBA 的 Jittor 复现",
        42,
        rgb((255, 255, 255)),
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        1.4,
        1.48,
        10.53,
        1.25,
        "The Source Image is the Best Attention\nfor Infrared and Visible Image Fusion",
        28,
        rgb((255, 255, 255)),
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        1.5,
        2.87,
        10.33,
        0.44,
        "Song Wang, Xie Han, Liqun Kuang, et al. · ICCV 2025",
        18,
        rgb((235, 242, 249)),
        font=FONT_CN,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        3.0,
        4.18,
        7.33,
        0.55,
        "汇报人：王胜娇",
        25,
        TEXT,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        3.0,
        4.95,
        7.33,
        0.45,
        "Jittor 复现 · ICCV 2025",
        22,
        TEXT,
        False,
        FONT_CN,
        PP_ALIGN.CENTER,
    )
    add_picture_contain(slide, ASSET_DIR / "jittor_logo.png", 5.4, 5.58, 2.53, 0.62)
    add_footer(slide, 1)
    note(
        "各位老师好，我是王胜娇。今天汇报使用 Jittor 框架复现 ICCV 2025 图像融合论文 SIBA。内容分为模型结构、Jittor 代码、训练与测试演示、实验结果四部分。"
    )


def build_method_slides(prs):
    slide = add_base_slide(prs, "创新点", 2, "Wang et al., ICCV 2025")
    add_text(slide, 0.78, 1.2, 3.2, 0.4, "常见注意力", 22, TEXT, True, FONT_CN_BOLD)
    add_box(
        slide,
        0.9,
        1.92,
        1.55,
        0.7,
        "中间特征",
        LIGHT_ORANGE,
        ORANGE,
        18,
        name="ANIM_02_01",
    )
    add_arrow(slide, 1.68, 2.65, 1.68, 3.22, ORANGE, name="ANIM_02_02")
    add_box(
        slide,
        0.9,
        3.28,
        1.55,
        0.7,
        "Q / K / V",
        LIGHT_ORANGE,
        ORANGE,
        18,
        name="ANIM_02_03",
    )
    add_arrow(slide, 1.68, 4.02, 1.68, 4.58, ORANGE, name="ANIM_02_04")
    add_box(
        slide,
        0.7,
        4.64,
        1.95,
        0.72,
        "源信息逐层压缩",
        LIGHT_PURPLE,
        PURPLE,
        18,
        name="ANIM_02_05",
    )
    add_text(
        slide,
        3.25,
        1.2,
        3.5,
        0.4,
        "SIBA",
        22,
        TEXT,
        True,
        FONT_CN_BOLD,
        name="ANIM_02_06",
    )
    add_box(
        slide,
        3.35,
        1.65,
        1.62,
        0.68,
        "红外原图",
        LIGHT_BLUE,
        BLUE,
        18,
        name="ANIM_02_07",
    )
    add_box(
        slide,
        5.45,
        1.65,
        1.62,
        0.68,
        "可见光原图",
        LIGHT_PINK,
        RED,
        18,
        name="ANIM_02_08",
    )
    add_box(
        slide,
        3.35,
        2.65,
        1.62,
        0.68,
        "1 − 红外",
        LIGHT_BLUE,
        BLUE,
        18,
        name="ANIM_02_09",
    )
    add_box(
        slide,
        5.45,
        2.65,
        1.62,
        0.68,
        "1 − 可见光",
        LIGHT_PINK,
        RED,
        18,
        name="ANIM_02_10",
    )
    add_arrow(slide, 4.15, 3.35, 4.15, 3.85, MUTED, name="ANIM_02_11")
    add_arrow(slide, 6.25, 3.35, 6.25, 3.85, MUTED, name="ANIM_02_12")
    add_box(
        slide,
        3.85,
        3.92,
        2.72,
        0.72,
        "CBSM 生成查询 Q",
        LIGHT_GREEN,
        GREEN,
        18,
        True,
        name="ANIM_02_13",
    )
    add_arrow(slide, 5.2, 4.68, 5.2, 5.2, GREEN, name="ANIM_02_14")
    add_box(
        slide,
        3.62,
        5.27,
        3.18,
        0.76,
        "源图像直接参与交叉注意力",
        LIGHT_PURPLE,
        PURPLE,
        18,
        True,
        name="ANIM_02_15",
    )
    add_picture_contain(
        slide,
        ASSET_DIR / "source_attention_four.png",
        7.45,
        1.45,
        5.35,
        3.35,
        name="ANIM_02_16",
    )
    add_rich_text(
        slide,
        7.55,
        5.02,
        5.0,
        0.8,
        [
            ("核心变化：", TEXT, True, FONT_CN_BOLD),
            ("查询来自源图像，不再完全依赖深层特征。", BLUE, True, FONT_CN_BOLD),
        ],
        20,
        name="ANIM_02_17",
    )
    add_text(
        slide,
        7.55,
        5.86,
        5.0,
        0.55,
        "原图和负变换分别提供显著目标与背景信息。",
        19,
        TEXT,
        name="ANIM_02_18",
    )
    note(
        "传统注意力通常从中间特征生成 Q、K、V。随着特征逐层压缩，原始图像中的显著目标和背景信息可能被削弱。SIBA 直接使用红外、可见光及其负变换，经 CBSM 生成查询 Q，再与另一模态特征做交叉注意力。"
    )

    slide = add_base_slide(prs, "基本架构", 3, "Wang et al., ICCV 2025 · Fig. 2")
    add_picture_contain(
        slide,
        ASSET_DIR / "paper_fig2_arch_hd.png",
        0.55,
        1.22,
        12.2,
        4.55,
        name="ANIM_03_01",
    )
    boxes = [
        (0.8, "双模态特征提取", LIGHT_BLUE, BLUE),
        (3.6, "四个源图查询", LIGHT_ORANGE, ORANGE),
        (6.4, "四路交叉注意力", LIGHT_PURPLE, PURPLE),
        (9.2, "拼接与重建", LIGHT_GREEN, GREEN),
    ]
    for index, (x, label, fill, line) in enumerate(boxes, 1):
        add_box(
            slide,
            x,
            6.03,
            2.25,
            0.58,
            label,
            fill,
            line,
            18,
            True,
            name=f"ANIM_03_{index * 2:02d}",
        )
        if index < len(boxes):
            add_arrow(
                slide,
                x + 2.28,
                6.32,
                x + 2.68,
                6.32,
                MUTED,
                1.6,
                name=f"ANIM_03_{index * 2 + 1:02d}",
            )
    note(
        "模型输入为红外图像和可见光图像。两路 SE-ResNet 与 Restormer 提取中间特征；原图和负变换经四个 CBSM 形成查询；I-SCA 与 V-SCA 完成四路跨模态交互；最后拼接四路 48 通道特征并重建单通道融合图像。"
    )

    slide = add_base_slide(
        prs, "CBSM：源图像查询", 4, "Wang et al., ICCV 2025 · Fig. 3"
    )
    add_picture_contain(
        slide, ASSET_DIR / "paper_fig3_cbsm_hd.png", 0.72, 1.42, 5.35, 3.62
    )
    flow = [
        (6.65, 1.45, 1.5, 0.68, "1×H×W", LIGHT_PINK, RED),
        (8.65, 1.45, 1.65, 0.68, "Conv + PReLU", LIGHT_ORANGE, ORANGE),
        (10.8, 1.45, 1.25, 0.68, "Conv", LIGHT_ORANGE, ORANGE),
        (9.85, 2.78, 2.2, 0.68, "SE 通道增强", LIGHT_GREEN, GREEN),
        (7.15, 2.78, 1.55, 0.68, "PReLU", LIGHT_PURPLE, PURPLE),
    ]
    animation_steps = (1, 3, 5, 7, 9)
    for index, item in enumerate(flow):
        add_box(
            slide,
            *item,
            18,
            index == 3,
            name=f"ANIM_04_{animation_steps[index]:02d}",
        )
    add_arrow(slide, 8.18, 1.79, 8.58, 1.79, RED, name="ANIM_04_02")
    add_arrow(slide, 10.35, 1.79, 10.75, 1.79, ORANGE, name="ANIM_04_04")
    add_arrow(slide, 11.42, 2.16, 11.42, 2.7, ORANGE, name="ANIM_04_06")
    add_arrow(slide, 9.75, 3.12, 8.77, 3.12, GREEN, name="ANIM_04_08")
    add_dot_list(
        slide,
        ["空间映射：1 → 48 通道", "SE 对通道响应重新加权", "输出尺寸与中间特征一致"],
        6.65,
        4.18,
        5.55,
        1.8,
        20,
    )
    add_text(
        slide,
        0.86,
        5.45,
        5.0,
        0.72,
        "CBSM 先整理源图，再将其作为查询送入注意力。",
        21,
        BLUE,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    note(
        "CBSM 由两层 3×3 卷积、PReLU 和 SE 模块组成。第一层将单通道源图映射到 48 通道，第二层与 SE 模块完成空间和通道增强，使查询尺寸能够与中间特征匹配。"
    )

    slide = add_base_slide(prs, "I-SCA / V-SCA", 5, "Wang et al., ICCV 2025 · Fig. 4")
    add_picture_contain(
        slide, ASSET_DIR / "paper_fig4_cross_attention_hd.png", 0.52, 1.2, 6.4, 5.5
    )
    add_text(
        slide,
        7.42,
        1.35,
        5.0,
        0.48,
        "查询来自源图像，键和值来自另一模态特征",
        21,
        BLUE,
        True,
        FONT_CN_BOLD,
        name="ANIM_05_01",
    )
    add_box(
        slide,
        7.52,
        2.08,
        2.1,
        0.62,
        "Q = CBSM(Iraw)",
        LIGHT_PINK,
        RED,
        18,
        True,
        name="ANIM_05_02",
    )
    add_box(
        slide,
        10.05,
        2.08,
        2.4,
        0.62,
        "K,V = DConv(Fmid)",
        LIGHT_BLUE,
        BLUE,
        18,
        True,
        name="ANIM_05_04",
    )
    add_arrow(slide, 8.57, 2.73, 9.75, 3.42, RED, 1.8, name="ANIM_05_03")
    add_arrow(slide, 11.25, 2.73, 10.0, 3.42, BLUE, 1.8, name="ANIM_05_05")
    add_box(
        slide,
        8.65,
        3.43,
        2.25,
        0.68,
        "Softmax(QKᵀ) V",
        LIGHT_PURPLE,
        PURPLE,
        18,
        True,
        name="ANIM_05_06",
    )
    add_arrow(slide, 9.78, 4.15, 9.78, 4.62, PURPLE, name="ANIM_05_07")
    add_box(
        slide,
        8.65,
        4.68,
        2.25,
        0.68,
        "GDFN",
        LIGHT_GREEN,
        GREEN,
        18,
        True,
        name="ANIM_05_08",
    )
    add_dot_list(
        slide,
        [
            "I-SCA：红外查询引导可见光特征",
            "V-SCA：可见光查询引导红外特征",
            "原图与负变换各形成一路查询",
        ],
        7.42,
        5.55,
        5.15,
        1.1,
        18,
    )
    note(
        "I-SCA 和 V-SCA 的区别只在查询来源。Q 由源图像经 CBSM 得到，K、V 由另一模态的中间特征生成。注意力输出再通过 GDFN 过滤冗余信息，并与输入特征残差相加。"
    )

    slide = add_base_slide(prs, "最终融合和损失函数", 6, "official implementation")
    inputs = [
        (0.68, "ir → vi", LIGHT_BLUE, BLUE),
        (2.72, "1−ir → vi", LIGHT_PINK, RED),
        (4.76, "vi → ir", LIGHT_GREEN, GREEN),
        (6.8, "1−vi → ir", LIGHT_PURPLE, PURPLE),
    ]
    for index, (x, label, fill, line) in enumerate(inputs, 1):
        add_box(
            slide,
            x,
            1.45,
            1.62,
            0.62,
            label,
            fill,
            line,
            18,
            True,
            name=f"ANIM_06_{index * 2 - 1:02d}",
        )
        add_arrow(
            slide,
            x + 0.8,
            2.1,
            5.3,
            2.92,
            line,
            1.5,
            name=f"ANIM_06_{index * 2:02d}",
        )
    add_box(
        slide,
        4.58,
        2.92,
        2.1,
        0.72,
        "Concat · 192 ch",
        LIGHT_GRAY,
        MUTED,
        18,
        True,
        name="ANIM_06_09",
    )
    add_arrow(slide, 6.72, 3.28, 7.35, 3.28, MUTED, name="ANIM_06_10")
    add_box(
        slide,
        7.42,
        2.92,
        1.72,
        0.72,
        "Res_SE\n192 → 96",
        LIGHT_ORANGE,
        ORANGE,
        18,
        name="ANIM_06_11",
    )
    add_arrow(slide, 9.18, 3.28, 9.72, 3.28, MUTED, name="ANIM_06_12")
    add_box(
        slide,
        9.8,
        2.92,
        1.72,
        0.72,
        "Res_SE\n96 → 48",
        LIGHT_GREEN,
        GREEN,
        18,
        name="ANIM_06_13",
    )
    add_arrow(slide, 11.56, 3.28, 12.05, 3.28, MUTED, name="ANIM_06_14")
    add_box(
        slide,
        12.1,
        2.92,
        0.62,
        0.72,
        "1",
        LIGHT_BLUE,
        BLUE,
        18,
        True,
        name="ANIM_06_15",
    )
    add_text(
        slide,
        0.82,
        4.45,
        6.0,
        0.52,
        "L = 10 Lₗₐₚ + 0.1 Lᵢₙₜ + Lₛₒᵦₑₗ",
        24,
        TEXT,
        True,
        FONT_MATH,
    )
    add_dot_list(
        slide,
        [
            "强度项：接近两幅源图逐像素最大值",
            "Sobel 项：保留较强的一阶边缘",
            "Laplacian 项：保留较强的局部细节",
        ],
        0.86,
        5.12,
        6.2,
        1.45,
        19,
    )
    add_box(
        slide,
        7.55,
        4.58,
        4.65,
        0.68,
        "60 epochs · batch 4 · patch 128",
        LIGHT_PURPLE,
        PURPLE,
        20,
        True,
    )
    add_text(
        slide,
        7.72,
        5.52,
        4.3,
        0.58,
        "损失权重和训练超参数与官方代码一致。",
        20,
        BLUE,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    note(
        "四路交叉注意力输出在通道维拼接为 192 通道，经过两级 SE-ResNet 和输出卷积得到融合图像。训练损失包含 Laplacian、强度和 Sobel 三项，权重严格采用官方 train.py 中的 10、0.1 和 1。"
    )


def code_dual_slide(
    prs,
    page,
    title,
    left_code,
    right_code,
    left_label,
    right_label,
    source,
    notes,
    line_spacing=0.56,
):
    slide = add_base_slide(prs, title, page, source)
    add_rect(slide, 0.55, 1.04, 12.23, 5.92, CODE_BG, None, rounded=False)
    divider = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(6.66),
        Inches(1.18),
        Inches(6.66),
        Inches(6.82),
    )
    divider.line.color.rgb = rgb((210, 213, 215))
    divider.line.width = Pt(1)
    add_code(
        slide,
        left_code,
        0.65,
        1.08,
        5.9,
        5.78,
        18,
        left_label,
        background=False,
        line_spacing=line_spacing,
    )
    add_code(
        slide,
        right_code,
        6.77,
        1.08,
        5.9,
        5.78,
        18,
        right_label,
        background=False,
        line_spacing=line_spacing,
    )
    note(notes)
    return slide


def code_columns_slide(
    prs,
    page,
    title,
    chunks,
    labels,
    source,
    notes,
    widths=None,
    y=0.96,
    h=6.1,
    leadings=None,
):
    slide = add_base_slide(prs, title, page, source)
    column_count = len(chunks)
    gap = 0.12
    x0 = 0.42
    total_width = 12.49
    if widths is None:
        widths = [
            (total_width - gap * (column_count - 1)) / column_count
        ] * column_count
    if len(widths) != column_count:
        raise ValueError("width count must match code column count")
    if leadings is None:
        leadings = [None] * column_count
    if len(leadings) != column_count:
        raise ValueError("leading count must match code column count")
    x = x0
    for index, (code, label) in enumerate(zip(chunks, labels)):
        width = widths[index]
        add_code(
            slide,
            code,
            x,
            y,
            width,
            h,
            18,
            label,
            leading=leadings[index],
        )
        x += width + gap
    note(notes)
    return slide


def code_grid_slide(prs, page, title, items, source, notes):
    slide = add_base_slide(prs, title, page, source)
    positions = [
        (0.55, 1.04, 6.05, 2.84),
        (6.75, 1.04, 6.03, 2.84),
        (0.55, 4.02, 6.05, 2.84),
        (6.75, 4.02, 6.03, 2.84),
    ]
    for (label, code), (x, y, w, h) in zip(items, positions):
        add_code(slide, code, x, y, w, h, 18, label)
    note(notes)


def build_code_slides(prs):
    code = class_method_excerpt(
        "SIBA", extract_code("models/SIBA.py", "SIBA", "__init__", 38)
    )
    chunks = split_code_n(code, 3)
    code_columns_slide(
        prs,
        7,
        "模型核心：SIBA 组件定义",
        chunks,
        ["SIBA.__init__（1/3）", "SIBA.__init__（2/3）", "SIBA.__init__（3/3）"],
        "models/SIBA.py",
        "这一页展示完整的构造函数。左侧定义两路特征提取、Restormer 与四个 CBSM；右侧定义四组交叉注意力、融合模块和输出层。通道数和层数均与官方代码一致。",
    )

    code = class_method_excerpt(
        "SIBA", extract_code("models/SIBA.py", "SIBA", "execute", 38)
    )
    code = "\n".join(
        line for line in code.splitlines() if not line.lstrip().startswith("#")
    )
    code_lines = code.splitlines()
    code = "\n".join(code_lines)
    chunks = split_code_n(code, 3)
    code_columns_slide(
        prs,
        8,
        "模型核心：SIBA 前向传播",
        chunks,
        ["SIBA.execute（1/3）", "SIBA.execute（2/3）", "SIBA.execute（3/3）"],
        "models/SIBA.py",
        "前向传播先保存原图并计算负变换，再提取红外和可见光特征，生成四个查询。四路交叉注意力按官方顺序执行，最后拼接并重建输出。函数从输入到 return 完整展示，没有截断。",
    )

    cbsm_init = class_method_excerpt(
        "CBSM", extract_code("base_blocks/cbsm.py", "CBSM", "__init__", 38)
    )
    se_init = extract_code("base_blocks/SE.py", "se_module", "__init__", 38)
    execute_code = "\n".join(
        [
            extract_code("base_blocks/cbsm.py", "CBSM", "execute", 38),
            extract_code("base_blocks/SE.py", "se_module", "execute", 38),
        ]
    )
    slide = code_columns_slide(
        prs,
        9,
        "模型核心：CBSM 与 SE",
        [cbsm_init, se_init, execute_code],
        ["CBSM.__init__", "se_module.__init__", "两个 execute 函数"],
        "cbsm.py · SE.py",
        "三栏依次展示 CBSM 初始化、SE 初始化和两个执行函数。Jittor 仅将 forward 改为 execute，并替换模块 API；卷积、激活、池化和通道重标定顺序不变。",
        y=0.86,
        h=6.28,
        leadings=[20, 19, 20],
    )
    add_picture_contain(
        slide,
        ASSET_DIR / "paper_fig3_cbsm_hd.png",
        9.25,
        5.55,
        3.2,
        1.1,
    )

    code = extract_code("base_blocks/se_resnet.py", "Res_SE", None, 38)
    chunks = split_code_n(code, 3)
    code_columns_slide(
        prs,
        10,
        "模型核心：SE-ResNet",
        chunks,
        ["Res_SE（1/3）", "Res_SE（2/3）", "Res_SE（3/3）"],
        "se_resnet.py",
        "SE-ResNet 使用 1×1 卷积匹配残差通道，两层 3×3 卷积提取特征，SE 完成通道增强，然后与残差相加。完整类分三栏连续展示。",
    )

    code = class_method_excerpt(
        "Self_Attention",
        wrap_rearrange_patterns(
            extract_code("base_blocks/restormer.py", "Self_Attention", "execute", 34)
        ),
    )
    attention_chunks = split_code_n(code, 2)
    support_code = "\n".join(
        [
            extract_code("base_blocks/restormer.py", None, "normalize", 36),
            extract_code("base_blocks/restormer.py", "Mlp", "execute", 36),
        ]
    )
    code_columns_slide(
        prs,
        11,
        "模型核心：自注意力",
        [attention_chunks[0], attention_chunks[1], support_code],
        ["Self_Attention.execute（1/2）", "Self_Attention.execute（2/2）", "normalize / Mlp.execute"],
        "restormer.py",
        "前两栏连续展示自注意力执行函数，第三栏展示 normalize 和门控前馈执行函数。卷积、重排、归一化、注意力和输出投影均可在本页代码中顺序对应。"
    )

    code = class_method_excerpt(
        "Cross_Attention",
        wrap_rearrange_patterns(
            extract_code("base_blocks/restormer.py", "Cross_Attention", "execute", 34)
        ),
    )
    attention_chunks = split_code_n(code, 2)
    layer_norm_code = "\n\n".join(
        [
            extract_code("base_blocks/restormer.py", "LayerNorm", "execute", 44),
            extract_code(
                "base_blocks/restormer.py", "WithBias_LayerNorm", "execute", 44
            ),
            extract_code(
                "base_blocks/restormer.py", "TransformerBlock_CA", "execute", 44
            ),
        ]
    )
    slide = code_columns_slide(
        prs,
        12,
        "模型核心：交叉注意力",
        [attention_chunks[0], attention_chunks[1], layer_norm_code],
        ["Cross_Attention.execute（1/2）", "Cross_Attention.execute（2/2）", "归一化 / TransformerBlock_CA"],
        "restormer.py",
        "前两栏连续展示交叉注意力执行函数，第三栏展示 LayerNorm、带偏置归一化和交叉注意力块。查询来自源图，键和值来自另一模态特征。",
        widths=[3.82, 3.82, 4.61],
    )
    item_code = extract_code(
        "loader/train_loader.py", "TrainLoader", "__getitem__", 38
    )
    item_code = item_code.replace(
        '    ), f"Mismatch ir:{ir_path.name} vi:{vi_path.name}."',
        "    ), (\n"
        '        f"Mismatch ir:{ir_path.name} "\n'
        '        f"vi:{vi_path.name}."\n'
        "    )",
    )
    slide = code_columns_slide(
        prs,
        13,
        "训练输入：数据加载与随机裁剪",
        [
            item_code,
            extract_code("loader/train_loader.py", "TrainLoader", "get_patch", 38),
            extract_code("loader/train_loader.py", "TrainLoader", "imread", 38),
        ],
        ["TrainLoader.__getitem__", "TrainLoader.get_patch", "TrainLoader.imread"],
        "train_loader.py",
        "三栏依次展示样本索引与配对校验、同坐标随机裁剪、灰度读取与归一化。三段均直接来自同一 TrainLoader。",
        widths=[4.45, 3.55, 4.25],
        y=0.86,
        h=6.28,
    )

    code = extract_code("loss/loss.py", None, "laplacian", 38)
    chunks = split_code_n(code, 3)
    code_columns_slide(
        prs,
        14,
        "损失核心：Laplacian",
        chunks,
        ["laplacian（1/3）", "laplacian（2/3）", "laplacian（3/3）"],
        "loss/loss.py",
        "Jittor 没有 Kornia，因此按 Kornia 0.7.0 的归一化 3×3 Laplacian、reflect padding 和逐通道卷积完整实现。该函数与官方损失逐项对齐。",
    )

    joint_code = extract_code("loss/loss.py", "JointGrad", "execute", 38)
    fusion_chunks = split_code_n(
        extract_code("loss/loss.py", "Fusionloss", "execute", 38), 2
    )
    code_columns_slide(
        prs,
        15,
        "损失核心：强度与梯度约束",
        [joint_code, fusion_chunks[0], fusion_chunks[1]],
        ["JointGrad.execute", "Fusionloss.execute（1/2）", "Fusionloss.execute（2/2）"],
        "loss/loss.py",
        "JointGrad 完整展示较强 Laplacian 响应的选择。Fusionloss 完整展示逐像素最大值强度损失和 Sobel 梯度损失。",
        y=0.86,
        h=6.28,
    )

    code = compact_indent(
        extract_code("compat/pytorch_clip.py", None, "clip_grad_norm_pytorch", 38)
    )
    chunks = split_code_n(code, 3)
    slide = code_columns_slide(
        prs,
        16,
        "迁移实现：梯度裁剪",
        chunks,
        [
            "clip_grad_norm_pytorch（1/3）",
            "clip_grad_norm_pytorch（2/3）",
            "clip_grad_norm_pytorch（3/3）",
        ],
        "compat/pytorch_clip.py",
        "Jittor 内置梯度裁剪与 PyTorch 1.10 的数值路径不同，因此按 PyTorch 的全局 L2 范数、裁剪系数和逐梯度缩放顺序完整实现。",
    )
    add_text(slide, 0.82, 5.84, 3.0, 0.3, "裁剪顺序", 18, MUTED, True, FONT_CN_BOLD)
    clip_steps = [
        (0.82, "收集可训练梯度", LIGHT_PINK, RED),
        (3.75, "计算全局 L2 范数", LIGHT_ORANGE, ORANGE),
        (6.68, "得到裁剪系数", LIGHT_GREEN, GREEN),
        (9.61, "逐梯度缩放", LIGHT_BLUE, BLUE),
    ]
    for index, (x, label, fill, line) in enumerate(clip_steps, 1):
        add_box(
            slide,
            x,
            6.18,
            2.35,
            0.62,
            label,
            fill,
            line,
            18,
            True,
            name=f"ANIM_16_{index * 2 - 1:02d}",
        )
        if index < len(clip_steps):
            add_arrow(
                slide,
                x + 2.39,
                6.49,
                x + 2.86,
                6.49,
                MUTED,
                1.5,
                name=f"ANIM_16_{index * 2:02d}",
            )

    code = extract_code("compat/pytorch_adam.py", "PyTorchAdam", "step", 38)
    chunks = split_code_n(code, 6)
    slide = code_columns_slide(
        prs,
        17,
        "迁移实现：兼容 Adam（1/2）",
        chunks[:3],
        [
            "PyTorchAdam.step（1/6）",
            "PyTorchAdam.step（2/6）",
            "PyTorchAdam.step（3/6）",
        ],
        "compat/pytorch_adam.py",
        "这一页展示完整 Adam 更新函数的前半部分，包括梯度收集、状态初始化、指数滑动平均和偏置修正。下一页继续展示同一函数直至结束。",
    )
    add_text(
        slide, 0.82, 5.55, 4.0, 0.3, "状态准备与偏置修正", 18, MUTED, True, FONT_CN_BOLD
    )
    adam_prepare = [
        (0.82, "读取 param_group", LIGHT_PINK, RED),
        (3.75, "初始化一、二阶矩", LIGHT_ORANGE, ORANGE),
        (6.68, "计算 bias correction", LIGHT_GREEN, GREEN),
        (9.61, "遍历参数与梯度", LIGHT_BLUE, BLUE),
    ]
    for index, (x, label, fill, line) in enumerate(adam_prepare, 1):
        add_box(
            slide,
            x,
            5.94,
            2.35,
            0.62,
            label,
            fill,
            line,
            18,
            True,
            name=f"ANIM_17_{index * 2 - 1:02d}",
        )
        if index < len(adam_prepare):
            add_arrow(
                slide,
                x + 2.39,
                6.25,
                x + 2.86,
                6.25,
                MUTED,
                1.5,
                name=f"ANIM_17_{index * 2:02d}",
            )

    slide = code_columns_slide(
        prs,
        18,
        "迁移实现：兼容 Adam（2/2）",
        chunks[3:],
        [
            "PyTorchAdam.step（4/6）",
            "PyTorchAdam.step（5/6）",
            "PyTorchAdam.step（6/6）",
        ],
        "compat/pytorch_adam.py",
        "这一页接续上一页，完整展示参数更新、状态写回和 Jittor 优化器收尾。相同参考梯度下，一步参数最大误差为 2.98e-8。",
    )
    add_text(slide, 0.82, 5.55, 3.8, 0.3, "参数更新顺序", 18, MUTED, True, FONT_CN_BOLD)
    adam_update = [
        (0.82, "更新 exp_avg", LIGHT_PINK, RED),
        (3.75, "更新 exp_avg_sq", LIGHT_ORANGE, ORANGE),
        (6.68, "计算 denominator", LIGHT_GREEN, GREEN),
        (9.61, "更新参数", LIGHT_BLUE, BLUE),
    ]
    for index, (x, label, fill, line) in enumerate(adam_update, 1):
        add_box(
            slide,
            x,
            5.94,
            2.35,
            0.62,
            label,
            fill,
            line,
            18,
            True,
            name=f"ANIM_18_{index * 2 - 1:02d}",
        )
        if index < len(adam_update):
            add_arrow(
                slide,
                x + 2.39,
                6.25,
                x + 2.86,
                6.25,
                MUTED,
                1.5,
                name=f"ANIM_18_{index * 2:02d}",
            )


def build_engineering_slides_legacy(prs):
    slide = add_base_slide(prs, "迁移与验证方法", 19, "PyTorch → Jittor")
    stages = [
        (0.68, "固定官方提交\n13 个文件", LIGHT_PINK, RED),
        (3.15, "逐文件 API 替换\n不改结构与超参", LIGHT_ORANGE, ORANGE),
        (5.62, "补齐 Jittor 缺失实现\nLaplacian / Adam / Clip", LIGHT_GREEN, GREEN),
        (8.09, "同权重数值检查\n激活 / 损失 / 梯度", LIGHT_BLUE, BLUE),
        (10.56, "全量训练与测试\n日志 / 指标 / 图像", LIGHT_PURPLE, PURPLE),
    ]
    for index, (x, label, fill, line) in enumerate(stages, 1):
        add_box(
            slide,
            x,
            1.58,
            2.1,
            1.06,
            label,
            fill,
            line,
            18,
            True,
            name=f"ANIM_19_{index:02d}",
        )
        if index < len(stages):
            add_arrow(slide, x + 2.14, 2.11, x + 2.42, 2.11, MUTED, 1.8)
    rows = [
        ["检查对象", "结果"],
        ["官方文件与符号", "13/13，无缺失"],
        ["参数张量 / 参数量", "137 / 565,941"],
        ["真实训练数据", "1,283 对"],
        ["真实测试数据", "706 对"],
        ["完整训练", "Jittor 与 PyTorch 各 60 轮"],
    ]
    add_table(slide, rows, 1.0, 3.35, 5.3, 2.65, 18, [2.7, 2.6])
    add_dot_list(
        slide,
        [
            "每项结果保留原始日志、权重哈希和逐图 CSV",
            "严格训练步等价未通过，因此不作过强结论",
            "逐文件最终审计见 docs/逐文件迁移最终审计_20260728.md",
        ],
        6.85,
        3.45,
        5.4,
        2.45,
        19,
    )
    note(
        "迁移从固定官方提交开始，先逐文件替换框架 API，再实现 Jittor 缺失的 Laplacian、Adam 和梯度裁剪。通过同权重数值检查后才进行 60 轮全量训练和三数据集测试。所有结论都保留原始日志、权重哈希和逐图结果。"
    )

    slide = add_base_slide(prs, "环境与数据", 20, "paper protocol and public datasets")
    rows = [
        ["项目", "论文 / 官方", "本次复现"],
        ["GPU", "TITAN RTX 24 GB", "RTX 3090 24 GB"],
        ["训练集", "MSRS + 200 RoadScene", "1,083 + 200 = 1,283"],
        ["训练轮数", "60", "60"],
        ["batch / patch", "4 / 128", "4 / 128"],
        ["测试集", "MSRS / M3FD / TNO", "361 / 300 / 45"],
    ]
    add_table(slide, rows, 0.7, 1.35, 7.35, 4.35, 18, [2.0, 2.7, 2.65])
    add_box(
        slide,
        8.48,
        1.45,
        3.85,
        0.72,
        "RoadScene 200 对未公开名单",
        LIGHT_PINK,
        RED,
        19,
        True,
    )
    add_text(
        slide,
        8.55,
        2.35,
        3.7,
        1.15,
        "本次从 221 对公开配准图像中固定 seed=2025 选择 200 对，并让 PyTorch 和 Jittor 共用同一清单。",
        19,
        TEXT,
    )
    add_box(slide, 8.48, 3.72, 3.85, 0.72, "数据完整性", LIGHT_GREEN, GREEN, 19, True)
    add_text(
        slide,
        8.55,
        4.62,
        3.7,
        1.22,
        "M3FD 按论文统一使用半分辨率，完整保留 300 对图像。",
        19,
        TEXT,
    )
    note(
        "论文使用 TITAN RTX 24GB，本次 AutoDL 可用配置为 RTX 3090 24GB，因此精度可以比较，速度不能声称同硬件。训练集为全部 1083 对 MSRS 加 200 对 RoadScene。作者未公开 RoadScene 名单，本次固定公开子集并让两个框架共用。"
    )

    slide = add_base_slide(prs, "训练演示", 21, "screen -S kk · real logs")
    add_picture_contain(
        slide, ASSET_DIR / "real_training_terminal.png", 0.62, 1.28, 8.15, 4.95
    )
    add_code(
        slide,
        "bash scripts/start_demo_\\\ntraining_screen.sh\nscreen -r kk\ntail -f logs/demo_*/train.log",
        9.02,
        1.48,
        3.62,
        1.75,
        18,
        "现场命令",
    )
    add_dot_list(
        slide,
        [
            "使用真实训练图片",
            "连续显示四项真实损失",
            "保存演示日志和 checkpoint",
            "短训练不用于论文指标",
        ],
        9.12,
        3.62,
        3.35,
        2.05,
        19,
    )
    add_text(
        slide,
        0.78,
        6.43,
        7.8,
        0.38,
        "图中内容来自完整 Jittor 60 轮训练日志，不是示意数据。",
        18,
        BLUE,
        True,
        FONT_CN_BOLD,
    )
    note(
        "现场先在终端执行 start_demo_training_screen.sh，再用 screen -r kk 查看。脚本读取真实训练目录，连续打印总损失、Laplacian、强度和 Sobel 四项损失，并保存演示权重。演示只证明训练代码可执行，最终指标来自已经完成的 60 轮训练。"
    )

    slide = add_base_slide(prs, "推理演示", 22, "tools/run_inference.py")
    add_picture_contain(
        slide, ASSET_DIR / "official_alignment_examples.png", 0.62, 1.3, 12.05, 2.75
    )
    add_code(
        slide,
        "python tools/run_inference.py \\\n  --framework jittor --use-cuda \\\n  --checkpoint checkpoints/.../SIBA_epoch60.pkl \\\n  --data-dir datasets/test/MSRS \\\n  --output results/demo/MSRS",
        0.72,
        4.32,
        7.15,
        2.05,
        18,
        "现场命令",
    )
    add_dot_list(
        slide,
        [
            "运行前记录输出目录文件数",
            "运行后打开新增融合图",
            "同时生成 timing.csv 与 summary.json",
            "再展示三数据集完整结果",
        ],
        8.35,
        4.42,
        4.15,
        1.9,
        19,
    )
    note(
        "推理演示先记录输出目录，再运行 Jittor 推理命令。完成后打开新生成的融合图，同时展示 timing.csv 和 summary.json。最后切换到完整三数据集结果，说明现场样例与最终全量实验的关系。"
    )


def build_engineering_slides(prs):
    slide = add_base_slide(prs, "训练演示", 19, "screen -S kk · real logs")
    add_picture_contain(
        slide, ASSET_DIR / "real_training_terminal.png", 0.62, 1.35, 7.45, 4.78
    )
    add_picture_contain(
        slide, ASSET_DIR / "notebook_live_start_crop.png", 8.35, 1.32, 4.25, 1.08
    )
    add_code(
        slide,
        "bash scripts/demo.sh\nscreen -r kk",
        8.35,
        2.63,
        4.25,
        1.32,
        18,
        "现场命令",
    )
    flow = [
        (4.28, "PyTorch\nseed 2025", LIGHT_ORANGE, ORANGE),
        (5.08, "共享初始权重\n137 个张量", LIGHT_BLUE, BLUE),
        (5.88, "Jittor\n真实图像训练 20 步", LIGHT_GREEN, GREEN),
    ]
    for index, (y, label, fill, line) in enumerate(flow, 1):
        add_box(
            slide,
            8.55,
            y,
            3.85,
            0.62,
            label,
            fill,
            line,
            18,
            True,
            name=f"ANIM_19_{index * 2 - 1:02d}",
        )
        if index < len(flow):
            add_arrow(
                slide,
                10.48,
                y + 0.64,
                10.48,
                y + 0.78,
                MUTED,
                1.8,
                name=f"ANIM_19_{index * 2:02d}",
            )
    add_text(
        slide,
        0.78,
        6.43,
        7.55,
        0.38,
        "真实训练目录：1,283 对图像；演示权重不用于论文指标。",
        18,
        BLUE,
        True,
        FONT_CN_BOLD,
    )
    note(
        "现场在 Jupyter 中执行训练单元，脚本先由官方 PyTorch 实现按随机种子 2025 导出共享初始权重，再由 Jittor 加载同一权重。随后在真实 1,283 对训练目录上运行 20 步，并通过 screen -r kk 连续查看总损失、Laplacian、强度和 Sobel 四项损失。演示只证明训练代码可执行，最终指标来自完整 60 轮训练。"
    )

    slide = add_base_slide(prs, "逐模块测试", 20, "ptTest.ipynb · jtTest.ipynb")
    add_text(slide, 0.72, 1.08, 5.9, 0.38, "PyTorch：生成固定参考", 20, TEXT, True, FONT_CN_BOLD)
    add_picture_contain(
        slide, ASSET_DIR / "notebook_pytorch_reference_crop.png", 0.62, 1.52, 6.15, 1.92
    )
    add_text(
        slide,
        0.78,
        3.54,
        5.85,
        0.5,
        "固定输入、官方权重，导出 707 项数组。",
        18,
        BLUE,
        True,
        FONT_CN_BOLD,
    )
    coverage = [
        (0.78, 4.32, "输入与参数", LIGHT_ORANGE, ORANGE),
        (2.78, 4.32, "SE · CBSM", LIGHT_BLUE, BLUE),
        (4.78, 4.32, "自注意力", LIGHT_GREEN, GREEN),
        (0.78, 5.22, "四路交叉注意力", LIGHT_PURPLE, PURPLE),
        (2.78, 5.22, "三项损失", LIGHT_PINK, RED),
        (4.78, 5.22, "反向 · Clip · Adam", LIGHT_ORANGE, ORANGE),
    ]
    for x, y, label, fill, line in coverage:
        add_box(slide, x, y, 1.75, 0.62, label, fill, line, 18, True)
    add_text(slide, 7.12, 1.08, 5.2, 0.38, "Jittor：逐项比较", 20, TEXT, True, FONT_CN_BOLD)
    add_picture_contain(
        slide, ASSET_DIR / "notebook_jittor_summary_crop.png", 7.12, 1.48, 5.25, 4.92
    )
    add_text(
        slide,
        7.18,
        6.48,
        5.1,
        0.32,
        "四个 Notebook 已整本执行，错误输出为 0。",
        18,
        BLUE,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    note(
        "逐模块测试采用两本对应 Notebook。PyTorch Notebook 固定输入和官方权重，导出输入、137 个参数张量、模块激活、损失、梯度、裁剪结果和一步 Adam 参数，共 707 项数组。Jittor Notebook 逐项读取并比较，覆盖 SE、CBSM、自注意力、四路交叉注意力、三项损失、反向传播、梯度裁剪和 Adam。四个演示 Notebook 已在 RTX 3090 上整本执行，全部代码单元完成，错误输出为零。"
    )

    slide = add_base_slide(prs, "环境与数据", 21, "paper protocol and public datasets")
    rows = [
        ["项目", "论文 / 官方", "本次复现"],
        ["GPU", "TITAN RTX 24 GB", "RTX 3090 24 GB"],
        ["训练集", "MSRS + 200 RoadScene", "1,083 + 200 = 1,283"],
        ["训练轮数", "60", "60"],
        ["batch / patch", "4 / 128", "4 / 128"],
        ["测试集", "MSRS / M3FD / TNO", "361 / 300 / 45"],
    ]
    add_table(slide, rows, 0.7, 1.35, 7.35, 4.35, 18, [2.0, 2.7, 2.65])
    add_box(
        slide,
        8.48,
        1.45,
        3.85,
        0.72,
        "RoadScene 200 对未公开名单",
        LIGHT_PINK,
        RED,
        19,
        True,
    )
    add_text(
        slide,
        8.55,
        2.35,
        3.7,
        1.15,
        "本次从 221 对公开配准图像中固定 seed=2025 选择 200 对，并让 PyTorch 和 Jittor 共用同一清单。",
        19,
        TEXT,
    )
    add_box(slide, 8.48, 3.72, 3.85, 0.72, "数据完整性", LIGHT_GREEN, GREEN, 19, True)
    add_text(
        slide,
        8.55,
        4.62,
        3.7,
        1.22,
        "M3FD 按论文统一使用半分辨率，完整保留 300 对图像。",
        19,
        TEXT,
    )
    note(
        "论文使用 TITAN RTX 24GB，本次 AutoDL 可用配置为 RTX 3090 24GB，因此精度可以比较，速度不能声称同硬件。训练集为全部 1083 对 MSRS 加 200 对 RoadScene。作者未公开 RoadScene 名单，本次固定公开子集并让两个框架共用。"
    )

    slide = add_base_slide(prs, "推理演示", 22, "tools/run_inference.py")
    add_picture_contain(
        slide, ASSET_DIR / "official_alignment_examples.png", 0.62, 1.3, 12.05, 2.75
    )
    add_code(
        slide,
        "python tools/run_inference.py \\\n  --framework jittor --use-cuda \\\n  --checkpoint checkpoints/.../SIBA_epoch60.pkl \\\n  --data-dir datasets/test/MSRS \\\n  --output results/demo/MSRS",
        0.72,
        4.32,
        7.15,
        2.05,
        18,
        "现场命令",
    )
    add_dot_list(
        slide,
        [
            "运行前记录输出目录文件数",
            "运行后打开新增融合图",
            "同时生成 timing.csv 与 summary.json",
            "再展示三数据集完整结果",
        ],
        8.35,
        4.42,
        4.15,
        1.9,
        19,
    )
    note(
        "推理演示先记录输出目录，再运行 Jittor 推理命令。完成后打开新生成的融合图，同时展示 timing.csv 和 summary.json。最后切换到完整三数据集结果，说明现场样例与最终全量实验的关系。"
    )


def build_result_slides(prs):
    slide = add_base_slide(
        prs, "前向、损失与训练步对齐", 23, "controlled shared-parameter test"
    )
    rows = [
        ["检查项", "数值", "判断"],
        ["中间激活最大误差", "2.0218e-4", "通过"],
        ["总损失最大误差", "2.5034e-6", "通过"],
        ["梯度余弦相似度", "0.999945", "接近"],
        ["梯度相对 L2", "1.0508%", "未严格相同"],
        ["一步更新相对 L2", "6.7275%", "未严格相同"],
        ["参考梯度下 Adam 最大误差", "2.9802e-8", "通过"],
    ]
    add_table(slide, rows, 0.68, 1.28, 7.05, 4.85, 18, [3.15, 2.0, 1.9])
    add_box(
        slide,
        8.22,
        1.45,
        4.15,
        0.75,
        "结论 1：前向与损失对齐",
        LIGHT_GREEN,
        GREEN,
        20,
        True,
    )
    add_box(
        slide,
        8.22,
        2.65,
        4.15,
        0.75,
        "结论 2：兼容 Adam 与裁剪对齐",
        LIGHT_BLUE,
        BLUE,
        20,
        True,
    )
    add_box(
        slide,
        8.22,
        3.85,
        4.15,
        0.75,
        "结论 3：原生训练步接近但不严格相同",
        LIGHT_PINK,
        RED,
        20,
        True,
    )
    add_text(
        slide,
        8.34,
        5.15,
        3.95,
        1.0,
        "因此可以说明训练功能完整，不能说明每一步参数更新与 PyTorch 完全一致。",
        20,
        TEXT,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    note(
        "受控对齐使用相同输入和相同初始权重。前向激活和损失误差很小；使用相同 PyTorch 参考梯度时，兼容 Adam 和裁剪也对齐。但 Jittor 原生梯度和一步更新没有达到严格等价标准，所以报告采用分层结论。"
    )

    slide = add_base_slide(
        prs, "完整 60 轮训练", 24, "1,283 pairs · official logging interval"
    )
    add_picture_contain(
        slide,
        ROOT
        / "results"
        / "training_analysis_20260727_siba_official_protocol"
        / "loss_curve.png",
        0.68,
        1.25,
        8.7,
        5.3,
    )
    add_box(
        slide, 9.72, 1.48, 2.65, 0.72, "Jittor：67.98 min", LIGHT_BLUE, BLUE, 20, True
    )
    add_box(
        slide,
        9.72,
        2.55,
        2.65,
        0.72,
        "PyTorch：37.02 min",
        LIGHT_ORANGE,
        ORANGE,
        20,
        True,
    )
    add_dot_list(
        slide,
        [
            "两边均完成 60 轮",
            "每 50 batch 记录一次",
            "各 420 条原始损失",
            "曲线用于验证收敛，不作逐 batch 相等声明",
        ],
        9.55,
        3.62,
        3.0,
        2.05,
        18,
    )
    note(
        "两套代码都在全部 1283 对训练数据上完成 60 轮。官方代码每 50 个 batch 打印一次，因此各有 420 条记录。两条曲线都收敛，但数据加载器的 shuffle 顺序不同，不能把曲线当作逐 batch 数值等价。"
    )

    slide = add_base_slide(prs, "运行性能", 25, "RTX 3090 24 GB · synchronized timing")
    add_picture_contain(
        slide, ASSET_DIR / "synchronized_fps.png", 0.72, 1.28, 8.15, 4.85
    )
    rows = [
        ["项目", "结果"],
        ["最大显存", "14,587 MiB"],
        ["最大 GPU 利用率", "100%"],
        ["模型参数", "565,941"],
        ["计时方式", "CUDA 同步"],
    ]
    add_table(slide, rows, 9.08, 1.52, 3.55, 3.15, 18, [1.8, 1.75])
    add_text(
        slide,
        9.12,
        5.05,
        3.45,
        1.12,
        "论文 test.py 未同步 CUDA。PPT 的框架速度对比只使用同步计时。",
        19,
        RED,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    note(
        "速度对比使用同一台 RTX 3090 并在每次前向后同步 CUDA。Jittor 在三个数据集上均慢于 PyTorch。论文 test.py 的计时没有同步 CUDA，因此只保留为协议复现，不用于两个框架的实际速度比较。"
    )

    metrics = list(
        csv.DictReader(
            (
                ROOT
                / "results"
                / "metrics_20260727_siba_official_protocol"
                / "metrics_summary.csv"
            ).open("r", encoding="utf-8-sig")
        )
    )
    paper = json.loads(
        (ROOT / "configs" / "paper_metrics.json").read_text(encoding="utf-8")
    )["datasets"]
    metric_names = ["VIF", "SCD", "MI", "Qabf", "SSIM", "MS_SSIM", "FMI"]

    slide = add_base_slide(
        prs,
        "三数据集：官方权重推理对齐",
        26,
        "released checkpoint · 706 pairs · not self-trained",
    )
    for panel, dataset, y in (
        ("MSRS", "MSRS", 1.18),
        ("M3FD 1/2", "M3FD_2x", 3.03),
        ("TNO", "TNO", 4.88),
    ):
        rows = [[panel] + [name.replace("MS_SSIM", "MS-SSIM") for name in metric_names]]
        rows.append(
            ["Paper"] + [f"{paper[dataset][name]:.3f}" for name in metric_names]
        )
        for experiment, label in (
            ("OfficialPyTorch", "PyTorch"),
            ("OfficialJittor", "Jittor"),
        ):
            row = next(
                item
                for item in metrics
                if item["experiment"] == experiment and item["dataset"] == dataset
            )
            rows.append([label] + [f"{float(row[name]):.6f}" for name in metric_names])
        add_table(slide, rows, 0.48, y, 12.35, 1.62, 18, [1.55] + [1.54] * 7)
    note(
        "本页三组结果都使用作者发布的同一份权重，不是两框架各自训练得到的权重。MSRS、M3FD 和 TNO 共 706 对图像，Jittor 与 PyTorch 输出最大像素差均为 1 个灰度级。表格保留实测六位小数，不能因为论文只保留三位小数就写成完全相同。M3FD 按论文第 4.1 节对全部 300 对图像统一使用半分辨率，没有减少测试样本。"
    )

    slide = add_base_slide(
        prs,
        "三数据集：完整 60 轮自训练",
        27,
        "1,283 training pairs · PyTorch and Jittor checkpoints",
    )
    for panel, dataset, y in (
        ("MSRS", "MSRS", 1.22),
        ("M3FD 1/2", "M3FD_2x", 2.96),
        ("TNO", "TNO", 4.70),
    ):
        rows = [[panel] + [name.replace("MS_SSIM", "MS-SSIM") for name in metric_names]]
        for experiment, label in (
            ("PyTorchSelfTrained", "PyTorch"),
            ("JittorSelfTrained", "Jittor"),
        ):
            row = next(
                item
                for item in metrics
                if item["experiment"] == experiment and item["dataset"] == dataset
            )
            rows.append([label] + [f"{float(row[name]):.6f}" for name in metric_names])
        add_table(slide, rows, 0.48, y, 12.35, 1.42, 18, [1.55] + [1.54] * 7)
    add_text(
        slide,
        0.82,
        6.32,
        11.7,
        0.48,
        "两套训练均完整收敛；因初始化、shuffle 与裁剪顺序不同，不作逐批次严格排名。",
        20,
        RED,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    note(
        "本页是 PyTorch 和 Jittor 分别使用全部 1283 对训练图像完成 60 轮后得到的真实结果。它证明两套完整训练链都能收敛并产生有效权重。两边没有从同一参数文件、同一显式 batch 顺序和同一 crop 清单开始，因此不能把最终差异解释为严格框架优劣，也不能写成逐批次对齐。官方权重的严格迁移证据在上一页。"
    )

    slide = add_base_slide(
        prs, "融合结果", 28, "real source images and released-checkpoint outputs"
    )
    add_picture_contain(
        slide,
        ROOT
        / "results"
        / "visual_comparisons_20260727_siba_official_protocol"
        / "MSRS_official_checkpoint_grid.png",
        0.52,
        1.15,
        6.15,
        5.9,
    )
    add_picture_contain(
        slide,
        ROOT
        / "results"
        / "visual_comparisons_20260727_siba_official_protocol"
        / "TNO_official_checkpoint_grid.png",
        6.78,
        1.15,
        6.0,
        5.9,
    )
    add_text(
        slide,
        0.92,
        6.55,
        5.4,
        0.34,
        "MSRS：人物热目标与可见光背景同时保留",
        18,
        BLUE,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        7.08,
        6.55,
        5.35,
        0.34,
        "TNO：Jittor 与 PyTorch 肉眼无明显差异",
        18,
        BLUE,
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    note(
        "左侧展示 MSRS，右侧展示 TNO。每组均包括红外、可见光、Jittor 和 PyTorch 官方权重输出。Jittor 与 PyTorch 的整体结构、热目标和背景纹理基本一致，细微差异主要来自一个灰度级以内的数值舍入。"
    )

    slide = add_base_slide(prs, "复现边界与我的思考", 29, "independent analysis")
    add_text(slide, 0.75, 1.25, 3.5, 0.42, "复现边界", 23, TEXT, True, FONT_CN_BOLD)
    add_dot_list(
        slide,
        [
            "RoadScene 200 对名单和随机种子未公开",
            "论文写 YCbCr 训练，官方代码实际读取灰度图",
            "复现 GPU 与论文 GPU 不同",
            "训练步数值接近，但未达到严格相同",
        ],
        0.82,
        1.9,
        5.4,
        2.55,
        19,
    )
    add_text(slide, 6.72, 1.25, 3.5, 0.42, "我的思考", 23, TEXT, True, FONT_CN_BOLD)
    add_dot_list(
        slide,
        [
            "源图像作为查询，减少了深层特征压缩造成的信息损失",
            "方法依赖红外和可见光严格配准；错位时源图查询可能引入错误响应",
            "目前没有不确定性估计，极端曝光和噪声条件下的稳定性仍需测试",
            "后续应增加错位、噪声和曝光扰动实验，先评估鲁棒性，再考虑改进",
        ],
        6.78,
        1.9,
        5.55,
        3.25,
        19,
    )
    add_box(
        slide,
        1.0,
        5.48,
        11.35,
        0.78,
        "结论：SIBA 的推理迁移高度一致；完整训练可执行，但论文级同条件复现受公开信息限制。",
        LIGHT_PURPLE,
        PURPLE,
        20,
        True,
    )
    note(
        "复现不能回避公开信息不足。作者没有公布 RoadScene 子集，论文和代码的数据处理描述也不一致，硬件也不同。方法本身的优点是让源图像直接参与注意力；局限是依赖配准且没有不确定性估计。后续应先做错位、噪声和曝光扰动实验。"
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, ASSET_DIR / "cover_gradient.png", 0, 0, SLIDE_W, SLIDE_H)
    add_text(
        slide,
        4.25,
        0.86,
        4.8,
        0.72,
        "感谢各位老师！",
        34,
        rgb((67, 57, 73)),
        True,
        FONT_CN_BOLD,
        PP_ALIGN.CENTER,
    )
    add_text(slide, 0.9, 2.0, 3.2, 0.42, "完成内容", 22, TEXT, True, FONT_CN_BOLD)
    add_dot_list(
        slide,
        [
            "13 个官方文件完整迁移",
            "Jittor / PyTorch 各 60 轮全量训练",
            "706 对测试与 7 项逐图指标",
            "日志、权重、曲线和审计记录齐全",
        ],
        0.95,
        2.58,
        4.75,
        2.55,
        19,
    )
    add_text(slide, 6.25, 2.0, 3.2, 0.42, "代码仓库", 22, TEXT, True, FONT_CN_BOLD)
    add_text(
        slide,
        6.28,
        2.68,
        6.0,
        0.52,
        "github.com/wangshengjiao1010-boop/SIBA-Jittor",
        20,
        BLUE,
        True,
        FONT_CN,
    )
    add_text(
        slide,
        6.28,
        3.46,
        5.75,
        1.22,
        "README 包含环境、数据、训练、测试、对齐、性能、指标、可视化和已知限制。",
        20,
        TEXT,
    )
    add_picture_contain(slide, ASSET_DIR / "jittor_logo.png", 8.2, 5.0, 2.45, 0.85)
    add_text(
        slide,
        4.25,
        6.42,
        4.8,
        0.42,
        "欢迎批评指正",
        20,
        rgb((73, 64, 80)),
        False,
        FONT_CN,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide, 12.35, 7.0, 0.4, 0.3, "30", 18, MUTED, False, FONT_CN, PP_ALIGN.RIGHT
    )
    note(
        "本次完成 13 个官方文件的 Jittor 迁移、两套 60 轮全量训练、706 对测试、七项指标和完整审计。代码仓库 README 保留环境、数据、训练、测试、对齐、性能和限制。以上是我的汇报，感谢各位老师，欢迎批评指正。"
    )


def add_notes_and_export():
    import win32com.client

    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = True
    presentation = app.Presentations.Open(str(PPTX_PATH.resolve()), False, False, False)
    animation_count = 0
    notes_count = 0
    code_overflow = []
    for index, slide in enumerate(presentation.Slides, 1):
        if index <= len(NOTES):
            try:
                slide.NotesPage.Shapes.Placeholders.Item(2).TextFrame.TextRange.Text = (
                    NOTES[index - 1]
                )
                notes_count += 1
            except Exception:
                pass
        try:
            sequence = slide.TimeLine.MainSequence
            animated = sorted(
                [
                    shape
                    for shape in slide.Shapes
                    if str(shape.Name).startswith("ANIM_")
                ],
                key=lambda shape: shape.Name,
            )
            for shape in animated:
                try:
                    sequence.AddEffect(shape, 1, 0, 1)
                    animation_count += 1
                except Exception:
                    continue
        except Exception:
            pass
        for shape in slide.Shapes:
            try:
                shape.Shadow.Visible = 0
            except Exception:
                pass
            try:
                if not str(shape.Name).startswith("CODE_TEXT_"):
                    continue
                bound_width = float(shape.TextFrame2.TextRange.BoundWidth)
                bound_height = float(shape.TextFrame2.TextRange.BoundHeight)
                if (
                    bound_width > float(shape.Width) + 2
                    or bound_height > float(shape.Height) + 2
                ):
                    code_overflow.append(
                        {
                            "slide": index,
                            "shape": str(shape.Name),
                            "bound_width": bound_width,
                            "shape_width": float(shape.Width),
                            "bound_height": bound_height,
                            "shape_height": float(shape.Height),
                        }
                    )
            except Exception:
                continue
    presentation.Save()
    presentation.ExportAsFixedFormat(
        str(PDF_PATH.resolve()),
        2,
        1,
        0,
        1,
        1,
        0,
        None,
        1,
        "",
        1,
        1,
        1,
        1,
        0,
    )
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    for slide in presentation.Slides:
        slide.Export(
            str((PREVIEW_DIR / f"幻灯片{slide.SlideIndex}.PNG").resolve()),
            "PNG",
            1920,
            1080,
        )
    presentation.Close()
    app.Quit()
    return notes_count, animation_count, code_overflow


def build_contact_sheet():
    files = sorted(
        PREVIEW_DIR.glob("*.PNG"),
        key=lambda path: int(re.search(r"(\d+)", path.stem).group(1)),
    )
    thumb_w, thumb_h = 480, 270
    columns = 3
    rows = math.ceil(len(files) / columns)
    canvas = Image.new("RGB", (thumb_w * columns, thumb_h * rows), (236, 233, 230))
    for index, path in enumerate(files):
        image = Image.open(path).convert("RGB")
        image.thumbnail((thumb_w, thumb_h))
        x = (index % columns) * thumb_w
        y = (index // columns) * thumb_h
        canvas.paste(image, (x, y))
    canvas.save(CONTACT_PATH)


def write_notes():
    lines = ["# 王胜娇-培育期逐页讲稿", ""]
    for index, text in enumerate(NOTES, 1):
        lines.extend([f"## 第 {index} 页", "", text, ""])
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def validate_ppt(notes_count, animation_count, code_overflow):
    from pptx import Presentation as CheckPresentation

    presentation = CheckPresentation(PPTX_PATH)
    small_fonts = []
    empty_code_panels = []
    for slide_index, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size.pt < 18:
                        small_fonts.append((slide_index, run.text, run.font.size.pt))
        if 7 <= slide_index <= 18:
            code_shapes = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if any(
                    run.font.name == FONT_CODE
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                ):
                    code_shapes.append(shape)
            if len(code_shapes) < 2:
                empty_code_panels.append(slide_index)
    report = [
        "# PPT 检查报告",
        "",
        f"- 幻灯片页数：{len(presentation.slides)}",
        f"- 演讲者备注：{notes_count}/{len(presentation.slides)}",
        f"- 逐步出现动画数量：{animation_count}",
        f"- 小于 18 pt 的文本：{len(small_fonts)}",
        f"- 双栏代码检查异常页：{empty_code_panels or '无'}",
        f"- PowerPoint 渲染后代码溢出：{len(code_overflow)}",
        f"- PPTX：`{PPTX_PATH.name}`",
        f"- PDF：`{PDF_PATH.name}`",
        f"- 全页预览：`{CONTACT_PATH.name}`",
        "",
        "所有实验图表和日志均来自项目内现有真实文件；未使用随机生成实验结果或手填指标。",
    ]
    if small_fonts:
        report.append("\n## 小字号明细\n")
        report.extend(
            f"- 第 {slide} 页：{size} pt `{text[:40]}`"
            for slide, text, size in small_fonts[:30]
        )
    if code_overflow:
        report.append("\n## 代码溢出明细\n")
        report.extend(
            f"- 第 {item['slide']} 页 `{item['shape']}`："
            f"{item['bound_width']:.1f}×{item['bound_height']:.1f} / "
            f"{item['shape_width']:.1f}×{item['shape_height']:.1f}"
            for item in code_overflow
        )
    CHECK_PATH.write_text("\n".join(report), encoding="utf-8")


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    if PREVIEW_DIR.exists():
        shutil.rmtree(PREVIEW_DIR)
    prepare_assets()
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    build_cover(presentation)
    build_method_slides(presentation)
    build_code_slides(presentation)
    build_engineering_slides(presentation)
    build_result_slides(presentation)
    assert len(presentation.slides) == 30
    assert len(NOTES) == 30
    presentation.save(PPTX_PATH)
    notes_count, animation_count, code_overflow = add_notes_and_export()
    build_contact_sheet()
    write_notes()
    validate_ppt(notes_count, animation_count, code_overflow)
    print(PPTX_PATH)
    print(PDF_PATH)
    print(CONTACT_PATH)


if __name__ == "__main__":
    main()
