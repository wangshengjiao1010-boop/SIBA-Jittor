from pathlib import Path
from shutil import copy2

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT.parent / "播种期PPT" / "TemplateMC-PPT.pptx"
WORK = ROOT / "ppt_work"
PAPER_DIR = WORK / "paper"
ASSETS = ROOT / "ppt" / "assets_generated"
DELIVERABLE = ROOT / "deliverables" / "SIBA_Jittor_培育期_20260727"
OUT = DELIVERABLE / "姓名-培育期-SIBA-Jittor-可编辑.pptx"
OUTLINE = DELIVERABLE / "SIBA-Jittor-PPT完整大纲.md"

BLUE = RGBColor(0, 103, 177)
NAVY = RGBColor(23, 47, 78)
ORANGE = RGBColor(222, 117, 48)
TEAL = RGBColor(35, 139, 145)
TEXT = RGBColor(45, 54, 64)
MUTED = RGBColor(91, 100, 111)
LIGHT_BLUE = RGBColor(235, 244, 251)
LIGHT_ORANGE = RGBColor(251, 242, 233)
LIGHT_GRAY = RGBColor(245, 247, 249)
MID_GRAY = RGBColor(217, 223, 229)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(186, 43, 43)


def find_font(size):
    for path in (
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/arial.ttf"),
    ):
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def contain(image, size, background="white"):
    image = image.convert("RGB")
    image.thumbnail(size, Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", size, background)
    canvas.paste(image, ((size[0] - image.width) // 2, (size[1] - image.height) // 2))
    return canvas


def label_image(path, label, size=(520, 340)):
    image = contain(Image.open(path), (size[0], size[1] - 46))
    canvas = Image.new("RGB", size, "white")
    canvas.paste(image, (0, 0))
    ImageDraw.Draw(canvas).text((12, size[1] - 38), label, font=find_font(21), fill=(35, 45, 55))
    return canvas


def crop_paper_assets():
    ASSETS.mkdir(parents=True, exist_ok=True)
    crops = {
        "paper_fig2_arch.png": (PAPER_DIR / "page_04.png", (80, 85, 850, 430)),
        "paper_fig3_cbsm.png": (PAPER_DIR / "page_04.png", (105, 765, 445, 925)),
        "paper_fig4_cross_attention.png": (PAPER_DIR / "page_05.png", (80, 175, 455, 495)),
        "paper_ablation_tables.png": (PAPER_DIR / "page_08.png", (465, 90, 850, 610)),
    }
    for name, (source, box) in crops.items():
        Image.open(source).crop(box).save(ASSETS / name)

    page1 = {
        "Infrared source": PAPER_DIR / "img_x322_p1_640x480.png",
        "Visible source": PAPER_DIR / "img_x325_p1_640x480.png",
        "Infrared colormap": PAPER_DIR / "img_x326_p1_640x480.jpeg",
        "Negative colormap": PAPER_DIR / "img_x327_p1_640x480.jpeg",
        "Grad-CAM example 1": PAPER_DIR / "img_x324_p1_500x425.jpeg",
        "Grad-CAM example 2": PAPER_DIR / "img_x323_p1_500x417.jpeg",
    }
    cells = [label_image(path, label, (420, 315)) for label, path in page1.items()]
    sheet = Image.new("RGB", (1260, 630), (242, 244, 246))
    for index, cell in enumerate(cells):
        sheet.paste(cell, ((index % 3) * 420, (index // 3) * 315))
    sheet.save(ASSETS / "source_attention_observation.png")


def first_pair(dataset):
    ir_dir = ROOT / "datasets" / "test" / dataset / "ir"
    vi_dir = ROOT / "datasets" / "test" / dataset / "vi"
    ir = sorted(path for path in ir_dir.iterdir() if path.is_file())[0]
    return ir, vi_dir / ir.name


def compose_dataset_examples():
    rows = []
    for dataset in ("MSRS", "M3FD_2x", "TNO"):
        ir, vi = first_pair(dataset)
        rows.append((label_image(ir, f"{dataset} · Infrared"), label_image(vi, f"{dataset} · Visible")))
    sheet = Image.new("RGB", (1040, 1020), (242, 244, 246))
    for row, pair in enumerate(rows):
        for col, image in enumerate(pair):
            sheet.paste(image, (col * 520, row * 340))
    sheet.save(ASSETS / "dataset_examples.png")


def compose_alignment_examples():
    dataset = "MSRS"
    ir_dir = ROOT / "datasets" / "test" / dataset / "ir"
    vi_dir = ROOT / "datasets" / "test" / dataset / "vi"
    base = ROOT / "results" / "official_checkpoint_alignment_20260727_siba_official_protocol"
    jittor_dir = base / "jittor" / dataset
    pytorch_dir = base / "pytorch" / dataset
    names = [path.name for path in sorted(ir_dir.iterdir()) if path.is_file()][:1]
    sheet = Image.new("RGB", (1760, 335), (242, 244, 246))
    labels = ("Infrared", "Visible", "Jittor released", "PyTorch released")
    for row, name in enumerate(names):
        paths = (ir_dir / name, vi_dir / name, jittor_dir / name, pytorch_dir / name)
        for col, (path, label) in enumerate(zip(paths, labels)):
            sheet.paste(label_image(path, f"{label} · {name}", (440, 335)), (col * 440, row * 335))
    sheet.save(ASSETS / "official_alignment_examples.png")


def compose_self_trained_examples():
    dataset = "TNO"
    ir_dir = ROOT / "datasets" / "test" / dataset / "ir"
    vi_dir = ROOT / "datasets" / "test" / dataset / "vi"
    base = ROOT / "results" / "full_20260727_siba_official_protocol"
    jittor_dir = base / "jittor" / dataset
    pytorch_dir = base / "pytorch" / dataset
    names = ["01.png"]
    sheet = Image.new("RGB", (1760, 335), (242, 244, 246))
    labels = ("Infrared", "Visible", "Jittor 60 epochs", "PyTorch 60 epochs")
    for row, name in enumerate(names):
        paths = (ir_dir / name, vi_dir / name, jittor_dir / name, pytorch_dir / name)
        for col, (path, label) in enumerate(zip(paths, labels)):
            sheet.paste(label_image(path, f"{label} · {name}", (440, 335)), (col * 440, row * 335))
    sheet.save(ASSETS / "self_trained_examples.png")


def prepare_assets():
    crop_paper_assets()
    compose_dataset_examples()
    compose_alignment_examples()
    compose_self_trained_examples()
    copy2(ROOT / "results" / "training_analysis_20260727_siba_official_protocol" / "loss_curve.png", ASSETS / "loss_curve.png")
    for dataset in ("MSRS", "M3FD_2x", "TNO"):
        source = ROOT / "results" / "metrics_20260727_siba_official_protocol" / "plots" / f"{dataset}_metric_ratio.png"
        copy2(source, ASSETS / f"{dataset}_metric_ratio.png")


def remove_all_slides(prs):
    while prs.slides:
        r_id = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[-1]


def update_master(prs):
    for master in prs.slide_masters:
        for shape in master.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if "大规模图像" in text:
                shape.text = "SIBA 的 Jittor 复现 · ICCV 2025"
            elif "程明明" in text:
                shape.text = "汇报人：________"
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = "Microsoft YaHei"


def set_text_style(paragraph, size=18, color=TEXT, bold=False, font="Microsoft YaHei", align=None):
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def add_text(slide, x, y, w, h, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei", margin=0.03):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(4)
        set_text_style(paragraph, size, color, bold, font, align)
    return shape


def add_bullets(slide, items, x, y, w, h, size=18, color=TEXT, leading=7):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.space_after = Pt(leading)
        paragraph.line_spacing = 1.12
        set_text_style(paragraph, size, color)
    return shape


def add_panel(slide, x, y, w, h, fill=LIGHT_GRAY, line=MID_GRAY, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_tag(slide, x, y, w, text, fill=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    set_text_style(paragraph, 12, WHITE, True, align=PP_ALIGN.CENTER)
    return shape


def add_caption(slide, x, y, w, text):
    return add_text(slide, x, y, w, 0.25, text, 10, MUTED)


def add_picture(slide, path, x, y, w, h=None):
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    image = Image.open(path)
    ratio = image.width / image.height
    box_ratio = w / h
    if ratio > box_ratio:
        height = w / ratio
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y + (h - height) / 2), width=Inches(w), height=Inches(height))
    width = h * ratio
    return slide.shapes.add_picture(str(path), Inches(x + (w - width) / 2), Inches(y), width=Inches(width), height=Inches(h))


def add_line(slide, x1, y1, x2, y2, color=BLUE, width=1.7, arrow=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_code(slide, x, y, w, h, title, lines, highlights=()):
    add_panel(slide, x, y, w, h, RGBColor(248, 249, 251), RGBColor(205, 211, 218), False)
    add_text(slide, x + 0.15, y + 0.08, w - 0.3, 0.3, title, 13, NAVY, True)
    shape = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.43), Inches(w - 0.3), Inches(h - 0.52))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = 0
    frame.margin_right = 0
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(1)
        set_text_style(paragraph, 10.5, RED if index in highlights else TEXT, index in highlights, "Consolas")
    return shape


def add_table(slide, x, y, w, h, data, col_widths=None, font_size=13, highlight_rows=()):
    shape = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if col_widths:
        for index, width in enumerate(col_widths):
            table.columns[index].width = Inches(width)
    for row_index, row in enumerate(data):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            if row_index == 0:
                cell.fill.fore_color.rgb = NAVY
            elif row_index in highlight_rows:
                cell.fill.fore_color.rgb = LIGHT_ORANGE
            elif row_index % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(241, 245, 248)
            else:
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                set_text_style(paragraph, font_size, WHITE if row_index == 0 else TEXT, row_index == 0 or row_index in highlight_rows, align=PP_ALIGN.CENTER)
                paragraph.space_after = Pt(0)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def set_title(slide, title, subtitle=None):
    slide.shapes.title.text = title
    set_text_style(slide.shapes.title.text_frame.paragraphs[0], 27, NAVY, True)
    if subtitle:
        add_text(slide, 0.55, 0.75, 12.1, 0.28, subtitle, 11, MUTED)


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title, subtitle)
    return slide


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_flow_box(slide, x, y, w, h, head, body="", fill=LIGHT_BLUE, accent=BLUE):
    add_panel(slide, x, y, w, h, fill, accent, True)
    add_text(slide, x + 0.12, y + 0.12, w - 0.24, 0.34, head, 16, NAVY, True, PP_ALIGN.CENTER)
    if body:
        add_text(slide, x + 0.12, y + 0.55, w - 0.24, h - 0.65, body, 13, TEXT, False, PP_ALIGN.CENTER)


outline = []


def record(number, title, text, images):
    outline.append({"number": number, "title": title, "text": text, "images": images})


def build_deck():
    prepare_assets()
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)
    update_master(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "SIBA 的 Jittor 完整复现"
    set_text_style(slide.shapes.title.text_frame.paragraphs[0], 38, WHITE, True, align=PP_ALIGN.CENTER)
    subtitle = slide.placeholders[1]
    subtitle.text = "The Source Image is the Best Attention for Infrared and Visible Image Fusion\nICCV 2025 · 培育期学习汇报 · 完整训练与对齐实验\n汇报人：________"
    for index, paragraph in enumerate(subtitle.text_frame.paragraphs):
        set_text_style(paragraph, 20 if index == 0 else 16, NAVY if index == 0 else TEXT, index == 0, align=PP_ALIGN.CENTER)
        paragraph.space_after = Pt(9)
    add_notes(slide, "本次汇报复现 ICCV 2025 的 SIBA。工作内容包括完整源码迁移、全量训练、测试、数值对齐、论文指标复现和运行记录。PPT 中只使用真实论文图、真实代码和真实实验结果。")
    record(1, "SIBA 的 Jittor 完整复现", ["论文全名", "ICCV 2025", "完整训练与对齐实验", "汇报人留空"], [])

    slide = new_slide(prs, "汇报主线", "先讲论文，再讲迁移和完整实验")
    items = [
        ("01", "论文问题与方法", "源图像为什么可以作为注意力"),
        ("02", "Jittor 忠实迁移", "逐文件镜像与框架替换"),
        ("03", "完整训练与对齐", "全量数据、60 轮训练、706 张测试图"),
        ("04", "结果与分析", "论文指标、框架差值、速度和局限"),
    ]
    for index, (number, head, body) in enumerate(items):
        x = 0.75 + index * 3.08
        add_tag(slide, x, 1.35, 0.75, number, ORANGE if index == 0 else BLUE)
        add_text(slide, x, 1.85, 2.7, 0.45, head, 19, NAVY, True)
        add_text(slide, x, 2.45, 2.65, 1.15, body, 15, TEXT)
        if index < 3:
            add_line(slide, x + 2.55, 3.75, x + 3.0, 3.75, BLUE, 2, True)
    add_text(slide, 0.8, 5.35, 11.8, 0.5, "方法讲解控制在前半部分；代码、训练、测试和对齐实验是本次复现的重点。", 19, BLUE, True, PP_ALIGN.CENTER)
    add_notes(slide, "汇报分为四部分。前半部分用论文图说明方法，后半部分重点说明迁移过程中做了什么、怎样验证没有改动算法，以及完整实验结果。")
    record(2, "汇报主线", [head + "：" + body for _, head, body in items], [])

    slide = new_slide(prs, "论文信息与选题核验")
    add_panel(slide, 0.65, 1.15, 7.5, 4.9, LIGHT_BLUE, MID_GRAY, False)
    add_text(slide, 0.95, 1.42, 6.9, 1.15, "Song Wang et al.\nThe Source Image is the Best Attention for Infrared and Visible Image Fusion", 18, NAVY, True)
    add_text(slide, 0.95, 2.72, 6.9, 0.42, "Proceedings of ICCV 2025, pp. 13513–13522", 17, BLUE, True)
    add_bullets(slide, [
        "发表时间不超过 2 年，主题与红外—可见光融合一致",
        "官方 PyTorch 代码、训练脚本、测试脚本和权重均公开",
        "Jittor-Sprouts 固定提交中没有 SIBA",
        "公开检索未发现 SIBA 的 Jittor 实现",
    ], 0.95, 3.38, 6.9, 2.2, 16)
    checks = [("顶会", "ICCV 2025"), ("近两年", "满足"), ("官方代码", "完整"), ("Jittor 重复", "未发现")]
    for index, (head, value) in enumerate(checks):
        y = 1.35 + index * 1.15
        add_panel(slide, 8.65, y, 3.8, 0.85, WHITE, BLUE, True)
        add_text(slide, 8.9, y + 0.14, 1.25, 0.35, head, 15, MUTED, True)
        add_text(slide, 10.1, y + 0.12, 2.0, 0.4, value, 18, TEAL, True, PP_ALIGN.RIGHT)
    add_caption(slide, 0.7, 6.45, 11.9, "核验记录：docs/SELECTION_AUDIT.md；官方代码冻结提交 880a1ddf9eaa610c64e5f25f87fbb146448addc9。")
    add_notes(slide, "先说明选题合规性。SIBA 正式发表于 ICCV 2025，官方仓库包含训练和测试代码，也提供了发布权重。选题时固定了官方提交，并检查 Jittor-Sprouts 和公开仓库，未发现重复的 Jittor 实现。")
    record(3, "论文信息与选题核验", ["完整论文出处", "四项任务合规检查", "官方提交哈希"], [])

    slide = new_slide(prs, "论文关注的问题", "跨模态注意力需要先解决特征差异")
    add_flow_box(slide, 0.8, 1.35, 3.1, 1.45, "红外深层特征", "热目标显著\n纹理较弱", LIGHT_ORANGE, ORANGE)
    add_flow_box(slide, 0.8, 3.5, 3.1, 1.45, "可见光深层特征", "边缘与纹理丰富\n受光照影响", LIGHT_BLUE, BLUE)
    add_flow_box(slide, 5.0, 2.3, 3.25, 1.65, "直接跨模态注意力", "Q、K、V 都来自深层特征\n模态差异会增加匹配难度", LIGHT_GRAY, MUTED)
    add_flow_box(slide, 9.3, 2.3, 3.05, 1.65, "论文的问题", "能否使用更直接的源图像先验\n引导跨模态信息选择", LIGHT_BLUE, TEAL)
    add_line(slide, 3.9, 2.05, 5.0, 2.75, ORANGE, 2, True)
    add_line(slide, 3.9, 4.2, 5.0, 3.5, BLUE, 2, True)
    add_line(slide, 8.25, 3.1, 9.3, 3.1, TEAL, 2.3, True)
    add_text(slide, 0.8, 5.7, 11.6, 0.55, "SIBA 不再只依赖深层特征生成注意力，而是把源图像本身用于查询。", 20, NAVY, True, PP_ALIGN.CENTER)
    add_notes(slide, "传统跨模态注意力通常让查询、键和值都来自深层特征。红外和可见光的统计分布不同，直接匹配会增加学习难度。SIBA 的切入点是：源图像本身已经包含明显的区域权重，能不能直接拿它来生成查询。")
    record(4, "论文关注的问题", ["红外和可见光深层特征差异", "直接跨模态注意力的匹配难度", "引出源图像先验"], [])

    slide = new_slide(prs, "核心观察：红外图像具有注意力属性")
    add_picture(slide, ASSETS / "source_attention_observation.png", 0.65, 1.12, 8.1, 4.15)
    add_bullets(slide, [
        "红外热目标区域与 Grad-CAM 的高响应区域相似",
        "源图像保留像素级位置，不需要从深层特征重新估计空间权重",
        "负变换图像补充暗区与背景结构",
        "这一观察直接形成 I-SCA 与 V-SCA 的查询设计",
    ], 9.0, 1.35, 3.65, 3.8, 16)
    add_caption(slide, 0.7, 5.55, 8.0, "图像来源：Wang et al., ICCV 2025, Fig. 1；仅提取图像区域，未截取论文正文与图注。")
    add_text(slide, 9.0, 5.45, 3.65, 0.55, "源图像不是额外标签，而是网络已有输入。", 17, TEAL, True, PP_ALIGN.CENTER)
    add_notes(slide, "论文用 Grad-CAM 做了一个直观观察：热目标集中的红外图像和分类网络的注意力图在空间分布上相似。作者因此把原始源图像和它的负变换作为查询输入，避免再用一个复杂模块从深层特征中预测注意力。")
    record(5, "核心观察：红外图像具有注意力属性", ["红外与 Grad-CAM 的相似性", "原图和负变换的作用", "源图像无需额外标签"], ["论文 Fig.1 的六张原始图像"])

    slide = new_slide(prs, "SIBA 的核心设计", "源图像生成 Q，深层特征生成 K 和 V")
    add_flow_box(slide, 0.65, 1.3, 2.5, 1.05, "源图像", "ir / vi", LIGHT_ORANGE, ORANGE)
    add_flow_box(slide, 0.65, 3.1, 2.5, 1.05, "负变换", "1-ir / 1-vi", LIGHT_BLUE, BLUE)
    add_flow_box(slide, 4.0, 2.05, 2.45, 1.35, "CBSM", "通道增强 + 空间映射\n生成查询 Q", LIGHT_GRAY, TEAL)
    add_flow_box(slide, 7.35, 1.3, 2.45, 1.05, "另一模态特征", "Restormer 输出", LIGHT_BLUE, BLUE)
    add_flow_box(slide, 7.35, 3.1, 2.45, 1.05, "线性变换", "生成 K、V", LIGHT_GRAY, MUTED)
    add_flow_box(slide, 10.7, 1.9, 2.0, 1.65, "Cross-Attn", "Q × K × V\n得到源图像导向特征", LIGHT_ORANGE, ORANGE)
    add_line(slide, 3.15, 1.82, 4.0, 2.4, ORANGE, 2, True)
    add_line(slide, 3.15, 3.62, 4.0, 3.05, BLUE, 2, True)
    add_line(slide, 6.45, 2.72, 10.7, 2.72, TEAL, 2.4, True)
    add_line(slide, 9.8, 1.82, 10.7, 2.38, BLUE, 2, True)
    add_line(slide, 9.8, 3.62, 10.7, 3.05, MUTED, 2, True)
    add_text(slide, 1.0, 5.2, 11.3, 0.62, "四个查询分支：ir、1-ir、vi、1-vi；每个分支都保留独立参数。", 20, NAVY, True, PP_ALIGN.CENTER)
    add_notes(slide, "核心计算关系是：CBSM 处理源图像并生成查询 Q；另一模态的深层特征生成 K 和 V。原图和负变换各形成一条查询分支，红外和可见光一共四条交叉注意力分支。")
    record(6, "SIBA 的核心设计", ["源图像和负变换", "CBSM 生成 Q", "另一模态特征生成 K、V", "四个独立查询分支"], [])

    slide = new_slide(prs, "总体结构", "双分支特征提取 + 四路源图像交叉注意力 + 通道融合")
    add_picture(slide, ASSETS / "paper_fig2_arch.png", 0.55, 1.1, 12.2, 5.15)
    add_caption(slide, 0.65, 6.35, 12.0, "来源：Wang et al., ICCV 2025, Fig. 2。图中 ir_n、vi_n 为负变换；C 表示通道拼接。")
    add_notes(slide, "整体结构从上下两路输入开始。红外和可见光先经过 SE-ResNet 和自注意力特征提取；原图和负变换经过 CBSM；四条交叉注意力结果拼接后，再用两层 SE-ResNet 融合并输出。")
    record(7, "总体结构", ["双分支编码", "四路交叉注意力", "拼接与输出"], ["论文 Fig.2 总体结构图"])

    slide = new_slide(prs, "特征提取：SE-ResNet 与 Restormer")
    add_flow_box(slide, 0.7, 1.35, 2.2, 1.0, "输入", "N×1×H×W", LIGHT_GRAY, MUTED)
    add_flow_box(slide, 3.45, 1.2, 3.0, 1.3, "SE-ResNet Block", "3×3 卷积 + PReLU\n通道重标定 + 残差", LIGHT_BLUE, BLUE)
    add_flow_box(slide, 7.0, 1.2, 3.0, 1.3, "Restormer Block", "自注意力 + GDFN\n保留全局依赖", LIGHT_ORANGE, ORANGE)
    add_flow_box(slide, 10.55, 1.35, 2.1, 1.0, "深层特征", "N×48×H×W", LIGHT_GRAY, TEAL)
    add_line(slide, 2.9, 1.85, 3.45, 1.85, BLUE, 2, True)
    add_line(slide, 6.45, 1.85, 7.0, 1.85, ORANGE, 2, True)
    add_line(slide, 10.0, 1.85, 10.55, 1.85, TEAL, 2, True)
    add_code(slide, 0.75, 3.25, 5.9, 2.45, "official_pytorch/models/SIBA.py", [
        "self.ir_conv = Res_SE(in_cha, mid_cha)",
        "self.vi_conv = Res_SE(in_cha, mid_cha)",
        "self.ir_sa.append(SA(mid_cha, mid_cha))",
        "self.vi_sa.append(SA(mid_cha, mid_cha))",
        "ir_sa = layer_ir(ir_sa)",
        "vi_sa = layer_vi(vi_sa)",
    ], (4, 5))
    add_bullets(slide, ["两路结构相同，但参数独立", "空间尺寸不变，通道由 1 提升到 48", "迁移保留模块顺序、残差路径和初始化"], 7.25, 3.45, 5.0, 2.1, 16)
    add_notes(slide, "特征提取部分由 SE-ResNet 和 Restormer 自注意力组成。两路模态使用相同结构，但没有共享参数。迁移时保留了 48 通道、模块顺序、残差路径和初始化方式。")
    record(8, "特征提取：SE-ResNet 与 Restormer", ["SE-ResNet 结构", "Restormer 自注意力", "两路参数独立", "对应源码片段"], ["可编辑流程图"])

    slide = new_slide(prs, "CBSM：把源图像映射为查询")
    add_picture(slide, ASSETS / "paper_fig3_cbsm.png", 0.7, 1.35, 5.4, 2.2)
    add_caption(slide, 0.8, 3.55, 5.2, "来源：Wang et al., ICCV 2025, Fig. 3。")
    add_flow_box(slide, 6.7, 1.25, 2.2, 1.0, "Conv 3×3", "通道提升", LIGHT_BLUE, BLUE)
    add_flow_box(slide, 9.35, 1.25, 2.2, 1.0, "SE Block", "抑制冗余通道", LIGHT_ORANGE, ORANGE)
    add_flow_box(slide, 9.35, 3.05, 2.2, 1.0, "空间映射", "PReLU + Conv", LIGHT_BLUE, TEAL)
    add_flow_box(slide, 6.7, 3.05, 2.2, 1.0, "残差相加", "保留原始响应", LIGHT_GRAY, MUTED)
    add_line(slide, 8.9, 1.75, 9.35, 1.75, BLUE, 2, True)
    add_line(slide, 10.45, 2.25, 10.45, 3.05, ORANGE, 2, True)
    add_line(slide, 9.35, 3.55, 8.9, 3.55, TEAL, 2, True)
    add_code(slide, 6.65, 4.5, 5.3, 1.5, "CBSM.execute", ["x = self.conv1(x)", "res = self.se(self.conv2(x))", "x = self.prelu(res + x)", "return self.conv3(x)"], (1, 2))
    add_notes(slide, "CBSM 是一个很小的卷积模块。它先提升通道，再用 SE 抑制冗余响应，最后做空间映射。论文把 CBSM 输出作为查询 Q，不是把源图像直接与深层特征相乘。")
    record(9, "CBSM：把源图像映射为查询", ["通道提升", "SE 通道抑制", "空间映射", "残差保留"], ["论文 Fig.3", "实际 CBSM 代码"])

    slide = new_slide(prs, "I-SCA / V-SCA：源图像交叉注意力")
    add_picture(slide, ASSETS / "paper_fig4_cross_attention.png", 0.55, 1.15, 6.2, 4.55)
    add_bullets(slide, ["Q：CBSM 处理后的源图像", "K、V：另一模态的深层特征", "Dconv 保持空间位置并降低计算量", "GDFN 继续筛选交互后的有效信息", "I-SCA 与 V-SCA 结构相同，仅输入内容不同"], 7.2, 1.35, 5.1, 3.7, 16)
    add_text(slide, 7.2, 5.25, 5.1, 0.62, "复杂度随像素数线性增长，避免窗口划分。", 18, TEAL, True, PP_ALIGN.CENTER)
    add_caption(slide, 0.7, 5.95, 6.0, "来源：Wang et al., ICCV 2025, Fig. 4；仅截取模块图。")
    add_notes(slide, "交叉注意力中，查询来自源图像，键和值来自另一模态的深层特征。Dconv 和逐点卷积用于生成 K、V，后面接 GDFN。I-SCA 和 V-SCA 结构完全相同，区别只是哪个源图像生成查询。")
    record(10, "I-SCA / V-SCA：源图像交叉注意力", ["Q/K/V 来源", "Dconv", "GDFN", "线性复杂度"], ["论文 Fig.4"])

    slide = new_slide(prs, "四路交叉注意力与最终融合")
    branch_info = [("ir → vi", "Q=ir", ORANGE), ("1-ir → vi", "Q=1-ir", BLUE), ("vi → ir", "Q=vi", TEAL), ("1-vi → ir", "Q=1-vi", NAVY)]
    for index, (head, body, color) in enumerate(branch_info):
        x = 0.65 + index * 2.65
        add_flow_box(slide, x, 1.3, 2.25, 1.25, head, body, LIGHT_GRAY, color)
        add_line(slide, x + 1.12, 2.55, 6.65, 3.6, color, 1.8, True)
    add_flow_box(slide, 5.35, 3.45, 2.6, 1.15, "通道拼接", "48×4 = 192 通道", LIGHT_BLUE, BLUE)
    add_flow_box(slide, 8.7, 3.45, 2.55, 1.15, "两层 Res_SE", "192 → 96 → 48", LIGHT_ORANGE, ORANGE)
    add_flow_box(slide, 11.65, 3.45, 1.1, 1.15, "输出", "48 → 1", LIGHT_GRAY, TEAL)
    add_line(slide, 7.95, 4.02, 8.7, 4.02, BLUE, 2.2, True)
    add_line(slide, 11.25, 4.02, 11.65, 4.02, TEAL, 2.2, True)
    add_text(slide, 0.85, 5.35, 11.7, 0.6, "原图与负变换不是共享一个注意力模块，而是四组独立参数。", 20, NAVY, True, PP_ALIGN.CENTER)
    add_notes(slide, "模型一共有四路交叉注意力：原红外、负红外、原可见光、负可见光。四路输出在通道维拼接，然后通过两层 Res_SE 压缩到 48 通道，最后输出单通道融合亮度。")
    record(11, "四路交叉注意力与最终融合", ["四个独立分支", "通道拼接 192", "Res_SE 逐级压缩", "单通道输出"], ["可编辑流程图"])

    slide = new_slide(prs, "损失函数", "官方代码的三个约束与权重保持不变")
    panels = [
        (0.7, "强度损失", "L_int = ‖max(I_ir, I_vi) − I_f‖₁", "保留两路输入中更强的像素响应", LIGHT_BLUE, BLUE),
        (4.8, "Sobel 梯度损失", "L_S = ‖max(∇S I_ir, ∇S I_vi) − ∇S I_f‖₁", "保留方向边缘与局部变化", LIGHT_ORANGE, ORANGE),
        (8.9, "Laplacian 梯度损失", "L_L = ‖JGrad(I_ir,I_vi) − ∇L I_f‖₁", "强化细节并抑制噪声", LIGHT_GRAY, TEAL),
    ]
    for x, head, equation, body, fill, line in panels:
        add_panel(slide, x, 1.25, 3.75, 3.6, fill, line, True)
        add_text(slide, x + 0.25, 1.55, 3.25, 0.4, head, 20, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.25, 2.2, 3.25, 0.85, equation, 15 if "Sobel" in head else 17, TEXT, True, PP_ALIGN.CENTER)
        add_text(slide, x + 0.3, 3.35, 3.15, 0.9, body, 16, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, 1.0, 5.35, 11.3, 0.55, "官方训练代码：L_total = 0.1 L_int + 1.0 L_S + 10.0 L_L", 21, BLUE, True, PP_ALIGN.CENTER)
    add_notes(slide, "损失函数完全跟随官方训练代码。强度损失取两路输入的逐像素最大值；Sobel 损失约束方向边缘；Laplacian 项权重为 10，用于细节保持和噪声抑制。迁移没有重新设计损失。")
    record(12, "损失函数", ["强度损失", "Sobel 梯度损失", "Laplacian 梯度损失", "总损失权重 0.1/1/10"], [])

    slide = new_slide(prs, "论文实验设置", "后续复现以官方代码和论文共同给出的设置为准")
    data = [["项目", "论文 / 官方设置", "本次复现"], ["训练数据", "MSRS 1083 + RoadScene 200", "完整 1283 对"], ["训练轮数", "60 epochs", "60 epochs"], ["Batch / Patch", "4 / 128×128", "4 / 128×128"], ["优化器", "Adam, lr=1e-4", "保持一致"], ["学习率", "每 25 轮 ×0.5", "保持一致"], ["测试集", "MSRS / M3FD_2x / TNO", "361 / 300 / 45"], ["硬件", "TITAN RTX 24GB", "RTX 3090 24GB"]]
    add_table(slide, 0.7, 1.25, 8.2, 4.95, data, [1.5, 3.35, 3.35], 14, (1, 2, 6))
    add_picture(slide, ASSETS / "paper_ablation_tables.png", 9.15, 1.35, 3.45, 3.8)
    add_caption(slide, 9.25, 5.25, 3.2, "论文 Table 3–5：模块消融与效率。")
    add_text(slide, 9.15, 5.55, 3.45, 0.55, "准确率可直接比较；速度只标注复现硬件。", 15, ORANGE, True, PP_ALIGN.CENTER)
    add_notes(slide, "论文训练数据为完整 MSRS 训练集和随机选取的 200 对 RoadScene。训练 60 轮，batch size 4，patch 128。论文使用 TITAN RTX 24GB；AutoDL 没有同型号，所以使用同显存级别的 RTX 3090，准确率可比，但速度不能声称同硬件复现。")
    record(13, "论文实验设置", ["数据、轮数、batch、patch、优化器、测试集和硬件对照"], ["论文消融与效率表局部截图"])

    slide = new_slide(prs, "迁移原则", "只替换框架，不改变算法、数据和训练逻辑")
    rules = [("同路径镜像", "官方 13 个 Python 文件均有同路径 Jittor 对应文件"), ("同结构", "模块、通道、残差、拼接顺序和初始化保持一致"), ("同训练", "损失权重、Adam、StepLR、梯度裁剪和 60 轮保持一致"), ("同数据", "完整数据集、同一配对清单、同一预处理和测试分辨率"), ("可审计", "所有框架强制修改、调试问题、日志和哈希均保留")]
    for index, (head, body) in enumerate(rules):
        y = 1.18 + index * 1.0
        add_tag(slide, 0.75, y, 1.65, head, BLUE if index < 4 else ORANGE)
        add_text(slide, 2.65, y - 0.02, 9.75, 0.55, body, 18, TEXT, index == 4)
    add_text(slide, 0.85, 6.3, 11.7, 0.4, "数值对齐通过后才启动 60 轮训练，避免把不完整迁移投入长时间训练。", 18, TEAL, True, PP_ALIGN.CENTER)
    add_notes(slide, "迁移的原则很简单：同路径、同结构、同训练、同数据，并且每处框架差异都可审计。先做源码覆盖和数值对齐，通过后才进行长时间训练。")
    record(14, "迁移原则", [body for _, body in rules], [])

    slide = new_slide(prs, "逐文件镜像与源码审计")
    tree = "official_pytorch/                 siba_jittor/\n├─ args/args_SIBA.py              ├─ args/args_SIBA.py\n├─ base_blocks/                    ├─ base_blocks/\n├─ loader/                         ├─ loader/\n├─ loss/loss.py                    ├─ loss/loss.py\n├─ models/SIBA.py                  ├─ models/SIBA.py\n├─ train.py / test.py              ├─ train.py / test.py\n└─ utils/                          └─ utils/\n                                    + compat/"
    add_code(slide, 0.65, 1.15, 7.2, 4.95, "目录对应关系", tree.splitlines())
    audit = [["审计项", "结果"], ["官方 Python 文件", "13"], ["缺失文件", "0"], ["缺失类 / 函数", "0"], ["模型参数张量", "137"], ["模型参数量", "565,941"], ["新增文件", "3 个 compat 文件"]]
    add_table(slide, 8.25, 1.25, 4.35, 3.9, audit, [2.55, 1.8], 15, (2, 3))
    add_text(slide, 8.35, 5.45, 4.15, 0.6, "新增 compat 只复现 PyTorch 的 Adam 与梯度裁剪数值规则。", 15, ORANGE, True, PP_ALIGN.CENTER)
    add_notes(slide, "官方仓库的 13 个 Python 文件全部建立同路径镜像。机器审计检查文件、类和函数符号，缺失项为零。额外的三个 compat 文件只负责复现 PyTorch 的优化器和梯度裁剪细节，没有新增模型功能。")
    record(15, "逐文件镜像与源码审计", ["13 个官方文件同路径对应", "0 缺失文件", "0 缺失类/函数", "137 个参数张量", "565,941 参数"], ["可编辑目录树", "源码审计表"])

    slide = new_slide(prs, "PyTorch 与 Jittor 代码对应", "核心前向逻辑逐行保持")
    pytorch_lines = ["def forward(self, ir, vi):", "    ir_raw_invert = 1 - ir", "    vi_raw_invert = 1 - vi", "    ir = self.ir_conv(ir)", "    vi = self.vi_conv(vi)", "    ir2vi_ca = layer_ir(vi_sa, w_ir)", "    vi2ir_ca = layer_vi(ir_sa, w_vi)", "    mixed = torch.cat([...], dim=1)", "    return self.out_conv(self.fuse_conv(mixed))"]
    jittor_lines = ["def execute(self, ir, vi):", "    ir_raw_invert = 1 - ir", "    vi_raw_invert = 1 - vi", "    ir = self.ir_conv(ir)", "    vi = self.vi_conv(vi)", "    ir2vi_ca = layer_ir(vi_sa, w_ir)", "    vi2ir_ca = layer_vi(ir_sa, w_vi)", "    mixed = jt.concat([...], dim=1)", "    return self.out_conv(self.fuse_conv(mixed))"]
    add_code(slide, 0.6, 1.2, 6.05, 4.95, "official_pytorch/models/SIBA.py", pytorch_lines, (0, 7))
    add_code(slide, 6.75, 1.2, 5.95, 4.95, "siba_jittor/models/SIBA.py", jittor_lines, (0, 7))
    add_text(slide, 2.2, 6.28, 9.0, 0.35, "变化只发生在框架 API：forward→execute，torch.cat→jt.concat。", 17, TEAL, True, PP_ALIGN.CENTER)
    add_notes(slide, "这里给出模型前向的直接对照。计算顺序、四条分支和拼接顺序不变，只把 PyTorch 的 forward 改为 Jittor 的 execute，把 torch.cat 改为 jt.concat。")
    record(16, "PyTorch 与 Jittor 代码对应", ["前向函数逐行对应", "仅替换 forward/execute 和 concat API"], ["两侧可编辑代码"])

    slide = new_slide(prs, "框架强制修改", "非机械替换均做了独立数值验证")
    substitutions = [["位置", "PyTorch", "Jittor 处理", "验证"], ["模块调用", "forward", "execute", "结构一致"], ["L2 normalize", "F.normalize", "同公式显式实现", "误差约 1e-12"], ["Laplacian", "Kornia 0.7.0", "同核、reflect padding、分组卷积", "损失对齐"], ["Adam", "torch.optim.Adam", "compat/pytorch_adam.py", "参数误差 2.98e-8"], ["梯度裁剪", "clip_grad_norm_", "同归约顺序", "误差 4.89e-9"], ["Dataset", "DataLoader", "set_attrs", "真实 crop 像素一致"], ["保存格式", ".pth", ".pkl", "状态字典结构一致"]]
    add_table(slide, 0.55, 1.15, 12.25, 5.25, substitutions, [1.65, 2.3, 4.9, 3.4], 13, (4, 5))
    add_notes(slide, "需要重点说明的不是普通 API 替换，而是 Laplacian、Adam 和梯度裁剪。Jittor 与 PyTorch 的默认实现存在数值细节差异，因此按 PyTorch 1.10 的规则实现兼容层，并用同一梯度验证更新后的参数。")
    record(17, "框架强制修改", ["normalize", "Kornia Laplacian", "Adam", "梯度裁剪", "Dataset", "保存格式"], ["可编辑对照表"])

    slide = new_slide(prs, "迁移完整性验收", "文件覆盖只是第一步，还要验证计算结果")
    stages = [("源码覆盖", "13/13 文件\n0 缺失符号"), ("模型结构", "137 参数张量\n565,941 参数"), ("受控数值", "前向、损失、梯度\n裁剪、Adam"), ("真实数据", "crop 像素一致\n同一清单"), ("端到端", "官方权重 706 张\n最大像素差 1")]
    for index, (head, body) in enumerate(stages):
        x = 0.55 + index * 2.55
        add_flow_box(slide, x, 1.55, 2.15, 1.45, head, body, LIGHT_BLUE if index % 2 == 0 else LIGHT_ORANGE, BLUE if index % 2 == 0 else ORANGE)
        if index < len(stages) - 1:
            add_line(slide, x + 2.15, 2.27, x + 2.55, 2.27, TEAL, 2, True)
    add_panel(slide, 1.1, 4.05, 11.15, 1.4, LIGHT_GRAY, TEAL, True)
    add_text(slide, 1.45, 4.24, 10.45, 0.72, "最终自动验收：base_complete=true · gpu_complete=true · metrics_complete=true · complete=true", 19, NAVY, True, PP_ALIGN.CENTER)
    add_caption(slide, 1.45, 5.05, 10.45, "正式报告：docs/final_validation.json。")
    add_notes(slide, "完整性验收分五层。文件和参数数量一致后，还要做受控前向、损失、梯度和一步更新对齐；真实数据 crop 做像素比较；最后用官方权重在全部 706 张测试图上做端到端输出比较。")
    record(18, "迁移完整性验收", ["五级验收流程", "最终 complete=true"], ["可编辑验收流程"])

    slide = new_slide(prs, "数值对齐结果", "同输入、同参数、同损失、同梯度")
    data = [["检查项", "结果", "说明"], ["初始参数最大误差", "0", "权重完全一致"], ["激活最大绝对误差", "2.0206e-4", "所有主要中间层"], ["总损失最大绝对误差", "2.9802e-6", "三项损失加权后"], ["梯度余弦相似度", "0.999945", "全部参数梯度"], ["裁剪后最大误差", "4.8894e-9", "同 PyTorch 归约顺序"], ["Adam 更新最大误差", "2.9802e-8", "一步参数更新"], ["官方权重 PNG 最大误差", "1/255", "真实全分辨率图像"]]
    add_table(slide, 0.75, 1.2, 8.45, 5.1, data, [3.05, 2.05, 3.35], 14, (4, 6, 7))
    add_panel(slide, 9.55, 1.35, 2.75, 2.2, LIGHT_BLUE, BLUE, True)
    add_text(slide, 9.8, 1.65, 2.25, 0.38, "端到端输出", 18, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, 9.8, 2.05, 2.25, 1.1, "706 / 706\n文件名和尺寸一致", 22, TEAL, True, PP_ALIGN.CENTER)
    add_panel(slide, 9.55, 4.0, 2.75, 1.55, LIGHT_ORANGE, ORANGE, True)
    add_text(slide, 9.8, 4.25, 2.25, 0.9, "这部分用于证明\n迁移没有改变模型计算。", 17, NAVY, True, PP_ALIGN.CENTER)
    add_notes(slide, "数值对齐结果表明，受控条件下前向、损失、梯度、裁剪和 Adam 更新都与 PyTorch 高度一致。官方发布权重在三套测试集的 706 张图上，文件名和尺寸全部一致，最大像素差只有一个 uint8 灰度级。")
    record(19, "数值对齐结果", ["参数、激活、损失、梯度、裁剪、Adam、官方权重输出误差"], ["可编辑数值表"])

    slide = new_slide(prs, "完整数据准备与可追溯性")
    add_picture(slide, ASSETS / "dataset_examples.png", 0.65, 1.15, 5.8, 5.15)
    data = [["数据集", "数量", "用途"], ["MSRS train", "1083", "训练"], ["RoadScene", "200 / 221", "训练补充"], ["MSRS test", "361", "测试"], ["M3FD_2x", "300", "测试"], ["TNO", "45", "测试"]]
    add_table(slide, 6.8, 1.3, 5.5, 3.4, data, [2.2, 1.3, 2.0], 15, (1, 2))
    add_bullets(slide, ["每个配对记录文件名、尺寸和 SHA256", "M3FD 严格执行官方半分辨率预处理", "TNO 使用官方 45 对，不使用常见 25 对子集", "RoadScene 200 个文件名未公开，使用 seed=2025 的确定性清单"], 6.95, 4.75, 5.3, 1.72, 14)
    add_notes(slide, "训练集包含 1083 对 MSRS 和 200 对 RoadScene，共 1283 对。测试集为 361 张 MSRS、300 张半分辨率 M3FD 和 45 张完整 TNO。每个文件都记录哈希。唯一无法与作者完全一致的是作者没有公开 RoadScene 的 200 个文件名和随机种子。")
    record(20, "完整数据准备与可追溯性", ["1283 对训练图", "706 对测试图", "哈希清单", "RoadScene 未公开子集说明"], ["三套真实数据的红外/可见光样例"])

    slide = new_slide(prs, "完整训练与日志", "Jittor 和 PyTorch 均按官方设置训练 60 轮")
    add_flow_box(slide, 0.65, 1.3, 2.35, 1.2, "启动前检查", "源码审计\n数值对齐", LIGHT_BLUE, BLUE)
    add_flow_box(slide, 3.45, 1.3, 2.35, 1.2, "Jittor 训练", "60 轮\n4079 s", LIGHT_ORANGE, ORANGE)
    add_flow_box(slide, 6.25, 1.3, 2.35, 1.2, "PyTorch 基线", "60 轮\n2221 s", LIGHT_BLUE, TEAL)
    add_flow_box(slide, 9.05, 1.3, 2.35, 1.2, "保存产物", "checkpoint\n420 条 loss", LIGHT_GRAY, MUTED)
    add_line(slide, 3.0, 1.9, 3.45, 1.9, BLUE, 2, True)
    add_line(slide, 5.8, 1.9, 6.25, 1.9, ORANGE, 2, True)
    add_line(slide, 8.6, 1.9, 9.05, 1.9, TEAL, 2, True)
    add_code(slide, 0.75, 3.25, 6.15, 2.2, "训练会话", ["RUN_TAG=20260727_siba_official_protocol", "bash scripts/train_full_sequence_screen.sh", "screen -r kk", "tail -f logs/.../train.log"], (2,))
    run_data = [["框架", "轮数", "时间", "日志条数"], ["Jittor", "60", "67.98 min", "420"], ["PyTorch", "60", "37.02 min", "420"]]
    add_table(slide, 7.3, 3.25, 5.1, 2.2, run_data, [1.35, 1.0, 1.65, 1.1], 15, (1,))
    add_text(slide, 0.85, 5.9, 11.5, 0.45, "GPU 任务结束后远程主机已关机；后续汇总与 PPT 均在本地完成。", 17, TEAL, True, PP_ALIGN.CENTER)
    add_notes(slide, "训练在 screen -S kk 中执行，方便断开连接后继续运行。Jittor 60 轮用时约 68 分钟，PyTorch 约 37 分钟，两者各记录 420 条官方间隔的 loss。训练、推理和计时完成后远程主机已关闭。")
    record(21, "完整训练与日志", ["screen -S kk", "Jittor/PyTorch 60 轮", "训练时间", "420 条日志", "GPU 完成后关机"], [])

    slide = new_slide(prs, "训练曲线", "两套完整训练均收敛，不声称逐批轨迹相同")
    add_picture(slide, ASSETS / "loss_curve.png", 0.55, 1.05, 9.1, 5.65)
    add_panel(slide, 9.85, 1.35, 2.55, 1.55, LIGHT_BLUE, BLUE, True)
    add_text(slide, 10.05, 1.62, 2.15, 0.35, "Epoch 60 均值", 16, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, 10.05, 1.98, 2.15, 0.82, "Jittor 0.03463\nPyTorch 0.04364", 18, TEAL, True, PP_ALIGN.CENTER)
    add_panel(slide, 9.85, 3.25, 2.55, 2.35, LIGHT_ORANGE, ORANGE, True)
    add_text(slide, 10.05, 3.52, 2.15, 1.65, "数据加载器的 shuffle 实现不同，因此完整训练轨迹会分叉。\n\n受控数值实验用于判断框架一致性。", 15, TEXT, False, PP_ALIGN.CENTER)
    add_notes(slide, "两条曲线都在 20 轮左右明显下降，并在后半程稳定。这里不把两条曲线画得完全重合，因为两个框架的数据加载器打乱顺序不同。迁移是否正确由上一页的受控数值对齐和官方权重输出对齐判断。")
    record(22, "训练曲线", ["60 轮 loss 曲线", "Epoch 60 均值", "shuffle 导致轨迹分叉说明"], ["真实训练 loss 曲线"])

    slide = new_slide(prs, "完整推理与可视化", "发布权重用于严格对齐，自训练权重用于完整流程验证")
    add_picture(slide, ASSETS / "official_alignment_examples.png", 0.55, 1.0, 12.2)
    add_caption(slide, 0.65, 3.36, 12.0, "MSRS：官方发布权重，Jittor 与 PyTorch 输出最大 uint8 差值为 1。")
    add_picture(slide, ASSETS / "self_trained_examples.png", 0.55, 3.7, 12.2)
    add_caption(slide, 0.65, 6.06, 12.0, "TNO：两套 60 轮自训练权重；保留真实差异，不进行人工修图。")
    add_text(slide, 0.8, 6.34, 11.7, 0.32, "四个实验分支均完成 MSRS 361 + M3FD 300 + TNO 45 张推理。", 14, TEAL, True, PP_ALIGN.CENTER)
    add_notes(slide, "上方左侧展示发布权重的输出，Jittor 和 PyTorch 几乎不可区分。右侧展示两套自训练权重在 TNO 上的结果。所有图像直接来自推理目录，没有人工修改或生成。")
    record(23, "完整推理与可视化", ["四个实验分支完成全部 706 张", "发布权重严格对齐", "自训练真实差异"], ["真实 MSRS 发布权重对比", "真实 TNO 自训练对比"])

    slide = new_slide(prs, "论文指标复现", "发布权重在三套数据集上均复现论文三位小数")
    metrics = ["VIF", "SCD", "MI", "Qabf", "SSIM", "MS-SSIM", "FMI"]
    rows = [["数据集 / 实现"] + metrics]
    values = {"MSRS Paper": [1.061, 1.700, 5.111, 0.715, 0.981, 0.975, 0.933], "MSRS Jittor": [1.060896, 1.700205, 5.110839, 0.715261, 0.980834, 0.975080, 0.932669], "M3FD Paper": [0.759, 1.660, 3.771, 0.654, 0.964, 0.923, 0.883], "M3FD Jittor": [0.759335, 1.660086, 3.770833, 0.653689, 0.964495, 0.923055, 0.882752], "TNO Paper": [0.836, 1.724, 3.508, 0.588, 0.932, 0.904, 0.914], "TNO Jittor": [0.835515, 1.724506, 3.507112, 0.587575, 0.931762, 0.904372, 0.914384]}
    for name, numbers in values.items():
        rows.append([name] + [f"{number:.6f}" if "Jittor" in name else f"{number:.3f}" for number in numbers])
    add_table(slide, 0.4, 1.15, 12.55, 4.7, rows, [1.85] + [1.52] * 7, 11.5, (2, 4, 6))
    add_text(slide, 0.8, 6.02, 11.7, 0.42, "Jittor 发布权重结果按论文精度取三位小数后，21 个数值全部一致。", 19, TEAL, True, PP_ALIGN.CENTER)
    add_notes(slide, "这一页是最直接的论文指标复现。三套数据集、七项指标共 21 个数值，Jittor 使用发布权重计算后，按论文表格保留三位小数全部一致。每张图的七项指标和汇总 CSV 都已保留。")
    record(24, "论文指标复现", ["MSRS/M3FD/TNO 七项论文指标", "Jittor 发布权重 21 个数值三位小数全部一致"], ["可编辑指标表"])

    slide = new_slide(prs, "Jittor 与 PyTorch 的端到端指标差值")
    add_picture(slide, ASSETS / "MSRS_metric_ratio.png", 0.55, 1.05, 8.15, 4.65)
    deltas = [["指标", "最大绝对差"], ["VIF", "5.42e-5"], ["SCD", "4.08e-5"], ["MI", "9.07e-4"], ["Qabf", "3.42e-5"], ["SSIM", "6.52e-6"], ["MS-SSIM", "4.24e-6"], ["FMI", "3.31e-6"]]
    add_table(slide, 9.0, 1.2, 3.35, 3.85, deltas, [1.65, 1.7], 14, (3,))
    add_text(slide, 8.9, 5.35, 3.55, 0.8, "差值统计基于发布权重。\n自训练权重因 shuffle 和未公开子集而单独报告。", 14, ORANGE, True, PP_ALIGN.CENTER)
    add_caption(slide, 0.7, 5.85, 7.8, "柱高为“指标值 / 论文值”；原始 CSV：results/metrics_20260727_siba_official_protocol/。")
    add_notes(slide, "发布权重下，Jittor 和 PyTorch 的七项指标差值都很小，最大的 MI 差值不到 0.001。自训练权重不能作为逐批数值对齐，因为打乱顺序不同，而且作者没有公开 RoadScene 200 个文件名，所以单独保存和解释。")
    record(25, "Jittor 与 PyTorch 的端到端指标差值", ["发布权重七项最大绝对差", "自训练差异的合理边界"], ["真实 MSRS 指标比值图"])

    slide = new_slide(prs, "运行效率与资源记录", "同步计时用于框架比较，论文异步计时单独保留")
    timing = [["数据集", "Jittor FPS", "PyTorch FPS", "论文 FPS*"], ["MSRS", "9.12", "11.99", "132.303"], ["M3FD_2x", "12.66", "19.90", "137.271"], ["TNO", "6.93", "12.31", "126.537"]]
    add_table(slide, 0.75, 1.3, 7.0, 3.0, timing, [1.6, 1.65, 1.75, 2.0], 16, (1, 2, 3))
    add_panel(slide, 8.2, 1.3, 4.2, 3.0, LIGHT_BLUE, BLUE, True)
    add_text(slide, 8.5, 1.62, 3.6, 0.4, "资源记录", 19, NAVY, True, PP_ALIGN.CENTER)
    add_bullets(slide, ["RTX 3090 24GB", "峰值显存 14,587 MiB", "Jittor 训练 67.98 min", "PyTorch 训练 37.02 min"], 8.65, 2.2, 3.35, 1.65, 15)
    add_text(slide, 0.85, 4.85, 11.5, 0.95, "* 论文 test.py 未做 CUDA 同步，计时只覆盖异步调用；因此论文 FPS 不与同步 FPS 直接比较。\n本次框架速度结论：在 RTX 3090 上，Jittor 同步推理慢于 PyTorch。", 16, ORANGE, True, PP_ALIGN.CENTER)
    add_notes(slide, "速度部分区分两种口径。论文代码在 CUDA 调用前后直接计时，没有同步，所以会得到 100 多 FPS。框架比较使用同步计时，Jittor 在这台 RTX 3090 上慢于 PyTorch。这里不把不同计时口径混在一起。")
    record(26, "运行效率与资源记录", ["三套测试集同步 FPS", "论文异步 FPS 说明", "峰值显存与训练时间"], ["可编辑速度表"])

    slide = new_slide(prs, "复现边界与我的思考")
    topics = [("官方信息不完整", "RoadScene 的 200 个文件名和随机种子未公开；自训练结果不能声称与作者训练轨迹完全相同。"), ("论文与代码有差异", "论文描述 RGB→YCbCr 训练，发布的 train_loader 实际直接读取灰度图；本次跟随发布代码。"), ("源图像先验的条件", "方法依赖输入配准和源图像质量。错位、饱和或传感器噪声可能使查询产生错误引导。"), ("负变换是固定先验", "1−I 能补充暗区信息，但也可能放大过曝或噪声。后续可研究可靠性权重，但不属于本次复现。"), ("评价仍以图像指标为主", "七项指标验证了融合质量；在检测、分割等下游任务上是否同样有效，还需要独立实验。")]
    for index, (head, body) in enumerate(topics):
        y = 1.05 + index * 1.05
        add_tag(slide, 0.65, y, 2.25, head, ORANGE if index < 2 else BLUE)
        add_text(slide, 3.15, y - 0.02, 9.45, 0.72, body, 16, TEXT, index >= 2)
    add_notes(slide, "这部分只讨论真实边界。第一，作者没有公开 RoadScene 子集；第二，论文预处理描述和发布代码不完全一致，本次遵循代码。方法层面，源图像作为查询依赖配准和传感器质量；负变换也可能放大异常。最后，图像指标之外还需要下游任务验证。")
    record(27, "复现边界与我的思考", [body for _, body in topics], [])

    slide = new_slide(prs, "开源仓库与提交材料", "README 已覆盖任务要求的代码与实验记录")
    repo_tree = "SIBA-Jittor/\n├─ official_pytorch/     冻结官方源码\n├─ siba_jittor/          Jittor 同路径镜像\n├─ scripts/              环境、数据、训练、推理\n├─ tools/                对齐、指标、可视化、验收\n├─ data_manifests/       完整数据哈希清单\n├─ logs/                 训练与性能日志\n├─ results/              曲线、指标和真实融合图\n├─ README.md\n└─ MIGRATION_LOG.md"
    add_code(slide, 0.65, 1.15, 6.35, 5.1, "仓库结构", repo_tree.splitlines())
    add_bullets(slide, ["环境配置与版本记录", "完整数据准备脚本和 SHA256 清单", "Jittor / PyTorch 训练与测试命令", "420 条训练日志、Loss 曲线、性能日志", "逐模块数值对齐与完整指标 CSV", "最终验收报告 complete=true"], 7.35, 1.25, 5.0, 3.4, 16)
    add_panel(slide, 7.35, 4.92, 5.0, 1.12, LIGHT_ORANGE, ORANGE, True)
    add_text(slide, 7.6, 5.12, 4.5, 0.68, "GitHub 链接：待用户提供个人空仓库后推送\n当前不伪造链接。", 15, NAVY, True, PP_ALIGN.CENTER)
    add_notes(slide, "README 已包含环境、数据、训练、测试、日志、曲线、结果、性能和对齐方法。仓库本地已初始化并完成发布前检查。因为没有个人 GitHub 空仓库地址，PPT 中保留待填位置，不生成虚假链接。")
    record(28, "开源仓库与提交材料", ["仓库结构", "README 六类必备内容", "GitHub 链接待用户提供"], ["可编辑仓库树"])

    slide = new_slide(prs, "总结")
    summary = [("完整迁移", "13 个官方 Python 文件全部镜像，模型和训练逻辑保持不变。"), ("完整实验", "1283 对训练图、60 轮双框架训练、四个实验分支各完成 706 张推理。"), ("严格对齐", "受控数值、官方权重像素和七项指标均完成验证。"), ("论文复现", "发布权重的 21 个论文指标按三位小数全部一致。")]
    for index, (head, body) in enumerate(summary):
        y = 1.25 + index * 1.15
        add_tag(slide, 1.0, y, 1.8, head, BLUE if index < 3 else ORANGE)
        add_text(slide, 3.15, y - 0.02, 8.9, 0.72, body, 19, NAVY if index == 3 else TEXT, index == 3)
    add_text(slide, 1.0, 6.15, 11.3, 0.45, "感谢聆听，欢迎批评指正。", 24, BLUE, True, PP_ALIGN.CENTER)
    add_notes(slide, "总结来说，本次完成了完整源码迁移、完整数据训练和测试、严格数值对齐以及论文指标复现。所有结果都有原始日志和文件支撑，没有使用生成数据或人工修改结果。感谢聆听。")
    record(29, "总结", [body for _, body in summary] + ["感谢聆听，欢迎批评指正"], [])

    DELIVERABLE.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))


def write_outline():
    lines = ["# SIBA-Jittor PPT 完整大纲", "", "- 模板：TemplateMC-PPT.pptx", "- 页数：29 页", "- 主线：论文问题 → 方法 → Jittor 迁移 → 数值对齐 → 完整训练与测试 → 结果与思考", "- 所有实验图和数据均来自真实运行；没有生成或伪造实验结果。", ""]
    for entry in outline:
        lines.extend([f"## 第 {entry['number']} 页：{entry['title']}", "", "**页面文字**"])
        lines.extend(f"- {text}" for text in entry["text"])
        lines.extend(["", "**图片 / 图表**"])
        lines.extend(f"- {image}" for image in entry["images"] or ["无外部图片；使用可编辑文字、流程或表格。"])
        lines.append("")
    OUTLINE.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    write_outline()
    print(OUT)
    print(OUTLINE)
